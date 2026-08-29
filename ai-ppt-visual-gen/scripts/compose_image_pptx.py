#!/usr/bin/env python3
"""Assemble manifest-listed raster slide pages into an image-only PPTX."""
from __future__ import annotations

import argparse
import json
from pathlib import Path

from atomic_output import atomic_save_presentation


RATIOS = {
    "16:9": (13.333333, 7.5),
    "3:2": (12.0, 8.0),
}


def resolve(base: Path, value: str) -> Path:
    path = Path(value)
    return path.resolve() if path.is_absolute() else (base / path).resolve()


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("manifest")
    parser.add_argument("output")
    parser.add_argument("--ratio", choices=sorted(RATIOS))
    args = parser.parse_args()
    manifest_path = Path(args.manifest).resolve()
    data = json.loads(manifest_path.read_text(encoding="utf-8"))
    slides = data.get("slides")
    if not isinstance(slides, list) or not slides:
        raise SystemExit("manifest slides are required")
    ordered = sorted(slides, key=lambda item: item.get("slide_no", 0))
    numbers = [item.get("slide_no") for item in ordered]
    if numbers != list(range(1, len(ordered) + 1)):
        raise SystemExit(f"slide numbers must be contiguous from 1: {numbers}")
    ratios = {
        ((item.get("canvas") or {}).get("ratio"))
        for item in ordered
        if isinstance(item, dict)
    }
    ratio = args.ratio or (next(iter(ratios)) if len(ratios) == 1 else None)
    if ratio not in RATIOS:
        raise SystemExit("one consistent 16:9 or 3:2 ratio is required")

    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as exc:
        raise SystemExit(f"python-pptx unavailable: {exc}") from exc

    presentation = Presentation()
    presentation.slide_width = Inches(RATIOS[ratio][0])
    presentation.slide_height = Inches(RATIOS[ratio][1])
    blank = presentation.slide_layouts[6]
    base = manifest_path.parent
    for item in ordered:
        value = item.get("copied_to")
        if not isinstance(value, str) or not value:
            raise SystemExit(f"slide {item.get('slide_no')} copied_to is required")
        image = resolve(base, value)
        if not image.is_file():
            raise SystemExit(f"slide image missing: {image}")
        slide = presentation.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(image),
            0,
            0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )
    output = atomic_save_presentation(presentation, Path(args.output).resolve())
    print(json.dumps({
        "schema": "ai-ppt-visual-gen/image-pptx/v1",
        "valid": True,
        "output": str(output),
        "slide_count": len(ordered),
        "ratio": ratio,
        "editability": "image-only",
    }, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
