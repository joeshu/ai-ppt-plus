#!/usr/bin/env python3
"""Regression coverage for page-scoped editable object auditing."""
from __future__ import annotations

import json
import subprocess
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]


def main() -> int:
    with tempfile.TemporaryDirectory(prefix="editable-object-audit-") as temp:
        root = Path(temp)
        deck = root / "deck.pptx"
        manifest = root / "objects.json"
        report = root / "report.json"

        prs = Presentation()
        for slide_no in (1, 2):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
            shape.name = "shared-title"
            shape.text = f"Slide {slide_no}"
        prs.save(deck)
        manifest.write_text(json.dumps({
            "schema": "ai-ppt-plus/slide-object-manifest/v1",
            "slides": [
                {"slide_no": 1, "objects": [{"object_id": "shared-title", "role": "formal-text"}]},
                {"slide_no": 2, "objects": [{"object_id": "shared-title", "role": "formal-text"}]},
            ],
        }), encoding="utf-8")

        result = subprocess.run([
            sys.executable, "scripts/inspect_editable_objects.py", str(deck),
            "--object-manifest", str(manifest), "--report", str(report),
        ], cwd=ROOT, capture_output=True, text=True, check=False)
        assert result.returncode == 0, result.stdout + result.stderr
        data = json.loads(report.read_text(encoding="utf-8"))
        assert data["valid"] is True, data
        assert data["expected_object_count"] == 2, data
        assert data["observed_shape_count"] == 2, data
        assert data["errors"] == [], data
    print("page-scoped editable object audit: ok")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
