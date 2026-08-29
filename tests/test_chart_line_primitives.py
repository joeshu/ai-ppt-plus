#!/usr/bin/env python3
"""Regression tests for diagonal static-line primitives."""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SCRIPT_DIR = Path(__file__).resolve().parents[1] / "scripts"
sys.path.insert(0, str(SCRIPT_DIR))

from pptx_primitives import add_shapes  # noqa: E402


def main() -> int:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    deck = {
        "units": "px",
        "ref_width": 100,
        "ref_height": 100,
        "strict_input": True,
    }
    add_shapes(
        slide,
        [
            {"object_id": "descending", "type": "line", "x": 80, "y": 80, "x2": 20, "y2": 20, "line": "#667A45"},
            {"object_id": "legacy", "type": "line", "x": 20, "y": 20, "w": 30, "h": 10, "line": "#EAA035"},
        ],
        deck,
        100,
        100,
        presentation.slide_width,
        presentation.slide_height,
    )
    connectors = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.LINE]
    assert len(connectors) == 2, "explicit-endpoint and legacy line forms must both author"
    assert {shape.name for shape in connectors} == {"descending", "legacy"}
    print("diagonal line primitives: ok")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
