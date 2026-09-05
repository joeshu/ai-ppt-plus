"""Three-way coverage against a separately observed, source-bound inventory.

This contract checks provenance and completeness, not whether a model truly
looked at the image. Observation evidence must be supplied by the host.
"""
from __future__ import annotations

from .graph_ir import PageGraph
from hashlib import sha256
from pathlib import Path


def extract_pptx_objects(path: str, *, slide_index: int = 0) -> dict:
    """Read actual shape names recursively; authoring must bind stable IDs to names."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    deck = Path(path)
    before = sha256(deck.read_bytes()).hexdigest()
    slide = Presentation(deck).slides[slide_index]
    ids, records = [], []

    def visit(shapes):
        for shape in shapes:
            ids.append(shape.name)
            # A card/label autoshape may contain text but remains a native shape.
            # Classify from the OOXML container first so text-frame presence does
            # not let raster/native structure audits report the wrong object type.
            kind = ("group" if shape.shape_type == MSO_SHAPE_TYPE.GROUP else
                    "table" if shape.has_table else "chart" if shape.has_chart else
                    "image" if shape.shape_type == MSO_SHAPE_TYPE.PICTURE else
                    "connector" if shape.shape_type == MSO_SHAPE_TYPE.LINE else
                    "text" if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX else "shape")
            records.append({"id": shape.name, "type": kind,
                            "text": shape.text if shape.has_text_frame else None})
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                visit(shape.shapes)

    visit(slide.shapes)
    if sha256(deck.read_bytes()).hexdigest() != before:
        raise ValueError("PPTX changed during extraction")
    return {"method": "pptx-extraction", "deck_sha256": before, "object_ids": ids, "objects": records}


def audit_source_coverage(inventory: dict, graph: PageGraph,
                          observed: dict | None = None,
                          *, require_spatial_evidence: bool = False) -> dict:
    errors = []
    source_hash = inventory.get("source_sha256")
    if (not isinstance(source_hash, str) or len(source_hash) != 64
            or any(c not in "0123456789abcdef" for c in source_hash)):
        errors.append("invalid source SHA-256")
    if source_hash != graph.metadata.get("source_sha256"):
        errors.append("inventory/graph source mismatch")
    observation = inventory.get("observation_id")
    if (not observation or not graph.metadata.get("planning_observation_id")
            or observation == graph.metadata.get("planning_observation_id")):
        errors.append("independent observation required")
    if inventory.get("method") != "source-image-observation" or not inventory.get("evidence"):
        errors.append("source observation evidence required")
    if require_spatial_evidence:
        verification = inventory.get("spatial_verification")
        if not isinstance(verification, dict):
            errors.append("independent spatial verification required")
        else:
            verifier = verification.get("observation_id")
            if not verifier or verifier in {observation, graph.metadata.get("planning_observation_id")}:
                errors.append("spatial verification must use an independent observation")
            canvas = verification.get("canvas")
            if (not isinstance(canvas, (list, tuple)) or len(canvas) != 2
                    or any(isinstance(v, bool) or not isinstance(v, (int, float)) or v <= 0 for v in canvas)):
                errors.append("valid spatial verification canvas required")
            unexplained = verification.get("unexplained_regions")
            if not isinstance(unexplained, list):
                errors.append("spatial verification unexplained_regions list required")
            elif unexplained:
                errors.append("unexplained source regions remain")
    objects = inventory.get("objects")
    if not isinstance(objects, list) or not objects:
        return {"valid": False, "errors": errors + ["non-empty source objects required"]}
    if any(not isinstance(item, dict) for item in objects):
        return {"valid": False, "errors": errors + ["invalid source object record"]}
    source_ids = [item.get("id") for item in objects]
    if require_spatial_evidence:
        for item in objects:
            bbox = item.get("bbox")
            if (not isinstance(bbox, (list, tuple)) or len(bbox) != 4
                    or any(isinstance(v, bool) or not isinstance(v, (int, float)) for v in bbox)
                    or bbox[2] <= 0 or bbox[3] <= 0):
                errors.append(f"missing spatial bbox: {item.get('id')}")
    if any(not isinstance(i, str) or not i for i in source_ids) or len(set(source_ids)) != len(source_ids):
        return {"valid": False, "errors": errors + ["invalid or duplicate source IDs"]}
    bindings = {}
    for node in graph.nodes:
        sid = node.source.get("source_object_id")
        if not isinstance(sid, str):
            errors.append(f"invalid source binding: {node.id}")
            continue
        if sid not in source_ids:
            errors.append(f"unbound planned object: {node.id}")
        bindings.setdefault(sid, []).append(node.id)
    for item in objects:
        if item.get("uncertain"):
            errors.append(f"unresolved source object: {item.get('id')}")
        if item.get("id") not in bindings:
            errors.append(f"missing planned source object: {item.get('id')}")
    if observed is not None:
        if observed.get("method") != "pptx-extraction" or not observed.get("deck_sha256"):
            errors.append("independent PPTX extraction required")
        actual = observed.get("object_ids", [])
        if not isinstance(actual, list) or any(not isinstance(i, str) or not i for i in actual):
            return {"valid": False, "errors": errors + ["invalid extracted IDs"]}
        planned = {node.id for node in graph.nodes}
        if len(actual) != len(set(actual)):
            errors.append("duplicate extracted object IDs")
        for missing in sorted(planned - set(actual)):
            errors.append(f"missing final object: {missing}")
        for extra in sorted(set(actual) - planned):
            errors.append(f"undeclared final object: {extra}")
        if "objects" in observed:
            records = {item["id"]: item for item in observed["objects"]}
            for node in graph.nodes:
                record = records.get(node.id)
                if record is None:
                    continue
                expected = "image" if node.type in {"icon", "illustration", "image", "decoration", "background"} else node.type
                if record["type"] != expected:
                    errors.append(f"wrong final type: {node.id}")
                if node.type == "text" and "text" in node.semantic and record.get("text") != node.semantic["text"]:
                    errors.append(f"wrong final text: {node.id}")
    return {"valid": not errors, "errors": errors,
            "source_count": len(objects), "planned_count": len(graph.nodes),
            "scope": "three-way" if observed is not None else "source-to-plan"}
