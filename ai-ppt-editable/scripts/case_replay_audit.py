#!/usr/bin/env python3
"""Audit case-level native structure and merge topology.

This is deliberately a case replay gate rather than a generic unit test.  It
binds a final PPTX to its manifest, inspects the real package objects, reads
the OOXML merge topology, and checks that formal text is still native text.
"""
from __future__ import annotations

import argparse
import hashlib
import json
import zipfile
from pathlib import Path

from atomic_output import atomic_write_json


def digest(path: Path) -> str:
    value = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            value.update(chunk)
    return value.hexdigest()


def actual_kind(shape) -> str:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if getattr(shape, "has_table", False):
        return "editable_table"
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        return "native_group"
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.TEXT_BOX:
        return "native_text"
    if getattr(shape, "has_text_frame", False):
        return "native_text"
    return "native_shape"


def walk(shapes, parent=None):
    for shape in shapes:
        yield shape, parent
        if actual_kind(shape) == "native_group":
            yield from walk(shape.shapes, shape.name)


def merge_topology(shape) -> list[list[int]]:
    from pptx.oxml.ns import qn

    table = shape._element.find(f".//{qn('a:tbl')}")
    if table is None:
        return []
    rows = len(shape.table.rows)
    columns = len(shape.table.columns)
    found = []
    for row_index, row_xml in enumerate(table.findall(qn("a:tr"))):
        column_index = 0
        for cell_xml in row_xml.findall(qn("a:tc")):
            try:
                grid_span = max(1, int(cell_xml.get("gridSpan", "1")))
            except (TypeError, ValueError):
                grid_span = 1
            try:
                row_span = max(1, int(cell_xml.get("rowSpan", "1")))
            except (TypeError, ValueError):
                row_span = 1
            continuation = cell_xml.get("hMerge") == "1" or cell_xml.get("vMerge") == "1"
            if not continuation and (grid_span > 1 or row_span > 1):
                found.append([
                    row_index,
                    column_index,
                    min(rows - 1, row_index + row_span - 1),
                    min(columns - 1, column_index + grid_span - 1),
                ])
            column_index += grid_span
    return [list(item) for item in sorted({tuple(item) for item in found})]


def table_values(shape):
    return [[str(cell.text or "") for cell in row.cells] for row in shape.table.rows]


def normalize_merges(value):
    output = []
    for merge in value or []:
        if isinstance(merge, (list, tuple)) and len(merge) == 4:
            output.append([int(part) for part in merge])
    return sorted(output)


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("deck")
    parser.add_argument("--object-manifest", required=True)
    parser.add_argument("--report", required=True)
    args = parser.parse_args()
    deck_path = Path(args.deck).resolve()
    manifest_path = Path(args.object_manifest).resolve()
    errors = []
    warnings = []
    try:
        manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
        from pptx import Presentation

        prs = Presentation(str(deck_path))
    except Exception as exc:
        result = {"schema": "ai-ppt-plus/case-replay-object-audit/v1", "valid": False, "status": "blocked", "errors": [{"code": "case_replay_input_error", "message": f"{type(exc).__name__}: {exc}"}]}
        atomic_write_json(Path(args.report).resolve(), result)
        print(json.dumps(result, ensure_ascii=False))
        return 2

    slides = []
    total_tables = 0
    total_native_panels = 0
    total_formal_text = 0
    total_formal_text_native = 0
    total_a_tbl = 0
    whole_slide_pictures = []
    manifest_slides = manifest.get("slides", []) if isinstance(manifest, dict) else []
    for slide_no, slide in enumerate(prs.slides, 1):
        expected_slide = next((item for item in manifest_slides if isinstance(item, dict) and int(item.get("slide_no", 0)) == slide_no), {})
        expected_objects = expected_slide.get("objects", []) if isinstance(expected_slide, dict) else []
        expected_by_id = {str(item.get("object_id")): item for item in expected_objects if isinstance(item, dict) and item.get("object_id")}
        observed = {str(shape.name): (shape, parent) for shape, parent in walk(slide.shapes)}
        table_records = []
        panel_records = []
        text_records = []
        for object_id, expected in expected_by_id.items():
            item = observed.get(object_id)
            if item is None:
                errors.append({"code": "case_replay_object_missing", "slide": slide_no, "object_id": object_id})
                continue
            shape, parent = item
            kind = actual_kind(shape)
            expected_type = str(expected.get("object_type", ""))
            if expected_type == "editable_table":
                expected_merges = normalize_merges(expected.get("merges") or expected.get("merge_topology"))
                observed_merges = normalize_merges(merge_topology(shape)) if kind == "editable_table" else []
                record = {
                    "object_id": object_id,
                    "actual_type": kind,
                    "expected_merges": expected_merges,
                    "observed_merges": observed_merges,
                    "merge_count": len(observed_merges),
                    "rows": len(shape.table.rows) if kind == "editable_table" else 0,
                    "columns": len(shape.table.columns) if kind == "editable_table" else 0,
                    "values": table_values(shape) if kind == "editable_table" else [],
                    "rich_text_runs": sum(len(paragraph.runs) for row in shape.table.rows for cell in row.cells for paragraph in cell.text_frame.paragraphs) if kind == "editable_table" else 0,
                    "rich_text_cells": sum(1 for row in shape.table.rows for cell in row.cells if sum(len(paragraph.runs) for paragraph in cell.text_frame.paragraphs) > 1) if kind == "editable_table" else 0,
                }
                table_records.append(record)
                total_tables += int(kind == "editable_table")
                if kind != "editable_table":
                    errors.append({"code": "case_replay_table_not_native", "slide": slide_no, "object_id": object_id, "actual_type": kind})
                if expected_merges != observed_merges:
                    errors.append({"code": "case_replay_merge_topology_mismatch", "slide": slide_no, "object_id": object_id, "expected": expected_merges, "observed": observed_merges})
                if expected.get("rich_text_required") is True and record["rich_text_cells"] == 0:
                    errors.append({"code": "case_replay_rich_text_missing", "slide": slide_no, "object_id": object_id})
            if expected.get("native_required") is True or expected.get("role") in {"semantic-panel", "panel", "card", "framework"}:
                record = {"object_id": object_id, "actual_type": kind, "parent_group": parent}
                panel_records.append(record)
                if kind in {"native_shape", "native_group"}:
                    total_native_panels += 1
                else:
                    errors.append({"code": "case_replay_panel_not_native", "slide": slide_no, "object_id": object_id, "actual_type": kind})
            if expected_type == "editable_text":
                text = str(getattr(shape, "text", "") or "")
                expected_text = str(((expected.get("text_spec") or {}).get("content") if isinstance(expected.get("text_spec"), dict) else expected.get("text", "")) or "")
                record = {"object_id": object_id, "actual_type": kind, "native": kind == "native_text", "expected": expected_text, "observed": text, "run_count": sum(len(paragraph.runs) for paragraph in shape.text_frame.paragraphs) if getattr(shape, "has_text_frame", False) else 0}
                text_records.append(record)
                total_formal_text += 1
                total_formal_text_native += int(kind == "native_text")
                if kind != "native_text":
                    errors.append({"code": "case_replay_formal_text_not_native", "slide": slide_no, "object_id": object_id, "actual_type": kind})
                elif expected_text != text:
                    errors.append({"code": "case_replay_formal_text_mismatch", "slide": slide_no, "object_id": object_id, "expected": expected_text, "observed": text})
        for shape, parent in walk(slide.shapes):
            if actual_kind(shape) == "picture" and shape.width >= prs.slide_width * 0.95 and shape.height >= prs.slide_height * 0.95:
                whole_slide_pictures.append({"slide": slide_no, "name": shape.name})
        slides.append({"slide": slide_no, "tables": table_records, "panels": panel_records, "formal_text": text_records, "top_level_objects": sum(1 for _, parent in walk(slide.shapes) if parent is None)})
    with zipfile.ZipFile(deck_path) as package:
        # Count the actual table element, not its child elements (a:tblPr
        # and a:tblGrid also begin with the bytes ``<a:tbl``).
        total_a_tbl = sum(package.read(name).count(b"<a:tbl>") for name in package.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml"))
    for table in [record for slide in slides for record in slide["tables"]]:
        if table["actual_type"] == "editable_table" and table["rows"] == 0:
            errors.append({"code": "case_replay_empty_table", "object_id": table["object_id"]})
    result = {
        "schema": "ai-ppt-plus/case-replay-object-audit/v1",
        "valid": not errors,
        "status": "passed" if not errors else "blocked",
        "deck": str(deck_path),
        "deck_sha256": digest(deck_path),
        "object_manifest": str(manifest_path),
        "object_manifest_sha256": digest(manifest_path),
        "native_table_count": total_tables,
        "native_panel_count": total_native_panels,
        "formal_text_count": total_formal_text,
        "formal_text_native_count": total_formal_text_native,
        "a_tbl_count": total_a_tbl,
        "whole_slide_pictures": whole_slide_pictures,
        "slides": slides,
        "errors": errors,
        "warnings": warnings,
        "human_visual_review_required": True,
    }
    atomic_write_json(Path(args.report).resolve(), result)
    print(json.dumps(result, ensure_ascii=False))
    return 0 if result["valid"] else 2


if __name__ == "__main__":
    raise SystemExit(main())
