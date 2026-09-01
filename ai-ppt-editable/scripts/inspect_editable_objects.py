#!/usr/bin/env python3
"""Reverse-audit final PPTX objects against a slide-object manifest."""
from __future__ import annotations

import argparse
import json
import math
from pathlib import Path
from typing import Any
from atomic_output import atomic_write_json


def _iter_shapes(shapes):
    """Yield top-level shapes and nested group children for identity checks."""
    for shape in shapes:
        yield shape
        if str(getattr(shape, "shape_type", "")).startswith("GROUP") or getattr(shape, "shapes", None) is not None:
            try:
                yield from _iter_shapes(shape.shapes)
            except (AttributeError, TypeError):
                continue


def _actual_object_type(shape) -> str:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if getattr(shape, "has_table", False):
        return "editable_table"
    if getattr(shape, "has_chart", False):
        return "editable_chart"
    shape_type = getattr(shape, "shape_type", None)
    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if shape_type == MSO_SHAPE_TYPE.GROUP:
        return "native_group"
    if shape_type in {MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.LINE}:
        return "native_shape"
    if shape_type == MSO_SHAPE_TYPE.TEXT_BOX and getattr(shape, "has_text_frame", False):
        return "native_text"
    return "native_text" if getattr(shape, "has_text_frame", False) else "other"


def _type_matches(expected: str, observed: str) -> bool:
    allowed = {
        "editable_text": {"native_text"},
        "native_shape": {"native_shape"},
        "native_group": {"native_group"},
        "editable_table": {"editable_table"},
        "editable_chart": {"editable_chart"},
        "independent_image": {"picture"},
        "traceable_static_graphic": {"picture", "native_shape"},
        "extracted_icon": {"picture", "native_shape"},
        "editable_vector": {"picture", "native_shape"},
    }
    return observed in allowed.get(expected, {expected})


def _manifest_geometry(obj: dict[str, Any], contract: dict[str, Any]) -> dict[str, float | str] | None:
    raw = obj.get("declared_geometry")
    if not isinstance(raw, dict):
        raw = obj.get("geometry") if isinstance(obj.get("geometry"), dict) else None
    if not isinstance(raw, dict):
        return None
    values = {key: raw.get(key) for key in ("x", "y", "w", "h")}
    if any(value is None for value in values.values()):
        return None
    try:
        result: dict[str, float | str] = {key: float(value) for key, value in values.items()}
    except (TypeError, ValueError):
        return None
    result["coordinate_space"] = str(raw.get("coordinate_space") or contract.get("coordinate_space") or "fraction")
    if result["coordinate_space"] == "px":
        try:
            result["reference_width"] = float(raw.get("reference_width") or contract.get("reference_width"))
            result["reference_height"] = float(raw.get("reference_height") or contract.get("reference_height"))
        except (TypeError, ValueError):
            return None
    return result


def _normalized_geometry(shape, prs, expected: dict[str, float | str]) -> tuple[float, float, float, float]:
    coordinate_space = str(expected.get("coordinate_space", "fraction"))
    width, height = float(prs.slide_width), float(prs.slide_height)
    if coordinate_space == "px":
        reference_width = float(expected["reference_width"])
        reference_height = float(expected["reference_height"])
        if reference_width <= 0 or reference_height <= 0:
            raise ValueError("pixel geometry requires positive reference dimensions")
        target = (float(expected["x"]) / reference_width, float(expected["y"]) / reference_height, float(expected["w"]) / reference_width, float(expected["h"]) / reference_height)
    else:
        target = tuple(float(expected[key]) for key in ("x", "y", "w", "h"))
    return target


def _observed_geometry(shape, prs) -> tuple[float, float, float, float]:
    return (float(shape.left) / float(prs.slide_width), float(shape.top) / float(prs.slide_height), float(shape.width) / float(prs.slide_width), float(shape.height) / float(prs.slide_height))


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("deck")
    ap.add_argument("--object-manifest")
    ap.add_argument("--require-independent-panels", action="store_true")
    ap.add_argument("--require-types", action="store_true", help="require manifest object_type to match the PPTX native object type")
    ap.add_argument("--require-geometry", action="store_true", help="require declared geometry and compare it with the rendered shape box")
    ap.add_argument("--geometry-tolerance", type=float, default=0.02, help="normalized x/y/w/h tolerance for object geometry")
    ap.add_argument("--report")
    args = ap.parse_args()
    from pptx import Presentation
    prs = Presentation(args.deck)
    expected_by_slide: dict[int, dict[str, dict]] = {}
    if args.object_manifest:
        data = json.loads(Path(args.object_manifest).read_text(encoding="utf-8"))
        for slide_index, manifest_slide in enumerate(data.get("slides", []), 1):
            slide_no = int(manifest_slide.get("slide_no", slide_index))
            page_objects = expected_by_slide.setdefault(slide_no, {})
            for obj in manifest_slide.get("objects", []):
                object_id = obj.get("object_id")
                if object_id in page_objects:
                    raise ValueError(f"duplicate object_id on slide {slide_no}: {object_id}")
                page_objects[object_id] = obj
    errors, warnings, slides = [], [], []
    geometry_checked = 0
    geometry_mismatches = 0
    type_checked = 0
    type_mismatches = 0
    try:
        if not math.isfinite(float(args.geometry_tolerance)) or float(args.geometry_tolerance) < 0:
            raise ValueError
    except (TypeError, ValueError):
        raise ValueError("--geometry-tolerance must be a non-negative finite number")
    manifest_contract: dict[str, Any] = {}
    if args.object_manifest:
        manifest_contract = data.get("geometry_contract", {}) if isinstance(data, dict) else {}
    for si, slide in enumerate(prs.slides, 1):
        shapes = []
        for shape in slide.shapes:
            kind = str(getattr(shape, "shape_type", "unknown"))
            text = getattr(shape, "text", "") if hasattr(shape, "text") else ""
            shapes.append({"name": shape.name, "type": kind, "has_text": bool(text.strip()), "text": text[:160], "left": shape.left, "top": shape.top, "width": shape.width, "height": shape.height})
        slides.append({"slide": si, "shape_count": len(shapes), "shapes": shapes})
        names = {s["name"] for s in shapes}
        all_shapes = list(_iter_shapes(slide.shapes))
        by_name: dict[str, list[Any]] = {}
        for shape in all_shapes:
            by_name.setdefault(str(getattr(shape, "name", "")), []).append(shape)
        expected = expected_by_slide.get(si, {})
        for oid, obj in expected.items():
            matches = by_name.get(str(oid), [])
            if not matches:
                errors.append({"code": "manifest_object_missing_in_pptx", "slide": si, "object_id": oid, "role": obj.get("role")})
                continue
            if args.require_types and obj.get("object_type"):
                type_checked += 1
                observed_type = _actual_object_type(matches[0])
                if not _type_matches(str(obj["object_type"]), observed_type):
                    type_mismatches += 1
                    errors.append({"code": "manifest_object_type_mismatch", "slide": si, "object_id": oid, "expected": obj.get("object_type"), "observed": observed_type})
            geometry = _manifest_geometry(obj, manifest_contract)
            if args.require_geometry and geometry is None:
                errors.append({"code": "manifest_object_geometry_missing", "slide": si, "object_id": oid, "role": obj.get("role")})
            elif geometry is not None:
                geometry_checked += 1
                try:
                    expected_box = _normalized_geometry(matches[0], prs, geometry)
                    observed_box = _observed_geometry(matches[0], prs)
                    deltas = [abs(a - b) for a, b in zip(expected_box, observed_box)]
                    if max(deltas) > float(args.geometry_tolerance):
                        geometry_mismatches += 1
                        errors.append({"code": "manifest_object_geometry_mismatch", "slide": si, "object_id": oid, "expected": list(expected_box), "observed": list(observed_box), "deltas": deltas, "tolerance": float(args.geometry_tolerance)})
                except (KeyError, TypeError, ValueError, ZeroDivisionError) as exc:
                    geometry_mismatches += 1
                    errors.append({"code": "manifest_object_geometry_invalid", "slide": si, "object_id": oid, "message": f"{type(exc).__name__}: {exc}"})
        if args.require_independent_panels:
            panel_ids = [oid for oid, obj in expected.items() if obj.get("role") in {"semantic-panel", "panel", "frame-panel"}]
            for oid in panel_ids:
                matches = [s for s in shapes if s["name"] == oid]
                if len(matches) != 1:
                    errors.append({"code": "panel_shape_count_not_one", "slide": si, "object_id": oid, "observed": len(matches)})
    whole_slide_images = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if str(getattr(shape, "shape_type", "")) == "PICTURE" and shape.width >= prs.slide_width * .95 and shape.height >= prs.slide_height * .95:
                whole_slide_images.append(shape.name)
    expected_object_count = sum(len(objects) for objects in expected_by_slide.values())
    if whole_slide_images and expected_object_count:
        warnings.append({"code": "whole_slide_picture_present", "names": whole_slide_images})
    result = {"schema": "ai-ppt-plus/editable-object-audit/v1", "valid": not errors, "deck": str(Path(args.deck).resolve()), "slides": slides, "expected_object_count": expected_object_count, "observed_shape_count": sum(s["shape_count"] for s in slides), "whole_slide_pictures": whole_slide_images, "geometry_checked": geometry_checked, "geometry_mismatch_count": geometry_mismatches, "type_checked": type_checked, "type_mismatch_count": type_mismatches, "errors": errors, "warnings": warnings, "human_visual_review_required": True}
    if args.report:
        out = Path(args.report); out.parent.mkdir(parents=True, exist_ok=True)
        atomic_write_json(out.resolve(), result)
    print(json.dumps(result, ensure_ascii=False))
    return 0 if result["valid"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
