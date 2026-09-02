#!/usr/bin/env python3
"""Verify that semantic panels, tables and text are native PPTX objects."""
from __future__ import annotations

import argparse
import json
from pathlib import Path

from atomic_output import atomic_write_json


PANEL_ROLES = {"semantic-panel", "panel", "frame-panel", "card", "card-frame", "framework"}
TABLE_ROLES = {"table", "data-table", "editable-table", "data-grid"}
NATIVE_PANEL_KINDS = {"native_shape", "native_group"}
FORMAL_TEXT_ROLES = {"formal-text", "title", "body", "label", "caption", "callout"}


def _shape_kind(shape) -> str:
    try:
        if bool(getattr(shape, "has_table", False)):
            return "editable_table"
    except Exception:
        pass
    raw = str(getattr(shape, "shape_type", "")).upper()
    if "GROUP" in raw:
        return "native_group"
    if "CHART" in raw:
        return "editable_chart"
    if "PICTURE" in raw:
        return "picture"
    if "CONNECTOR" in raw:
        return "native_line"
    try:
        if bool(getattr(shape, "has_text_frame", False)) and str(getattr(shape, "text", "")).strip():
            return "editable_text"
    except Exception:
        pass
    return "native_shape"


def _text(shape) -> str:
    try:
        return str(getattr(shape, "text", "") or "")
    except Exception:
        return ""


def _record(shape, kind: str, slide_width: int, slide_height: int, parent: str | None) -> dict:
    record = {
        "name": str(getattr(shape, "name", "")),
        "kind": kind,
        "parent": parent,
        "left": int(getattr(shape, "left", 0) or 0),
        "top": int(getattr(shape, "top", 0) or 0),
        "width": int(getattr(shape, "width", 0) or 0),
        "height": int(getattr(shape, "height", 0) or 0),
        "normalized_box": [
            round(float(getattr(shape, "left", 0) or 0) / slide_width, 6) if slide_width else 0,
            round(float(getattr(shape, "top", 0) or 0) / slide_height, 6) if slide_height else 0,
            round(float(getattr(shape, "width", 0) or 0) / slide_width, 6) if slide_width else 0,
            round(float(getattr(shape, "height", 0) or 0) / slide_height, 6) if slide_height else 0,
        ],
        "text": _text(shape)[:240],
    }
    if kind == "editable_table":
        table = getattr(shape, "table", None)
        if table is not None:
            record["rows"] = len(table.rows)
            record["columns"] = len(table.columns)
            record["cell_text"] = [
                [str(table.cell(row, column).text or "") for column in range(len(table.columns))]
                for row in range(len(table.rows))
            ]
    return record


def _walk_shapes(shapes, slide_width: int, slide_height: int, index: dict[str, list[dict]], parent: str | None = None) -> list[dict]:
    flat: list[dict] = []
    for shape in shapes:
        kind = _shape_kind(shape)
        current = _record(shape, kind, slide_width, slide_height, parent)
        flat.append(current)
        index.setdefault(current["name"], []).append(current)
        if kind == "native_group":
            try:
                flat.extend(_walk_shapes(shape.shapes, slide_width, slide_height, index, current["name"]))
            except Exception:
                current["children_unreadable"] = True
    return flat


def _load_manifest(path: Path, errors: list[dict]) -> dict:
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except Exception as exc:
        errors.append({"severity": "blocker", "code": "object_manifest_unreadable", "message": f"{type(exc).__name__}: {exc}"})
        return {}
    if not isinstance(data, dict):
        errors.append({"severity": "blocker", "code": "object_manifest_not_object"})
        return {}
    return data


def _matches(index: dict[str, list[dict]], object_id) -> list[dict]:
    return index.get(str(object_id), [])


def validate_native_editability(
    deck_path: Path,
    manifest_path: Path | None,
    *,
    require_native_panels: bool = False,
    require_native_tables: bool = False,
    forbid_whole_slide_pictures: bool = False,
    require_complete_manifest: bool = False,
) -> dict:
    errors: list[dict] = []
    manifest = {}
    if manifest_path is not None:
        manifest = _load_manifest(manifest_path, errors)
    elif require_native_panels or require_native_tables or require_complete_manifest:
        errors.append({"severity": "blocker", "code": "object_manifest_missing"})

    try:
        from pptx import Presentation
        prs = Presentation(str(deck_path))
    except Exception as exc:
        errors.append({"severity": "blocker", "code": "pptx_unreadable", "message": f"{type(exc).__name__}: {exc}"})
        return {
            "schema": "ai-ppt-plus/native-object-validation/v1",
            "valid": False,
            "status": "blocked",
            "deck": str(deck_path.resolve()),
            "errors": errors,
            "slides": [],
        }

    slide_objects: dict[int, dict[str, dict]] = {}
    manifest_ids: dict[int, set[str]] = {}
    for index, slide_spec in enumerate(manifest.get("slides", []) if isinstance(manifest, dict) else [], 1):
        if not isinstance(slide_spec, dict):
            continue
        slide_no = int(slide_spec.get("slide_no", index))
        objects = slide_spec.get("objects", [])
        if not isinstance(objects, list):
            errors.append({"severity": "blocker", "code": "manifest_objects_invalid", "slide": slide_no})
            continue
        slide_objects[slide_no] = {}
        manifest_ids[slide_no] = set()
        for obj in objects:
            if not isinstance(obj, dict) or not obj.get("object_id"):
                errors.append({"severity": "blocker", "code": "manifest_object_invalid", "slide": slide_no})
                continue
            object_id = str(obj["object_id"])
            if object_id in slide_objects[slide_no]:
                errors.append({"severity": "blocker", "code": "manifest_object_duplicate", "slide": slide_no, "object_id": object_id})
            slide_objects[slide_no][object_id] = obj
            manifest_ids[slide_no].add(object_id)

    slides = []
    whole_slide_pictures: list[dict] = []
    native_table_count = 0
    native_panel_count = 0
    for slide_no, slide in enumerate(prs.slides, 1):
        index: dict[str, list[dict]] = {}
        top_level = _walk_shapes(slide.shapes, int(prs.slide_width), int(prs.slide_height), index)
        tables = [record for record in index.values() for record in record if record["kind"] == "editable_table"]
        native_table_count += len(tables)
        for record in top_level:
            if record["kind"] == "picture" and record["width"] >= prs.slide_width * 0.95 and record["height"] >= prs.slide_height * 0.95:
                whole_slide_pictures.append({"slide": slide_no, "name": record["name"], "record": record})
        expected = slide_objects.get(slide_no, {})
        background_ids = {
            object_id for object_id, obj in expected.items()
            if obj.get("role") == "background" or object_id.casefold() == "background"
        }
        if forbid_whole_slide_pictures:
            for picture in whole_slide_pictures:
                if picture["slide"] != slide_no:
                    continue
                if picture["name"] not in background_ids:
                    errors.append({"severity": "blocker", "code": "whole_slide_picture_forbidden", "slide": slide_no, "name": picture["name"]})
        slide_errors_before = len(errors)
        for object_id, obj in expected.items():
            matches = _matches(index, object_id)
            if len(matches) != 1:
                errors.append({"severity": "blocker", "code": "manifest_object_count_mismatch", "slide": slide_no, "object_id": object_id, "observed": len(matches)})
                continue
            actual = matches[0]
            role = str(obj.get("role") or "")
            object_type = str(obj.get("object_type") or "")
            if role in FORMAL_TEXT_ROLES or object_type == "editable_text":
                if actual["kind"] != "editable_text":
                    errors.append({"severity": "blocker", "code": "formal_text_not_native", "slide": slide_no, "object_id": object_id, "actual_kind": actual["kind"]})
            is_table = role in TABLE_ROLES or object_type == "editable_table"
            if is_table:
                if actual["kind"] != "editable_table":
                    errors.append({"severity": "blocker", "code": "table_not_native", "slide": slide_no, "object_id": object_id, "actual_kind": actual["kind"]})
                elif require_native_tables:
                    native_table_count += 0
            is_panel = role in PANEL_ROLES or bool(obj.get("native_required"))
            if is_panel and (require_native_panels or obj.get("native_required") is True):
                if actual["kind"] not in NATIVE_PANEL_KINDS:
                    errors.append({"severity": "blocker", "code": "panel_not_native", "slide": slide_no, "object_id": object_id, "actual_kind": actual["kind"]})
                else:
                    native_panel_count += 1
            if require_complete_manifest and actual["name"] not in manifest_ids.get(slide_no, set()):
                errors.append({"severity": "blocker", "code": "undeclared_object", "slide": slide_no, "name": actual["name"]})
            if require_native_panels and role in {"frame", "framework", "whole-frame", "skeleton"} and actual["kind"] == "picture":
                errors.append({"severity": "blocker", "code": "semantic_frame_not_native", "slide": slide_no, "object_id": object_id})
        if require_native_panels:
            for record in top_level:
                if record["name"].casefold() in {"frame", "framework", "whole-frame", "skeleton"} and record["kind"] == "picture":
                    errors.append({"severity": "blocker", "code": "unmanifested_semantic_frame_picture", "slide": slide_no, "name": record["name"]})
        slides.append({
            "slide": slide_no,
            "top_level_shape_count": len(top_level),
            "object_count_recursive": len(index),
            "native_tables": sum(1 for values in index.values() for record in values if record["kind"] == "editable_table"),
            "objects": [record for values in index.values() for record in values],
            "errors_added": len(errors) - slide_errors_before,
        })

    result = {
        "schema": "ai-ppt-plus/native-object-validation/v1",
        "valid": not errors,
        "status": "passed" if not errors else "blocked",
        "deck": str(deck_path.resolve()),
        "object_manifest": str(manifest_path.resolve()) if manifest_path else None,
        "requirements": {
            "require_native_panels": require_native_panels,
            "require_native_tables": require_native_tables,
            "forbid_whole_slide_pictures": forbid_whole_slide_pictures,
            "require_complete_manifest": require_complete_manifest,
        },
        "native_table_count": native_table_count,
        "native_panel_count": native_panel_count,
        "whole_slide_pictures": whole_slide_pictures,
        "slides": slides,
        "errors": errors,
        "human_visual_review_required": True,
    }
    return result


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("deck")
    parser.add_argument("--object-manifest")
    parser.add_argument("--require-native-structure", action="store_true")
    parser.add_argument("--require-native-panels", action="store_true")
    parser.add_argument("--require-native-tables", action="store_true")
    parser.add_argument("--forbid-whole-slide-pictures", action="store_true")
    parser.add_argument("--require-complete-manifest", action="store_true")
    parser.add_argument("--report")
    args = parser.parse_args()
    result = validate_native_editability(
        Path(args.deck).resolve(),
        Path(args.object_manifest).resolve() if args.object_manifest else None,
        require_native_panels=bool(args.require_native_structure or args.require_native_panels),
        require_native_tables=bool(args.require_native_structure or args.require_native_tables),
        forbid_whole_slide_pictures=bool(args.forbid_whole_slide_pictures or args.require_native_structure),
        require_complete_manifest=args.require_complete_manifest,
    )
    if args.report:
        atomic_write_json(Path(args.report).resolve(), result)
    print(json.dumps(result, ensure_ascii=False))
    return 0 if result["valid"] else 2


if __name__ == "__main__":
    raise SystemExit(main())
