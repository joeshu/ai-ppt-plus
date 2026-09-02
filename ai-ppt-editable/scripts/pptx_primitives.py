#!/usr/bin/env python3
"""Native PowerPoint primitive writers.

This module owns text, shape, group, table and chart creation. It deliberately
does not know about the command line, component libraries, previews or font
embedding; the authoring backend supplies the slide and coordinate context.
"""
from __future__ import annotations

import sys
import math

from pptx.util import Pt


def _die(message: str, code: int = 2):
    print(f"Error: {message}", file=sys.stderr)
    raise SystemExit(code)


def text_size_pt(item: dict, slide_height_pt: float, reference_height: float, default=None) -> float:
    """Convert ratio, pixel or absolute text sizes to points."""
    if item.get("size") is not None:
        value = float(item["size"])
    elif item.get("size_ratio") is not None:
        value = float(item["size_ratio"]) * slide_height_pt
    elif item.get("size_pct") is not None:
        value = float(item["size_pct"]) / 100.0 * slide_height_pt
    elif item.get("size_px") is not None:
        if not reference_height:
            _die("size_px requires a positive reference height")
        value = float(item["size_px"]) * slide_height_pt / reference_height
    else:
        value = float(default) if default is not None else 18.0
    if not math.isfinite(value) or value <= 0:
        _die(f"font size must be a positive finite number: {value}")
    return value


def _set_run_fonts(run, name: str):
    """Set latin, east-asian and complex-script typefaces."""
    from pptx.oxml.ns import qn

    run.font.name = name
    run_properties = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        element = run_properties.find(qn(tag))
        if element is None:
            element = run_properties.makeelement(qn(tag), {})
            run_properties.append(element)
        element.set("typeface", name)


def _set_run_alpha(run, opacity: float):
    from pptx.oxml.ns import qn

    opacity = max(0.0, min(1.0, float(opacity)))
    run_properties = run._r.get_or_add_rPr()
    solid = run_properties.find(qn("a:solidFill"))
    if solid is None:
        return
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:
        return
    for old in srgb.findall(qn("a:alpha")):
        srgb.remove(old)
    srgb.append(srgb.makeelement(qn("a:alpha"), {"val": str(int(opacity * 100000))}))


def _normalized_hex(value: str) -> str:
    raw = str(value).strip().lstrip("#")
    if len(raw) == 3:
        raw = "".join(char * 2 for char in raw)
    if len(raw) != 6:
        raise ValueError(f"invalid color {value!r}; expected #RRGGBB")
    try:
        int(raw, 16)
    except ValueError as exc:
        raise ValueError(f"invalid color {value!r}; expected #RRGGBB") from exc
    return raw.upper()


def _hex_to_rgb(value: str):
    from pptx.dml.color import RGBColor

    raw = _normalized_hex(value)
    return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))


def _hex_to_tuple(value: str):
    raw = _normalized_hex(value)
    return (int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))


def _hex_to_rgba(value: str, opacity: float = 1.0):
    red, green, blue = _hex_to_tuple(value)
    return red, green, blue, int(max(0.0, min(1.0, float(opacity))) * 255)


def _set_fill_alpha(shape, opacity: float):
    """Add an alpha transform to a solid shape fill."""
    from pptx.oxml.ns import qn

    shape_properties = shape._element.spPr
    solid = shape_properties.find(qn("a:solidFill"))
    if solid is None:
        return
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:
        return
    for old in srgb.findall(qn("a:alpha")):
        srgb.remove(old)
    alpha = int(max(0.0, min(1.0, opacity)) * 100000)
    srgb.append(srgb.makeelement(qn("a:alpha"), {"val": str(alpha)}))


def _add_outer_shadow(shape, blur_pt=6.0, dist_pt=3.0, alpha=0.35):
    """Best-effort soft drop shadow for card-like shapes."""
    from pptx.oxml.ns import qn
    from pptx.util import Pt

    shape_properties = shape._element.spPr
    for old in shape_properties.findall(qn("a:effectLst")):
        shape_properties.remove(old)
    effects = shape_properties.makeelement(qn("a:effectLst"), {})
    shadow = effects.makeelement(qn("a:outerShdw"), {
        "blurRad": str(int(Pt(blur_pt))),
        "dist": str(int(Pt(dist_pt))),
        "dir": "5400000",
        "rotWithShape": "0",
    })
    color = shadow.makeelement(qn("a:srgbClr"), {"val": "000000"})
    color.append(color.makeelement(qn("a:alpha"), {"val": str(int(alpha * 100000))}))
    shadow.append(color)
    effects.append(shadow)
    shape_properties.append(effects)


def _set_gradient_fill(shape, gradient: dict):
    """Write a deterministic PresentationML linear gradient fill."""
    from pptx.oxml.ns import qn

    shape_properties = shape._element.spPr
    for tag in ("a:solidFill", "a:noFill", "a:gradFill"):
        old = shape_properties.find(qn(tag))
        if old is not None:
            shape_properties.remove(old)
    stops = gradient.get("stops", []) if isinstance(gradient, dict) else []
    if len(stops) < 2:
        _die("gradient fill requires at least two color stops")
    gradient_fill = shape_properties.makeelement(qn("a:gradFill"), {"rotWithShape": "1"})
    stop_list = gradient_fill.makeelement(qn("a:gsLst"), {})
    for stop in stops:
        if not isinstance(stop, dict) or not isinstance(stop.get("color"), str):
            _die("gradient stops require color fields")
        raw_position = float(stop.get("position", stop.get("pos", 0)))
        position = max(0, min(100000, int(raw_position * (100000 if 0 <= raw_position <= 1 else 1000))))
        gradient_stop = stop_list.makeelement(qn("a:gs"), {"pos": str(position)})
        color = gradient_stop.makeelement(qn("a:srgbClr"), {"val": _normalized_hex(stop["color"])})
        if stop.get("opacity") is not None:
            alpha = int(max(0, min(1, float(stop["opacity"]))) * 100000)
            color.append(color.makeelement(qn("a:alpha"), {"val": str(alpha)}))
        gradient_stop.append(color)
        stop_list.append(gradient_stop)
    gradient_fill.append(stop_list)
    angle = float(gradient.get("angle", 0))
    gradient_fill.append(gradient_fill.makeelement(qn("a:lin"), {"ang": str(int(angle * 60000) % 21600000), "scaled": "1"}))
    line = shape_properties.find(qn("a:ln"))
    if line is not None:
        shape_properties.insert(list(shape_properties).index(line), gradient_fill)
    else:
        shape_properties.append(gradient_fill)


def apply_shape_fill(shape, spec: dict):
    gradient = spec.get("gradient") or spec.get("fill_gradient")
    if gradient:
        _set_gradient_fill(shape, gradient)
    elif spec.get("fill"):
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(spec["fill"])
        if spec.get("opacity") is not None:
            _set_fill_alpha(shape, float(spec["opacity"]))
    else:
        shape.fill.background()


def set_alt_text(shape, text: str | None):
    if not text:
        return
    element = shape._element
    for attribute in ("nvSpPr", "nvPicPr", "nvGrpSpPr", "nvGraphicFramePr"):
        container = getattr(element, attribute, None)
        if container is not None and getattr(container, "cNvPr", None) is not None:
            container.cNvPr.set("descr", str(text))
            return


def shape_map():
    from pptx.enum.shapes import MSO_SHAPE

    return {
        "rect": MSO_SHAPE.RECTANGLE,
        "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "oval": MSO_SHAPE.OVAL,
        "ellipse": MSO_SHAPE.OVAL,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "chevron": MSO_SHAPE.CHEVRON,
        "right_arrow": MSO_SHAPE.RIGHT_ARROW,
        "left_arrow": MSO_SHAPE.LEFT_ARROW,
        "up_arrow": MSO_SHAPE.UP_ARROW,
        "down_arrow": MSO_SHAPE.DOWN_ARROW,
        "pentagon": MSO_SHAPE.REGULAR_PENTAGON,
        "hexagon": MSO_SHAPE.HEXAGON,
        "parallelogram": MSO_SHAPE.PARALLELOGRAM,
        "trapezoid": MSO_SHAPE.TRAPEZOID,
        "diamond": MSO_SHAPE.DIAMOND,
    }


def chart_map():
    from pptx.enum.chart import XL_CHART_TYPE

    return {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
    }


def add_shapes(slide, specs: list[dict], deck: dict, ref_w: float, ref_h: float, sw_emu: int, sh_emu: int):
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.util import Emu, Pt

    shape_types = shape_map()
    for index, spec in enumerate(specs, 1):
        fx = _frac(deck, spec, "x", ref_w)
        fy = _frac(deck, spec, "y", ref_h)
        if deck.get("strict_input") and "type" not in spec:
            _die(f"shape {index} requires an explicit type in strict input mode")
        shape_type = spec.get("type", "rounded_rect")
        if shape_type == "line":
            # Lines are not boxes: a descending segment needs a negative
            # delta, which cannot be represented by the historical positive
            # w/h contract.  Accept explicit x2/y2 endpoints while retaining
            # the old x/y/w/h form for backwards compatibility.
            if ("x2" in spec) != ("y2" in spec):
                _die(f"line shape {index} requires both x2 and y2 when either endpoint is provided")
            if "x2" in spec:
                ex = _frac(deck, spec, "x2", ref_w)
                ey = _frac(deck, spec, "y2", ref_h)
                if not spec.get("allow_bleed") is True and deck.get("strict_input"):
                    if any(value < 0 or value > 1 for value in (fx, fy, ex, ey)):
                        _die(f"line shape {index} endpoints must stay within the slide in strict input mode")
                x1, y1, x2, y2 = fx, fy, ex, ey
            else:
                fw = _frac(deck, spec, "w", ref_w)
                fh = _frac(deck, spec, "h", ref_h)
                _validate_box(deck, spec, fx, fy, fw, fh, f"shape {index}")
                x1, y1, x2, y2 = fx, fy, fx + fw, fy + fh
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Emu(int(x1 * sw_emu)),
                Emu(int(y1 * sh_emu)),
                Emu(int(x2 * sw_emu)),
                Emu(int(y2 * sh_emu)),
            )
            connector.line.color.rgb = _hex_to_rgb(spec.get("line", spec.get("fill", "#FFFFFF")))
            connector.line.width = Pt(float(spec.get("line_width", 1.5)))
            connector.name = str(spec.get("object_id") or spec.get("name") or f"shape-{index:02d}")
            continue

        fw = _frac(deck, spec, "w", ref_w)
        fh = _frac(deck, spec, "h", ref_h)
        _validate_box(deck, spec, fx, fy, fw, fh, f"shape {index}")
        left, top = Emu(int(fx * sw_emu)), Emu(int(fy * sh_emu))
        width, height = Emu(int(fw * sw_emu)), Emu(int(fh * sh_emu))
        if shape_type not in shape_types:
            _die(f"unsupported shape type: {shape_type}")
        shape = slide.shapes.add_shape(shape_types[shape_type], left, top, width, height)
        shape.name = str(spec.get("object_id") or spec.get("name") or f"shape-{index:02d}")
        if shape_type == "rounded_rect" and "radius" in spec:
            try:
                shape.adjustments[0] = float(spec["radius"])
            except (TypeError, ValueError, IndexError) as exc:
                _die(f"rounded_rect radius is invalid: {exc}")
        apply_shape_fill(shape, spec)
        if spec.get("line"):
            shape.line.color.rgb = _hex_to_rgb(spec["line"])
            shape.line.width = Pt(float(spec.get("line_width", 1.0)))
        else:
            shape.line.fill.background()
        if spec.get("shadow"):
            _add_outer_shadow(shape)
        else:
            shape.shadow.inherit = False
        if spec.get("rotation"):
            shape.rotation = float(spec["rotation"])
        set_alt_text(shape, spec.get("alt_text"))


def add_groups(slide, specs: list[dict], deck: dict, ref_w: float, ref_h: float, sw_emu: int, sh_emu: int):
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.util import Emu, Pt

    shape_types = shape_map()
    for group_index, group_spec in enumerate(specs, 1):
        group = slide.shapes.add_group_shape()
        group.name = str(group_spec.get("object_id") or group_spec.get("name") or f"group-{group_index:02d}")
        local = group_spec.get("children_coordinate_space") == "local"
        has_box = all(key in group_spec for key in ("x", "y", "w", "h"))
        gx = _frac(deck, group_spec, "x", ref_w) if has_box else 0
        gy = _frac(deck, group_spec, "y", ref_h) if has_box else 0
        gw = _frac(deck, group_spec, "w", ref_w) if has_box else 1
        gh = _frac(deck, group_spec, "h", ref_h) if has_box else 1
        if has_box:
            _validate_box(deck, group_spec, gx, gy, gw, gh, f"group {group_index}")
        for child_index, raw_child in enumerate(group_spec.get("children", []), 1):
            if not isinstance(raw_child, dict):
                if deck.get("strict_input"):
                    _die(f"group {group_index} child {child_index} must be an object")
                continue
            child = dict(raw_child)
            if local:
                child["x"] = gx + _frac(deck, child, "x", ref_w) * gw
                child["y"] = gy + _frac(deck, child, "y", ref_h) * gh
                child["w"] = _frac(deck, child, "w", ref_w) * gw
                child["h"] = _frac(deck, child, "h", ref_h) * gh
                child_deck = dict(deck, units="fraction")
            else:
                child_deck = deck
            cfx = _frac(child_deck, child, "x", ref_w)
            cfy = _frac(child_deck, child, "y", ref_h)
            child_name = str(child.get("object_id") or child.get("name") or f"{group.name}-child-{child_index:02d}")
            if child.get("type", "rounded_rect") == "line":
                if ("x2" in child) != ("y2" in child):
                    _die(f"group {group_index} line child {child_index} requires both x2 and y2 when either endpoint is provided")
                if "x2" in child:
                    cex = _frac(child_deck, child, "x2", ref_w)
                    cey = _frac(child_deck, child, "y2", ref_h)
                    if not child.get("allow_bleed") is True and deck.get("strict_input"):
                        if any(value < 0 or value > 1 for value in (cfx, cfy, cex, cey)):
                            _die(f"group {group_index} line child {child_index} endpoints must stay within the slide in strict input mode")
                    cx1, cy1, cx2, cy2 = cfx, cfy, cex, cey
                else:
                    cfw = _frac(child_deck, child, "w", ref_w)
                    cfh = _frac(child_deck, child, "h", ref_h)
                    _validate_box(deck, child, cfx, cfy, cfw, cfh, f"group {group_index} child {child_index}")
                    cx1, cy1, cx2, cy2 = cfx, cfy, cfx + cfw, cfy + cfh
                connector = group.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Emu(int(cx1 * sw_emu)),
                    Emu(int(cy1 * sh_emu)),
                    Emu(int(cx2 * sw_emu)),
                    Emu(int(cy2 * sh_emu)),
                )
                connector.line.color.rgb = _hex_to_rgb(child.get("line", child.get("fill", "#FFFFFF")))
                connector.line.width = Pt(float(child.get("line_width", 1.5)))
                connector.name = child_name
                continue
            cfw = _frac(child_deck, child, "w", ref_w)
            cfh = _frac(child_deck, child, "h", ref_h)
            _validate_box(deck, child, cfx, cfy, cfw, cfh, f"group {group_index} child {child_index}")
            if deck.get("strict_input") and "type" not in child:
                _die(f"group {group_index} child {child_index} requires an explicit type in strict input mode")
            child_type = child.get("type", "rounded_rect")
            if child_type not in shape_types:
                _die(f"unsupported group shape type: {child_type}")
            shape = group.shapes.add_shape(
                shape_types[child_type],
                Emu(int(cfx * sw_emu)),
                Emu(int(cfy * sh_emu)),
                Emu(int(cfw * sw_emu)),
                Emu(int(cfh * sh_emu)),
            )
            shape.name = child_name
            apply_shape_fill(shape, child)
            if child.get("line"):
                shape.line.color.rgb = _hex_to_rgb(child["line"])
                shape.line.width = Pt(float(child.get("line_width", 1.0)))
            else:
                shape.line.fill.background()
            if child.get("rotation"):
                shape.rotation = float(child["rotation"])
            set_alt_text(shape, child.get("alt_text"))
        group.shapes._recalculate_extents()
        set_alt_text(group, group_spec.get("alt_text"))


def _set_cell_border(cell, border: dict | None) -> None:
    """Write optional native table cell borders in DrawingML."""
    if not isinstance(border, dict):
        return
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Pt

    tc_pr = cell._tc.get_or_add_tcPr()
    for side, key in (("L", "left"), ("R", "right"), ("T", "top"), ("B", "bottom")):
        value = border.get(key, border.get("all"))
        if not isinstance(value, dict):
            continue
        line = tc_pr.find(qn(f"a:ln{side}"))
        if line is None:
            line = OxmlElement(f"a:ln{side}")
            tc_pr.append(line)
        for child in list(line):
            line.remove(child)
        width = value.get("width", border.get("width", 0.75))
        width_emu = max(1, int(Pt(float(width))))
        line.set("w", str(width_emu))
        if str(value.get("style", "solid")).casefold() == "none":
            line.append(OxmlElement("a:noFill"))
            continue
        solid = OxmlElement("a:solidFill")
        color = OxmlElement("a:srgbClr")
        color.set("val", _normalized_hex(value.get("color", border.get("color", "#D9D9D9"))))
        solid.append(color)
        line.append(solid)
        dash = OxmlElement("a:prstDash")
        dash.set("val", "solid")
        line.append(dash)


def _set_cell_margins(cell, margins) -> None:
    if not isinstance(margins, dict):
        return
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement

    tc_pr = cell._tc.get_or_add_tcPr()
    tc_mar = tc_pr.find(qn("a:tcMar"))
    if tc_mar is None:
        tc_mar = OxmlElement("a:tcMar")
        tc_pr.append(tc_mar)
    for side in ("left", "right", "top", "bottom"):
        if side not in margins:
            continue
        element = tc_mar.find(qn(f"a:mar{side[0].upper()}"))
        if element is None:
            element = OxmlElement(f"a:mar{side[0].upper()}")
            tc_mar.append(element)
        element.set("w", str(max(0, int(float(margins[side]) * 12700))))
        element.set("type", "dxa")


def _write_cell_runs(cell, value, font: str, size: float, color: str | None, bold: bool, style: dict) -> None:
    cell.text = ""
    text_frame = cell.text_frame
    runs = value.get("runs") if isinstance(value, dict) else None
    if isinstance(runs, list) and runs:
        paragraph = text_frame.paragraphs[0]
        for raw_run in runs:
            if not isinstance(raw_run, dict):
                continue
            run = paragraph.add_run()
            run.text = str(raw_run.get("text", ""))
            _set_run_fonts(run, str(raw_run.get("font") or style.get("font") or font))
            run.font.size = Pt(float(raw_run.get("size", style.get("size", size))))
            run.font.bold = bool(raw_run.get("bold", style.get("bold", bold)))
            run_color = raw_run.get("color", style.get("color", color))
            if run_color:
                run.font.color.rgb = _hex_to_rgb(run_color)
    else:
        paragraph = text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = str(value.get("text", "") if isinstance(value, dict) else ("" if value is None else value))
        _set_run_fonts(run, str(style.get("font") or font))
        run.font.size = Pt(float(style.get("size", size)))
        run.font.bold = bool(style.get("bold", bold))
        run_color = style.get("color", color)
        if run_color:
            run.font.color.rgb = _hex_to_rgb(run_color)
    if style.get("align"):
        paragraph.alignment = style["align"]


def add_tables(slide, specs: list[dict], deck: dict, theme: dict, ref_w: float, ref_h: float, sw_emu: int, sh_emu: int):
    from pptx.util import Emu, Pt

    for table_index, spec in enumerate(specs, 1):
        rows = spec.get("rows", [])
        columns = int(spec.get("columns") or (len(rows[0]) if rows and isinstance(rows[0], list) else 0))
        if not rows or columns <= 0:
            _die(f"table {table_index} requires non-empty rows and columns")
        if any(not isinstance(row, list) or len(row) > columns for row in rows):
            _die(f"table {table_index} rows must be lists no wider than columns")
        if not any(cell is not None and str(cell).strip() for row in rows for cell in row):
            _die(f"table {table_index} requires at least one non-empty cell")
        table_fx = _frac(deck, spec, "x", ref_w)
        table_fy = _frac(deck, spec, "y", ref_h)
        table_fw = _frac(deck, spec, "w", ref_w)
        table_fh = _frac(deck, spec, "h", ref_h)
        _validate_box(deck, spec, table_fx, table_fy, table_fw, table_fh, f"table {table_index}")
        table = slide.shapes.add_table(
            len(rows),
            columns,
            Emu(int(table_fx * sw_emu)),
            Emu(int(table_fy * sh_emu)),
            Emu(int(table_fw * sw_emu)),
            Emu(int(table_fh * sh_emu)),
        ).table
        graphic_frame = table._graphic_frame
        graphic_frame.name = str(spec.get("object_id") or spec.get("name") or f"table-{table_index:02d}")
        font = str(spec.get("font") or theme.get("font") or "Noto Sans CJK SC")
        header_fill = spec.get("header_fill") or theme.get("table_header_fill")
        body_fill = spec.get("fill") or theme.get("table_fill")
        for column_index, width in enumerate((spec.get("column_widths") or [])[:columns]):
            table.columns[column_index].width = Emu(int(float(width) * sw_emu)) if deck["units"] == "fraction" else int(width)
        row_heights = spec.get("row_heights") or []
        cell_styles = spec.get("cell_styles") or {}
        border = spec.get("border")
        margins = spec.get("cell_margins") or spec.get("padding")
        default_size = float(spec.get("size", theme.get("size", 12)))
        default_color = spec.get("color") or theme.get("text_color")
        for row_index, row in enumerate(rows):
            if row_index < len(row_heights):
                raw_height = float(row_heights[row_index])
                table.rows[row_index].height = Emu(int(raw_height * sh_emu)) if deck["units"] == "fraction" else Emu(int(raw_height))
            for column_index in range(columns):
                cell = table.cell(row_index, column_index)
                value = row[column_index] if column_index < len(row) else ""
                style = {}
                if isinstance(cell_styles, dict):
                    style = cell_styles.get(f"{row_index},{column_index}", cell_styles.get(str(row_index), {}))
                if not isinstance(style, dict):
                    style = {}
                fill = style.get("fill", header_fill if row_index == 0 else body_fill)
                if fill:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _hex_to_rgb(fill)
                _write_cell_runs(
                    cell,
                    value,
                    str(style.get("font") or font),
                    float(style.get("size", default_size)),
                    style.get("color", default_color),
                    bool(style.get("bold", row_index == 0 and spec.get("header_bold", True))),
                    style,
                )
                _set_cell_border(cell, style.get("border", border))
                _set_cell_margins(cell, style.get("margins", margins))
        for merge in spec.get("merges", []):
            if isinstance(merge, list) and len(merge) == 4:
                r1, c1, r2, c2 = [int(value) for value in merge]
                table.cell(r1, c1).merge(table.cell(r2, c2))
        set_alt_text(graphic_frame, spec.get("alt_text"))


def add_charts(slide, specs: list[dict], deck: dict, theme: dict, ref_w: float, ref_h: float, sw_emu: int, sh_emu: int):
    import math

    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_LEGEND_POSITION
    from pptx.util import Emu

    charts = chart_map()
    for chart_index, spec in enumerate(specs, 1):
        chart_type = str(spec.get("type", "column")).casefold()
        if chart_type not in charts:
            _die(f"unsupported chart type: {chart_type}")
        data = CategoryChartData()
        data.categories = [str(value) for value in spec.get("categories", [])]
        series = spec.get("series", [])
        if not data.categories or not series:
            _die("chart requires categories and series")
        for item in series:
            if not isinstance(item, dict):
                _die("chart series entries must be objects")
            raw_values = item.get("values", [])
            if len(raw_values) != len(data.categories):
                _die("chart series length must match categories")
            values = []
            for value in raw_values:
                try:
                    number = float(value)
                except (TypeError, ValueError):
                    _die("chart values must be numeric")
                if not math.isfinite(number):
                    _die("chart values must be finite")
                values.append(number)
            data.add_series(str(item.get("name", "Series")), values)
        chart_fx = _frac(deck, spec, "x", ref_w)
        chart_fy = _frac(deck, spec, "y", ref_h)
        chart_fw = _frac(deck, spec, "w", ref_w)
        chart_fh = _frac(deck, spec, "h", ref_h)
        _validate_box(deck, spec, chart_fx, chart_fy, chart_fw, chart_fh, f"chart {chart_index}")
        graphic_frame = slide.shapes.add_chart(
            charts[chart_type],
            Emu(int(chart_fx * sw_emu)),
            Emu(int(chart_fy * sh_emu)),
            Emu(int(chart_fw * sw_emu)),
            Emu(int(chart_fh * sh_emu)),
            data,
        )
        graphic_frame.name = str(spec.get("object_id") or spec.get("name") or f"chart-{chart_index:02d}")
        chart = graphic_frame.chart
        chart.has_title = bool(spec.get("title"))
        if chart.has_title:
            chart.chart_title.text_frame.text = str(spec["title"])
        chart.has_legend = bool(spec.get("legend", len(series) > 1))
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
        if spec.get("data_labels"):
            for series_item in chart.series:
                series_item.has_data_labels = True
                series_item.data_labels.show_value = True
        for series_item, color in zip(chart.series, spec.get("colors") or theme.get("chart_colors") or []):
            series_item.format.fill.solid()
            series_item.format.fill.fore_color.rgb = _hex_to_rgb(color)

        # Charts are often placed on the same dark navy canvas as the rest of
        # the editable slide.  python-pptx leaves tick labels, legends and
        # data labels at Office's black default, which makes an otherwise
        # native chart visually unreadable.  Apply the deck theme explicitly
        # while keeping all chart data and axes native.
        chart_text_color = theme.get("chart_text_color") or theme.get("text_color") or "#FFFFFF"
        chart_muted_color = theme.get("chart_muted_color") or chart_text_color
        try:
            chart.font.name = str(theme.get("font") or "Aptos")
            chart.font.color.rgb = _hex_to_rgb(chart_text_color)
        except Exception:
            pass
        for axis in (getattr(chart, "category_axis", None), getattr(chart, "value_axis", None)):
            if axis is None:
                continue
            try:
                axis.tick_labels.font.name = str(theme.get("font") or "Aptos")
                axis.tick_labels.font.color.rgb = _hex_to_rgb(chart_text_color)
                axis.tick_labels.font.size = Pt(float(theme.get("chart_tick_size", 8)))
            except Exception:
                pass
            try:
                axis.format.line.color.rgb = _hex_to_rgb(chart_muted_color)
                axis.format.line.width = Pt(0.7)
            except Exception:
                pass
            try:
                if axis.has_major_gridlines:
                    axis.major_gridlines.format.line.color.rgb = _hex_to_rgb(theme.get("chart_grid_color") or chart_muted_color)
                    axis.major_gridlines.format.line.width = Pt(0.5)
            except Exception:
                pass
        if chart.has_legend and chart.legend is not None:
            try:
                chart.legend.font.name = str(theme.get("font") or "Aptos")
                chart.legend.font.color.rgb = _hex_to_rgb(chart_text_color)
                chart.legend.font.size = Pt(float(theme.get("chart_legend_size", 8)))
            except Exception:
                pass
        for series_item in chart.series:
            try:
                series_item.data_labels.font.name = str(theme.get("font") or "Aptos")
                series_item.data_labels.font.color.rgb = _hex_to_rgb(chart_text_color)
                series_item.data_labels.font.size = Pt(float(theme.get("chart_label_size", 7)))
            except Exception:
                pass
        set_alt_text(graphic_frame, spec.get("alt_text"))


def add_texts(slide, specs: list[dict], deck: dict, theme: dict, ref_w: float, ref_h: float, sw_emu: int, sh_emu: int):
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Emu, Pt

    alignments = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    anchors = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "center": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }
    strict_input = bool(deck.get("strict_input"))
    slide_height_pt = float(deck["slide_height_in"]) * 72.0
    for spec in specs:
        fx = _frac(deck, spec, "x", ref_w)
        fy = _frac(deck, spec, "y", ref_h)
        fw = _frac(deck, spec, "w", ref_w)
        fh = _frac(deck, spec, "h", ref_h)
        _validate_box(deck, spec, fx, fy, fw, fh, f"text {spec.get('object_id') or spec.get('name') or 'unnamed'}")
        box = slide.shapes.add_textbox(
            Emu(int(fx * sw_emu)),
            Emu(int(fy * sh_emu)),
            Emu(int(fw * sw_emu)),
            Emu(int(fh * sh_emu)),
        )
        frame = box.text_frame
        frame.word_wrap = True
        valign = spec.get("valign", "top")
        if strict_input and valign not in anchors:
            _die(f"unsupported vertical alignment: {valign}")
        frame.vertical_anchor = anchors.get(valign, MSO_ANCHOR.TOP)
        frame.margin_left = Pt(float(spec.get("margin_left", 0)))
        frame.margin_right = Pt(float(spec.get("margin_right", 0)))
        frame.margin_top = Pt(float(spec.get("margin_top", 0)))
        frame.margin_bottom = Pt(float(spec.get("margin_bottom", 0)))

        size_pt = text_size_pt(spec, slide_height_pt, ref_h)
        color = _hex_to_rgb(spec.get("color", theme.get("text_color", "#111111")))
        font = str(spec.get("font", theme.get("font", "Noto Sans CJK SC")))
        bold = bool(spec.get("bold", False))
        italic = bool(spec.get("italic", False))
        align = spec.get("align", "left")
        if strict_input and align not in alignments:
            _die(f"unsupported text alignment: {align}")
        alignment = alignments.get(align, PP_ALIGN.LEFT)
        line_spacing = spec.get("line_spacing")
        opacity = float(spec.get("opacity", 1.0))
        if spec.get("name") or spec.get("object_id"):
            box.name = str(spec.get("name") or spec["object_id"])

        runs = spec.get("runs")
        if runs:
            paragraph = frame.paragraphs[0]
            paragraph.alignment = alignment
            if line_spacing:
                paragraph.line_spacing = float(line_spacing)
            for run_spec in runs:
                run = paragraph.add_run()
                run.text = str(run_spec.get("text", ""))
                run.font.size = Pt(text_size_pt(run_spec, slide_height_pt, ref_h, default=size_pt))
                run.font.bold = bool(run_spec.get("bold", bold))
                run.font.italic = bool(run_spec.get("italic", italic))
                run.font.color.rgb = _hex_to_rgb(run_spec["color"]) if run_spec.get("color") else color
                _set_run_fonts(run, str(run_spec.get("font", font)))
                _set_run_alpha(run, float(run_spec.get("opacity", opacity)))
        else:
            for line_index, line in enumerate(str(spec.get("text", "")).split("\n")):
                paragraph = frame.paragraphs[0] if line_index == 0 else frame.add_paragraph()
                paragraph.alignment = alignment
                if line_spacing:
                    paragraph.line_spacing = float(line_spacing)
                run = paragraph.add_run()
                run.text = line
                run.font.size = Pt(size_pt)
                run.font.bold = bold
                run.font.italic = italic
                run.font.color.rgb = color
                _set_run_fonts(run, font)
                _set_run_alpha(run, opacity)


def _frac(deck: dict, item: dict, key: str, reference: float) -> float:
    if key not in item:
        _die(f"object is missing coordinate: {key}")
    try:
        value = float(item[key])
    except (TypeError, ValueError):
        _die(f"coordinate {key} must be numeric")
    units = deck.get("units", "fraction")
    if units not in {"fraction", "px"}:
        _die(f"unsupported coordinate units: {units}")
    if units == "px":
        if not reference or not math.isfinite(reference):
            _die(f"pixel coordinate {key} requires a positive reference canvas")
        value = value / reference
    if not math.isfinite(value):
        _die(f"coordinate {key} must be finite")
    if key in {"w", "h"} and value <= 0:
        _die(f"coordinate {key} must be positive")
    return value


def _validate_box(deck: dict, item: dict, x: float, y: float, w: float, h: float, label: str) -> None:
    """Reject silent geometry clipping when strict input mode is enabled."""
    if not deck.get("strict_input") or item.get("allow_bleed") is True:
        return
    if x < 0 or y < 0 or x + w > 1 or y + h > 1:
        _die(f"{label} box must stay within the slide in strict input mode: {(x, y, w, h)}")


# Private aliases preserve imports used by older helper scripts while the
# public module names describe the primitive responsibilities directly.
_apply_shape_fill = apply_shape_fill
_set_alt_text = set_alt_text
