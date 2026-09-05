#!/usr/bin/env python3
"""Regression coverage for explicit post-table native foreground shapes."""
from __future__ import annotations

import sys
import tempfile
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
if str(SCRIPTS) not in sys.path:
    sys.path.insert(0, str(SCRIPTS))

from authoring_backend import _partition_shapes, build_pptx  # noqa: E402


def main() -> int:
    background, foreground = _partition_shapes(
        {
            "shapes": [
                {"object_id": "base", "type": "rect", "x": 0, "y": 0, "w": 1, "h": 1},
                {"object_id": "bar", "type": "rect", "x": 0.5, "y": 0.5, "w": 0.1, "h": 0.02, "z_layer": "foreground"},
            ]
        }
    )
    assert [item["object_id"] for item in background] == ["base"]
    assert [item["object_id"] for item in foreground] == ["bar"]

    with tempfile.TemporaryDirectory(prefix="foreground-shape-") as temp_dir:
        temp = Path(temp_dir)
        output = temp / "foreground-shape.pptx"
        deck = {
            "slide_width_in": 13.333,
            "slide_height_in": 7.5,
            "ref_width": 1600,
            "ref_height": 900,
            "units": "fraction",
            "assets_dir": str(temp),
            "font_family": "Noto Sans CJK SC",
            "theme": {
                "font": "Noto Sans CJK SC",
                "text_color": "#0A2B5E",
                "table_header_fill": "#0A4A86",
                "table_fill": "#FFFFFF",
                "size": 9,
            },
            "slides": [
                {
                    "layout_name": "Blank",
                    "shapes": [
                        {
                            "object_id": "base",
                            "type": "rect",
                            "x": 0,
                            "y": 0,
                            "w": 1,
                            "h": 1,
                            "fill": "#FFFFFF",
                        },
                        {
                            "object_id": "progress-track",
                            "type": "rect",
                            "x": 0.63,
                            "y": 0.43,
                            "w": 0.20,
                            "h": 0.035,
                            "fill": "#E7EEF8",
                            "line": "#B8C7DB",
                            "line_width": 0.6,
                            "z_layer": "foreground",
                        },
                        {
                            "object_id": "progress-fill",
                            "type": "rect",
                            "x": 0.63,
                            "y": 0.43,
                            "w": 0.12,
                            "h": 0.035,
                            "fill": "#357DE8",
                            "z_layer": "foreground",
                        },
                    ],
                    "groups": [],
                    "tables": [
                        {
                            "object_id": "test-table",
                            "native_required": True,
                            "x": 0.1,
                            "y": 0.25,
                            "w": 0.8,
                            "h": 0.45,
                            "columns": 2,
                            "rows": [["项目", "进度"], ["A", "60%"]],
                            "header_fill": "#0A4A86",
                            "fill": "#FFFFFF",
                            "color": "#0A2B5E",
                            "border": {"all": {"color": "#9EB3C8", "width": 0.7}},
                        }
                    ],
                    "charts": [],
                    "icons": [],
                    "texts": [
                        {
                            "object_id": "progress-label",
                            "text": "60%",
                            "x": 0.84,
                            "y": 0.43,
                            "w": 0.05,
                            "h": 0.035,
                            "size": 9,
                            "color": "#0A4A86",
                        }
                    ],
                }
            ],
        }
        build_pptx(deck, output)

        from pptx import Presentation

        presentation = Presentation(str(output))
        names = [shape.name for shape in presentation.slides[0].shapes]
        assert names.index("base") < names.index("test-table")
        assert names.index("test-table") < names.index("progress-track")
        assert names.index("progress-track") < names.index("progress-fill")
        assert names.index("progress-fill") < names.index("progress-label")

    print("foreground shape layer: passed")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
