#!/usr/bin/env python3
"""End-to-end regression coverage for native panels and tables."""
from __future__ import annotations

import json
import subprocess
import sys
import tempfile
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[2]
PNG_1X1 = bytes.fromhex(
    "89504e470d0a1a0a0000000d4948445200000001000000010806000000"
    "1f15c4890000000d49444154789c6360000000020001e221bc3300000000"
    "49454e44ae426082"
)


def run(*args: str) -> subprocess.CompletedProcess[str]:
    return subprocess.run([sys.executable, *args], cwd=ROOT, capture_output=True, text=True, check=False)


def write(path: Path, value: dict) -> None:
    path.write_text(json.dumps(value, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def main() -> int:
    with tempfile.TemporaryDirectory(prefix="native-structure-") as temp:
        root = Path(temp)
        layout = root / "layout.json"
        deck = root / "native.pptx"
        manifest = root / "slide-object-manifest.json"
        report = root / "native-validation.json"
        write(layout, {
            "slide_width_in": 4,
            "slide_height_in": 2.25,
            "units": "fraction",
            "editable_object_policy": "native-semantic-objects",
            "slides": [{
                "slide_no": 1,
                "native_panels": [{
                    "object_id": "panel-01",
                    "treatment": "native-group",
                    "x": 0.05,
                    "y": 0.08,
                    "w": 0.90,
                    "h": 0.84,
                    "children_coordinate_space": "local",
                    "children": [{
                        "object_id": "panel-fill",
                        "type": "rounded_rect",
                        "x": 0.0,
                        "y": 0.0,
                        "w": 1.0,
                        "h": 1.0,
                        "fill": "#F4F7FB",
                        "line": "#D9E2F0",
                        "line_width": 1.0,
                    }],
                }],
                "native_tables": [{
                    "object_id": "data-table-01",
                    "native_required": True,
                    "x": 0.12,
                    "y": 0.42,
                    "w": 0.76,
                    "h": 0.34,
                    "columns": 2,
                    "rows": [["Metric", "Value"], ["Editable", "Native"]],
                    "header_fill": "#EAF1FB",
                    "border": {"all": {"color": "#B7C9E2", "width": 0.75}},
                    "cell_margins": {"left": 0.04, "right": 0.04, "top": 0.02, "bottom": 0.02},
                }],
                "texts": [{
                    "object_id": "title",
                    "text": "Native structure",
                    "x": 0.12,
                    "y": 0.14,
                    "w": 0.70,
                    "h": 0.16,
                    "size": 14,
                    "bold": True,
                }],
            }],
        })
        write(manifest, {
            "schema": "ai-ppt-plus/slide-object-manifest/v1",
            "slides": [{
                "slide_no": 1,
                "objects": [
                    {"object_id": "panel-01", "role": "semantic-panel", "object_type": "native_group", "native_required": True},
                    {"object_id": "data-table-01", "role": "table", "object_type": "editable_table", "native_required": True},
                    {"object_id": "title", "role": "formal-text", "object_type": "editable_text"},
                ],
            }],
        })

        composed = run(
            "ai-ppt-editable/scripts/compose_pptx.py",
            str(layout), str(deck), "--strict-input", "--require-native-structure",
        )
        assert composed.returncode == 0, composed.stdout + composed.stderr
        prs = Presentation(str(deck))
        shapes = list(prs.slides[0].shapes)
        assert any(shape.name == "panel-01" and "GROUP" in str(shape.shape_type).upper() for shape in shapes), [shape.name for shape in shapes]
        table_shapes = [shape for shape in shapes if bool(getattr(shape, "has_table", False))]
        assert any(shape.name == "data-table-01" for shape in table_shapes), [shape.name for shape in shapes]
        assert any(shape.name == "title" and "Native structure" in shape.text for shape in shapes), [shape.name for shape in shapes]

        checked = run(
            "ai-ppt-editable/scripts/validate_native_editability.py",
            str(deck), "--object-manifest", str(manifest), "--require-native-structure", "--report", str(report),
        )
        assert checked.returncode == 0, checked.stdout + checked.stderr
        result = json.loads(report.read_text(encoding="utf-8"))
        assert result["valid"] is True, result
        assert result["native_table_count"] == 1, result
        assert result["native_panel_count"] == 1, result
        assert result["slides"][0]["objects"][0]["kind"] in {"native_group", "native_shape"}, result

        frame = root / "frame.png"
        frame.write_bytes(PNG_1X1)
        raster_layout = root / "raster-layout.json"
        write(raster_layout, {
            "slide_width_in": 4,
            "slide_height_in": 2.25,
            "assets_dir": str(root),
            "editable_object_policy": "native-semantic-objects",
            "slides": [{"slide_no": 1, "frame": "frame.png"}],
        })
        raster = run(
            "ai-ppt-editable/scripts/compose_pptx.py",
            str(raster_layout), str(root / "raster.pptx"), "--strict-input",
        )
        assert raster.returncode == 2, raster.stdout + raster.stderr
        assert "semantic frame images are disabled" in raster.stderr, raster.stdout + raster.stderr

        raster_panel_layout = root / "raster-panel-layout.json"
        write(raster_panel_layout, {
            "slide_width_in": 4,
            "slide_height_in": 2.25,
            "assets_dir": str(root),
            "editable_object_policy": "native-semantic-objects",
            "slides": [{
                "slide_no": 1,
                "native_panels": [{
                    "object_id": "bad-panel",
                    "native_required": True,
                    "file": "frame.png",
                    "x": 0.1,
                    "y": 0.1,
                    "w": 0.8,
                    "h": 0.8,
                }],
            }],
        })
        raster_panel = run(
            "ai-ppt-editable/scripts/compose_pptx.py",
            str(raster_panel_layout), str(root / "raster-panel.pptx"), "--strict-input",
        )
        assert raster_panel.returncode == 2, raster_panel.stdout + raster_panel.stderr
        assert "cannot be promoted from a raster file" in raster_panel.stderr, raster_panel.stdout + raster_panel.stderr

    print("native panels and tables: ok")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
