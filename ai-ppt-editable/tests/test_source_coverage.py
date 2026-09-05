import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
from reconstruction.graph_ir import PageGraph
from reconstruction.source_coverage import audit_source_coverage, extract_pptx_objects
from reconstruction.pipeline import ReconstructionPipeline, Stage


def fixtures():
    inventory = {"source_sha256": "a" * 64, "observation_id": "observation-1",
                 "method": "source-image-observation", "evidence": "reference-region-review",
                 "objects": [{"id": "source-title"}, {"id": "source-icon"}]}
    graph = PageGraph.from_dict({"metadata": {"source_sha256": "a" * 64,
        "planning_observation_id": "planning-2", "route": "reference-reconstruction"},
        "nodes": [{"id": "title", "type": "text", "bbox": [0, 0, .4, .1],
                   "source": {"source_object_id": "source-title"}},
                  {"id": "icon", "type": "icon", "bbox": [.5, .2, .1, .1],
                   "source": {"source_object_id": "source-icon"}}]})
    return inventory, graph


def test_complete_three_way():
    inventory, graph = fixtures()
    result = audit_source_coverage(inventory, graph, {"method": "pptx-extraction",
        "deck_sha256": "b" * 64, "object_ids": ["title", "icon"]})
    assert result["valid"]
    assert result["scope"] == "three-way"


def test_shared_omission_is_blocked():
    inventory, graph = fixtures()
    from dataclasses import replace
    graph = replace(graph, nodes=graph.nodes[:1])
    result = audit_source_coverage(inventory, graph, {"method": "pptx-extraction",
        "deck_sha256": "b" * 64, "object_ids": ["title"]})
    assert not result["valid"]
    assert "missing planned source object: source-icon" in result["errors"]


def test_stale_and_self_derived_inventory():
    inventory, graph = fixtures()
    inventory["source_sha256"] = "c" * 64
    inventory["observation_id"] = "planning-2"
    result = audit_source_coverage(inventory, graph)
    assert len(result["errors"]) == 2


def test_uncertainty_duplicates_and_final_extras():
    inventory, graph = fixtures()
    inventory["objects"][0]["uncertain"] = True
    result = audit_source_coverage(inventory, graph, {"method": "pptx-extraction",
        "deck_sha256": "b" * 64, "object_ids": ["title", "title", "extra"]})
    assert not result["valid"]
    assert len(result["errors"]) == 4


def test_pipeline_blocks_before_author_without_inventory():
    _, graph = fixtures()
    def forbidden(*args):
        raise AssertionError("must not author without source coverage")
    result = ReconstructionPipeline().run(understand=lambda: graph, author=forbidden,
        render=forbidden, inspect=forbidden, apply_repairs=forbidden, measure=forbidden)
    assert result.stage == Stage.BLOCKED


def test_actual_pptx_extraction(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(0, 0, Inches(1), Inches(1)).name = "title"
    group = slide.shapes.add_group_shape()
    group.name = "group"
    group.shapes.add_textbox(0, 0, Inches(1), Inches(1)).name = "nested"
    path = tmp_path / "actual.pptx"
    prs.save(path)
    observed = extract_pptx_objects(str(path))
    assert observed["object_ids"] == ["title", "group", "nested"]
    assert len(observed["deck_sha256"]) == 64


if __name__ == "__main__":
    from tempfile import TemporaryDirectory
    test_complete_three_way()
    test_shared_omission_is_blocked()
    test_stale_and_self_derived_inventory()
    test_uncertainty_duplicates_and_final_extras()
    test_pipeline_blocks_before_author_without_inventory()
    with TemporaryDirectory() as folder:
        test_actual_pptx_extraction(Path(folder))
    print("6 source coverage tests passed")
