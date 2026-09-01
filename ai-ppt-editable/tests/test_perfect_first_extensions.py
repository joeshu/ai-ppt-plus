#!/usr/bin/env python3
"""Regression tests for the perfect-first outer contracts."""
from __future__ import annotations

import hashlib
import json
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]


def digest(value: object) -> str:
    return hashlib.sha256(json.dumps(value, ensure_ascii=False, sort_keys=True, separators=(",", ":")).encode("utf-8")).hexdigest()


def run(*args: str) -> subprocess.CompletedProcess[str]:
    return subprocess.run([sys.executable, *args], cwd=ROOT, capture_output=True, text=True, check=False)


def chart_manifest(path: Path) -> None:
    snapshot = {
        "kind": "category_chart",
        "categories": ["一月", "二月", "三月"],
        "series": [
            {"series_id": "sales", "name": "销量", "values": [10, 12, 14]},
            {"series_id": "target", "name": "目标", "values": [9, 11, 13]},
        ],
    }
    value = {
        "schema": "ai-ppt-plus/chart-reconstruction/v1",
        "project_id": "perfect-first-fixture",
        "coordinate_space": "reference_pixels",
        "canvas": [400, 225],
        "charts": [{
            "chart_id": "native-chart",
            "slide_no": 1,
            "title": "月度趋势",
            "representation": "native_chart",
            "editability_level": "L1",
            "source_data_status": "verified",
            "data_source": {"kind": "approved_table", "method": "human_confirmed"},
            "categories": snapshot["categories"],
            "series": [
                {"series_id": "sales", "name": "销量", "color": "#2A6FDB", "values": [10, 12, 14], "value_labels": [{"content": "10"}]},
                {"series_id": "target", "name": "目标", "color": "#E58B27", "values": [9, 11, 13], "value_labels": [{"content": "9"}]},
            ],
            "missing_value_policy": "blank_not_zero",
            "data_snapshot_sha256": digest(snapshot),
            "geometry": {"source_bbox": [200, 22.5, 160, 90], "plot_bbox": [215, 40, 130, 55], "point_anchor_tolerance": 0.015},
            "required_elements": ["title", "legend", "category_labels"],
            "visible_elements": {"title": [{"content": "月度趋势"}], "legend": [{"content": "销量"}], "category_labels": [{"content": "一月"}]},
            "qa": {"reference_region": [200, 22.5, 160, 90]},
        }],
    }
    path.write_text(json.dumps(value, ensure_ascii=False), encoding="utf-8")


def main() -> int:
    with tempfile.TemporaryDirectory(prefix="perfect-first-extensions-") as temp:
        root = Path(temp)
        layout = root / "layout.json"
        layout.write_text(json.dumps({
            "project_id": "perfect-first-fixture",
            "slide_width_in": 4,
            "slide_height_in": 2.25,
            "units": "fraction",
            "theme": {"font_family": "Noto Sans CJK SC", "font_color": "#223344"},
            "slides": [{
                "shapes": [{
                    "object_id": "gradient-card", "type": "rounded_rect", "x": 0.05, "y": 0.05, "w": 0.4, "h": 0.25,
                    "gradient": {"angle_deg": 90, "stops": [{"position_pct": 0, "color": "#FF0000FF"}, {"position_pct": 100, "color": "#0000FF80"}]},
                }],
                "charts": [{"object_id": "native-chart", "type": "line", "x": 0.5, "y": 0.1, "w": 0.4, "h": 0.4}],
                "texts": [{
                    "object_id": "styled-title", "bbox": [0.05, 0.4, 0.4, 0.15], "content": "策略 精准",
                    "style": {"font_family": "Noto Sans CJK SC", "size_pt": 18, "font_color": "#223344", "bold": True},
                    "runs": [{"content": "策略 ", "style": {"font_family": "Noto Sans CJK SC", "size_pt": 18, "font_color": "#223344", "bold": True}}, {"content": "精准", "style": {"font_family": "Noto Sans CJK SC", "size_pt": 20, "font_color": "#D52B2B"}}],
                }],
            }],
        }, ensure_ascii=False), encoding="utf-8")
        chart = root / "chart-reconstruction.json"
        chart_manifest(chart)
        deck = root / "deck.pptx"
        adapter_report = root / "adapter.json"
        gradient_report = root / "gradient.json"
        typography_report = root / "typography.json"
        composed = run("scripts/compose_pptx.py", str(layout), str(deck), "--strict-input", "--chart-manifest", str(chart), "--adapter-report", str(adapter_report), "--gradient-report", str(gradient_report), "--typography-report", str(typography_report))
        assert composed.returncode == 0, composed.stdout + composed.stderr

        prs = Presentation(deck)
        slide = prs.slides[0]
        chart_shapes = [shape for shape in slide.shapes if getattr(shape, "has_chart", False)]
        assert len(chart_shapes) == 1
        native_chart = chart_shapes[0].chart
        assert [category.label for category in native_chart.plots[0].categories] == ["一月", "二月", "三月"]
        assert [list(series.values) for series in native_chart.series] == [[10.0, 12.0, 14.0], [9.0, 11.0, 13.0]]
        title = next(shape for shape in slide.shapes if shape.name == "styled-title")
        assert [run.text for run in title.text_frame.paragraphs[0].runs] == ["策略 ", "精准"]
        assert str(title.text_frame.paragraphs[0].runs[1].font.color.rgb) == "D52B2B"
        with zipfile.ZipFile(deck) as package:
            xml = package.read("ppt/slides/slide1.xml")
            assert b"<a:gradFill" in xml
            assert b'val="FF0000"' in xml and b'val="0000FF"' in xml

        adapter = json.loads(adapter_report.read_text(encoding="utf-8"))
        assert adapter["charts"]["charts"][0]["route"] == "native_chart"
        assert adapter["gradients"]["native_gradient_count"] == 1
        typography = json.loads(typography_report.read_text(encoding="utf-8"))
        assert typography["valid"] is True and typography["text_count"] == 3
        assert "#D52B2B" in typography["font_colors"]

        object_manifest = root / "objects.json"
        built = run("scripts/build_object_manifest.py", str(layout), "--chart-manifest", str(chart), "--output", str(object_manifest))
        assert built.returncode == 0, built.stdout + built.stderr
        audit = run("scripts/inspect_editable_objects.py", str(deck), "--object-manifest", str(object_manifest), "--require-types", "--require-geometry", "--report", str(root / "objects-audit.json"))
        assert audit.returncode == 0, audit.stdout + audit.stderr
        audit_data = json.loads((root / "objects-audit.json").read_text(encoding="utf-8"))
        assert audit_data["geometry_checked"] == 3 and audit_data["geometry_mismatch_count"] == 0, audit_data

        invalid = json.loads(chart.read_text(encoding="utf-8"))
        invalid["charts"][0]["source_data_status"] = "unverified"
        invalid_chart = root / "invalid-chart.json"
        invalid_chart.write_text(json.dumps(invalid, ensure_ascii=False), encoding="utf-8")
        blocked = run("scripts/compose_pptx.py", str(layout), str(root / "blocked.pptx"), "--chart-manifest", str(invalid_chart))
        assert blocked.returncode != 0
        assert "chart manifest is invalid" in blocked.stderr or "perfect-first contract failed" in blocked.stderr
    print("perfect-first extensions: ok")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
