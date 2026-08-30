#!/usr/bin/env python3
"""Regression test for the PresentationML font embedding adapter."""
from __future__ import annotations

import hashlib
import json
import struct
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path


ROOT = Path(__file__).resolve().parents[1]


def run(*args: str) -> subprocess.CompletedProcess[str]:
    return subprocess.run([sys.executable, *args], cwd=ROOT, capture_output=True, text=True, check=False)


def synthetic_sfnt(family: str, restricted: bool = False) -> bytes:
    """Build a tiny valid-enough OTTO SFNT for deterministic adapter tests."""
    head = bytearray(54)
    head[0:4] = struct.pack(">I", 0x00010000)
    head[18:20] = struct.pack(">H", 1000)
    head[44:46] = struct.pack(">H", 0)

    os2 = bytearray(86)
    os2[0:2] = struct.pack(">H", 4)
    os2[4:6] = struct.pack(">H", 400)
    os2[8:10] = struct.pack(">H", 0x0002 if restricted else 0)
    os2[32:42] = bytes((2, 11, 5, 2, 2, 2, 2, 2, 2, 4))
    os2[42:58] = b"\x00" * 16
    os2[62:64] = struct.pack(">H", 0)
    os2[78:86] = b"\x00" * 8

    records = []
    storage = bytearray()
    names = ((1, family), (2, "Regular"), (4, family + " Regular"), (5, "Version 1.0"))
    for name_id, value in names:
        encoded = value.encode("utf-16-be")
        records.append(struct.pack(">HHHHHH", 3, 1, 0x0804, name_id, len(encoded), len(storage)))
        storage.extend(encoded)
    name = struct.pack(">HHH", 0, len(records), 6 + 12 * len(records)) + b"".join(records) + storage

    tables = [(b"OS/2", bytes(os2)), (b"head", bytes(head)), (b"name", bytes(name))]
    directory_size = 12 + 16 * len(tables)
    offset = (directory_size + 3) & ~3
    directory = bytearray(struct.pack(">IHHHH", 0x4F54544F, len(tables), 0, 0, 0))
    payload = bytearray(b"\x00" * (offset - directory_size))
    for tag, table in tables:
        offset = (directory_size + len(payload) + 3) & ~3
        padding = offset - (directory_size + len(payload))
        payload.extend(b"\x00" * padding)
        record = struct.pack(">4sIII", tag, 0, offset, len(table))
        directory.extend(record)
        payload.extend(table)
    return bytes(directory + payload)


def write_manifest(font_dir: Path, filename: str, family: str) -> Path:
    font_path = font_dir / filename
    manifest = {
        "file": filename,
        "family": family,
        "sha256": hashlib.sha256(font_path.read_bytes()).hexdigest(),
        "license": "synthetic regression fixture",
        "license_url": "https://example.invalid/font-fixture-license",
    }
    path = font_dir / "font-manifest.json"
    path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    return path


def write_input_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = textbox.text_frame.paragraphs[0].add_run()
    run.text = "中文字体嵌入测试"
    run.font.size = Pt(24)
    presentation.save(path)


def rewrite_package(source: Path, target: Path, replacements: dict[str, bytes] | None = None, additions: dict[str, bytes] | None = None) -> None:
    replacements = replacements or {}
    additions = additions or {}
    with zipfile.ZipFile(source) as original, zipfile.ZipFile(target, "w", compression=zipfile.ZIP_DEFLATED) as changed:
        names = set()
        for info in original.infolist():
            names.add(info.filename)
            changed.writestr(info, replacements.get(info.filename, original.read(info.filename)))
        for name, data in additions.items():
            if name not in names:
                changed.writestr(name, data)


def main() -> int:
    family = "Synthetic CJK Fixture"
    with tempfile.TemporaryDirectory(prefix="font-embedding-") as temp:
        root = Path(temp)
        input_pptx = root / "input.pptx"
        write_input_pptx(input_pptx)

        good_dir = root / "good-fonts"
        good_dir.mkdir()
        (good_dir / "fixture.otf").write_bytes(synthetic_sfnt(family))
        write_manifest(good_dir, "fixture.otf", family)
        output = root / "embedded.pptx"
        embedding_report = root / "embedding.json"
        embedded = run("scripts/embed_fonts.py", str(input_pptx), str(output), "--font-dir", str(good_dir), "--report", str(embedding_report))
        assert embedded.returncode == 0, embedded.stdout + embedded.stderr
        report = json.loads(embedding_report.read_text(encoding="utf-8"))
        assert report["valid"] is True
        assert report["package"]["font_parts"] == ["ppt/fonts/font1.fntdata"]
        assert report["fonts"][0]["metadata"]["family"] == family

        inspection = root / "inspection.json"
        inspected = run("scripts/inspect_pptx.py", str(output), "--report", str(inspection))
        assert inspected.returncode == 0, inspected.stdout + inspected.stderr
        inspection_data = json.loads(inspection.read_text(encoding="utf-8"))
        evidence = inspection_data["embedded_fonts"]
        assert evidence["present"] is True
        assert evidence["part_evidence"][0]["valid"] is True
        with zipfile.ZipFile(output) as package:
            assert "ppt/fonts/font1.fntdata" in package.namelist()
            content_types = package.read("[Content_Types].xml").decode("utf-8")
            assert "application/x-fontdata" in content_types
            presentation = package.read("ppt/presentation.xml").decode("utf-8")
            assert "embeddedFontLst" in presentation
            assert "embedTrueTypeFonts=\"1\"" in presentation

        # Variable-font exports may expose a style suffix in name ID 1 while
        # fontconfig exposes the root family too.  The manifest should be able
        # to declare that stable root family and embed it under the same name.
        alias_dir = root / "alias-fonts"
        alias_dir.mkdir()
        alias_family = "Synthetic Variable"
        (alias_dir / "fixture.otf").write_bytes(synthetic_sfnt(alias_family + " Thin"))
        write_manifest(alias_dir, "fixture.otf", alias_family)
        alias_output = root / "alias-embedded.pptx"
        alias_report = root / "alias-embedding.json"
        aliased = run("scripts/embed_fonts.py", str(input_pptx), str(alias_output), "--font-dir", str(alias_dir), "--report", str(alias_report))
        assert aliased.returncode == 0, aliased.stdout + aliased.stderr
        alias_data = json.loads(alias_report.read_text(encoding="utf-8"))
        assert alias_data["valid"] is True
        assert alias_data["fonts"][0]["metadata"]["family"] == alias_family

        malformed = root / "malformed.pptx"
        rewrite_package(output, malformed, {"ppt/fonts/font1.fntdata": b"not-an-eot"})
        malformed_report = root / "malformed-inspection.json"
        malformed_checked = run("scripts/inspect_pptx.py", str(malformed), "--report", str(malformed_report))
        assert malformed_checked.returncode == 2
        malformed_data = json.loads(malformed_report.read_text(encoding="utf-8"))
        assert any(item["code"] == "malformed_embedded_font_part" for item in malformed_data["issues"])

        orphan = root / "orphan.pptx"
        rewrite_package(output, orphan, additions={"ppt/fonts/orphan.fntdata": b"orphan"})
        orphan_report = root / "orphan-inspection.json"
        orphan_checked = run("scripts/inspect_pptx.py", str(orphan), "--report", str(orphan_report))
        assert orphan_checked.returncode == 2
        orphan_data = json.loads(orphan_report.read_text(encoding="utf-8"))
        assert any(item["code"] == "orphan_embedded_font_part" for item in orphan_data["issues"])

        wrong_type = root / "wrong-type.pptx"
        with zipfile.ZipFile(output) as package:
            content_types = package.read("[Content_Types].xml").decode("utf-8").replace("application/x-fontdata", "application/octet-stream")
        rewrite_package(output, wrong_type, {"[Content_Types].xml": content_types.encode("utf-8")})
        wrong_type_report = root / "wrong-type-inspection.json"
        wrong_type_checked = run("scripts/inspect_pptx.py", str(wrong_type), "--report", str(wrong_type_report))
        assert wrong_type_checked.returncode == 2
        wrong_type_data = json.loads(wrong_type_report.read_text(encoding="utf-8"))
        assert any(item["code"] == "invalid_embedded_font_content_type" for item in wrong_type_data["issues"])

        layout = root / "layout.json"
        layout.write_text(json.dumps({
            "slide_width_in": 13.333,
            "slide_height_in": 7.5,
            "slides": [{"texts": [{
                "text": "中文字体嵌入测试",
                "x": 0.1, "y": 0.1, "w": 0.5, "h": 0.1,
                "font": family, "size": 24,
            }]}],
        }, ensure_ascii=False), encoding="utf-8")
        composed = root / "composed.pptx"
        composed_report = root / "composed-embedding.json"
        composed_result = run(
            "scripts/compose_pptx.py", str(layout), str(composed),
            "--font-dir", str(good_dir), "--embed-fonts",
            "--embedding-report", str(composed_report),
        )
        assert composed_result.returncode == 0, composed_result.stdout + composed_result.stderr
        composed_inspection = root / "composed-inspection.json"
        composed_checked = run("scripts/inspect_pptx.py", str(composed), "--report", str(composed_inspection))
        assert composed_checked.returncode == 0, composed_checked.stdout + composed_checked.stderr
        assert json.loads(composed_inspection.read_text(encoding="utf-8"))["embedded_fonts"]["present"] is True

        restricted_dir = root / "restricted-fonts"
        restricted_dir.mkdir()
        (restricted_dir / "restricted.otf").write_bytes(synthetic_sfnt(family, restricted=True))
        write_manifest(restricted_dir, "restricted.otf", family)
        restricted_output = root / "restricted-output.pptx"
        restricted_report = root / "restricted.json"
        blocked = run("scripts/embed_fonts.py", str(input_pptx), str(restricted_output), "--font-dir", str(restricted_dir), "--report", str(restricted_report))
        assert blocked.returncode == 2
        assert not restricted_output.exists()
        blocked_data = json.loads(restricted_report.read_text(encoding="utf-8"))
        assert blocked_data["valid"] is False
        assert "restricted-license" in blocked_data["issues"][0]["message"]
    print("PresentationML font embedding and restricted-rights gate: ok")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
