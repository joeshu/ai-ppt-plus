#!/usr/bin/env python3
"""Place raster/vector assets and manage SVG conversion lifetimes."""
from __future__ import annotations

import os
import posixpath
import subprocess
import sys
import tempfile
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

from atomic_output import atomic_rewrite_zip
from pptx_primitives import set_alt_text


def _die(message: str, code: int = 2):
    print(f"Error: {message}", file=sys.stderr)
    raise SystemExit(code)


def _resolve(assets_dir: Path, file: str) -> Path:
    path = Path(file)
    return path if path.is_absolute() else assets_dir / path


def _frac(deck: dict, item: dict, key: str, reference: float) -> float:
    value = item[key]
    if deck["units"] == "px":
        return value / reference
    return value


def replace_svg_media(pptx_path: Path, svg_assets: list[tuple[int, str, Path]]) -> None:
    """Replace temporary PNG media with native SVG package parts atomically."""
    if not svg_assets:
        return
    presentation_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    drawing_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    relationship_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    with zipfile.ZipFile(pptx_path, "r") as archive:
        entries = {name: archive.read(name) for name in archive.namelist()}

    replacements = []
    for slide_no, object_name, svg_path in svg_assets:
        slide_name = f"ppt/slides/slide{slide_no}.xml"
        relationships_name = f"ppt/slides/_rels/slide{slide_no}.xml.rels"
        root = ET.fromstring(entries[slide_name])
        relationships = ET.fromstring(entries[relationships_name])
        target = None
        for picture in root.findall(f".//{{{presentation_ns}}}pic"):
            non_visual = picture.find(f".//{{{presentation_ns}}}cNvPr")
            if non_visual is None or non_visual.get("name") != object_name:
                continue
            blip = picture.find(f".//{{{drawing_ns}}}blip")
            relationship_id = blip.get(f"{{{relationship_ns}}}embed") if blip is not None else None
            for relationship in relationships:
                if relationship.get("Id") == relationship_id:
                    target = posixpath.normpath(posixpath.join("ppt/slides", relationship.get("Target")))
                    relationship.set("Target", "../media/" + Path(target).with_suffix(".svg").name)
                    break
            break
        if target is None or target not in entries:
            _die(f"could not resolve SVG media relationship for {object_name}")
        new_target = str(Path(target).with_suffix(".svg")).replace("\\", "/")
        entries[relationships_name] = ET.tostring(relationships, encoding="utf-8", xml_declaration=True)
        entries[new_target] = svg_path.read_bytes()
        del entries[target]
        replacements.append((target, new_target))

    content_types = ET.fromstring(entries["[Content_Types].xml"])
    for old, new in replacements:
        for override in content_types:
            if override.get("PartName") == "/" + old:
                override.set("PartName", "/" + new)
                override.set("ContentType", "image/svg+xml")
    entries["[Content_Types].xml"] = ET.tostring(content_types, encoding="utf-8", xml_declaration=True)
    atomic_rewrite_zip(pptx_path, entries)


def svg_to_png(svg_path: Path) -> Path:
    """Rasterize an SVG into a uniquely named temporary PNG."""
    handle, name = tempfile.mkstemp(prefix="ai-ppt-svg-", suffix=".png")
    os.close(handle)
    output = Path(name)
    try:
        from cairosvg import svg2png

        svg2png(url=str(svg_path), write_to=str(output))
        return output
    except ImportError:
        for command in (
            ("inkscape", str(svg_path), "--export-filename", str(output)),
            ("convert", str(svg_path), str(output)),
        ):
            try:
                result = subprocess.run(command, capture_output=True, text=True, timeout=30, check=False)
            except (FileNotFoundError, subprocess.TimeoutExpired):
                continue
            if result.returncode == 0 and output.is_file():
                return output
        output.unlink(missing_ok=True)
        _die(f"SVG asset requires cairosvg, inkscape, or ImageMagick: {svg_path}")
    except Exception:
        output.unlink(missing_ok=True)
        raise
    return output


def cleanup_temporary_files(paths: list[Path]) -> None:
    for path in paths:
        path.unlink(missing_ok=True)


def add_background(slide, slide_spec: dict, assets_dir: Path, sw_emu: int, sh_emu: int):
    from pptx.util import Emu

    background = slide_spec.get("background")
    if not background:
        return
    path = _resolve(assets_dir, background)
    if not path.exists():
        _die(f"background not found: {path}")
    picture = slide.shapes.add_picture(str(path), 0, 0, width=Emu(sw_emu), height=Emu(sh_emu))
    picture.name = str(slide_spec.get("background_object_id", "background"))
    set_alt_text(picture, slide_spec.get("background_alt_text"))


def add_frame(slide, slide_spec: dict, assets_dir: Path, sw_emu: int, sh_emu: int):
    from pptx.util import Emu

    frame = slide_spec.get("frame")
    if not frame:
        return
    path = _resolve(assets_dir, frame)
    if not path.exists():
        _die(f"frame not found: {path}")
    picture = slide.shapes.add_picture(str(path), 0, 0, width=Emu(sw_emu), height=Emu(sh_emu))
    picture.name = str(slide_spec.get("frame_object_id", "frame"))
    set_alt_text(picture, slide_spec.get("frame_alt_text"))


def add_panels(slide, specs: list[dict], assets_dir: Path, deck: dict, ref_w: float, ref_h: float, sw_emu: int, sh_emu: int, slide_no: int):
    from pptx.util import Emu

    for panel in specs:
        path = _resolve(assets_dir, panel["file"])
        if not path.exists():
            _die(f"slide {slide_no}: panel not found: {path}")
        fx = _frac(deck, panel, "x", ref_w)
        fy = _frac(deck, panel, "y", ref_h)
        fw = _frac(deck, panel, "w", ref_w)
        fh = _frac(deck, panel, "h", ref_h)
        picture = slide.shapes.add_picture(
            str(path),
            Emu(int(fx * sw_emu)),
            Emu(int(fy * sh_emu)),
            width=Emu(int(fw * sw_emu)),
            height=Emu(int(fh * sh_emu)),
        )
        picture.name = str(panel.get("object_id") or panel.get("panel_id") or f"panel-{slide_no}")
        set_alt_text(picture, panel.get("alt_text"))


def add_icons(slide, specs: list[dict], assets_dir: Path, deck: dict, ref_w: float, ref_h: float, sw_emu: int, sh_emu: int, slide_no: int, svg_assets: list[tuple[int, str, Path]], temporary_files: list[Path]):
    from pptx.util import Emu

    for icon_index, icon in enumerate(specs, 1):
        path = _resolve(assets_dir, icon["file"])
        if not path.exists():
            _die(f"slide {slide_no}: icon not found: {path}")
        fx = _frac(deck, icon, "x", ref_w)
        fy = _frac(deck, icon, "y", ref_h)
        fw = _frac(deck, icon, "w", ref_w)
        fh = _frac(deck, icon, "h", ref_h)
        source_path = path
        if path.suffix.casefold() == ".svg":
            source_path = svg_to_png(path)
            temporary_files.append(source_path)
        picture = slide.shapes.add_picture(
            str(source_path),
            Emu(int(fx * sw_emu)),
            Emu(int(fy * sh_emu)),
            width=Emu(int(fw * sw_emu)),
            height=Emu(int(fh * sh_emu)),
        )
        picture.name = str(icon.get("name") or icon.get("object_id") or f"icon-{icon_index:02d}")
        set_alt_text(picture, icon.get("alt_text"))
        if path.suffix.casefold() == ".svg":
            svg_assets.append((slide_no, picture.name, path))