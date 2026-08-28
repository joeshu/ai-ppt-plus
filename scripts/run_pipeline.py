#!/usr/bin/env python3
"""Run the deterministic AI PPT Plus verification pipeline.

The runner orchestrates existing gates around an already-authored PPTX. It
does not author slides, change formal text, update handoff state, or claim
human approval. All run outputs are isolated under `pipeline-runs/` unless an
explicit output directory is provided.

Usage: run_pipeline.py PROJECT_DIR --deck DECK.pptx --expected-pages N
       [--expected-ratio 1.7777778] [--font-dir DIR]
       [--region name=x,y,w,h ...] [--reference IMAGE | --reference-dir DIR]
       [--visual-threshold N]
       [--ocr-lang LANG] [--require-ocr] [--revision-label R4] [--require-cjk]
       [--route-decision ROUTE.json] [--require-route] [--require-editability]
       [--dpi N] [--strict-layout]
       [--execution-mode dag|linear] [--cache-dir DIR] [--no-cache]
       [--parallel-workers N] [--affected-pages 1,3-4]
       [--page-cache-dir DIR]
       [--affected-region name=x,y,w,h]
       [--output-dir RUN_DIR]
"""
import argparse
import hashlib
import json
import os
import subprocess
import sys
import uuid
from datetime import datetime, timezone
from pathlib import Path

from pipeline_engine import PipelineExecutor, PipelineTask
from report_envelope import normalize_child
from render_review_html import write_review


SCRIPT_DIR = Path(__file__).resolve().parent


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def write_json_atomic(path: Path, value: dict) -> None:
    """Write a run contract atomically so readers never see partial JSON."""
    path = Path(path)
    temporary = path.with_name(f".{path.name}.{uuid.uuid4().hex}.tmp")
    try:
        temporary.write_text(json.dumps(value, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        os.replace(temporary, path)
    finally:
        if temporary.exists():
            temporary.unlink()


STEP_TIMEOUT_SECONDS = 600


def parse_page_selection(expression: str | None, expected_pages: int) -> list[int] | None:
    """Parse and bounds-check a comma/range page selection."""
    if not expression:
        return None
    selected: set[int] = set()
    try:
        for part in expression.split(","):
            token = part.strip()
            if not token:
                raise ValueError
            if "-" in token:
                lo, hi = (int(value.strip()) for value in token.split("-", 1))
                if lo > hi:
                    raise ValueError
                selected.update(range(lo, hi + 1))
            else:
                selected.add(int(token))
    except (TypeError, ValueError):
        raise ValueError("pages must be a comma-separated list of positive integers/ranges")
    if not selected or min(selected) < 1 or max(selected) > expected_pages:
        raise ValueError(f"pages must be between 1 and {expected_pages}")
    return sorted(selected)


def run_step(run_dir: Path, name: str, args, timeout: int = STEP_TIMEOUT_SECONDS):
    stdout_path = run_dir / f"{name}.stdout.txt"
    stderr_path = run_dir / f"{name}.stderr.txt"
    command = [sys.executable, *args]
    try:
        completed = subprocess.run(command, capture_output=True, text=True, timeout=timeout, check=False)
        stdout = completed.stdout
        stderr = completed.stderr
        exit_code = completed.returncode
        failure = None
    except subprocess.TimeoutExpired as exc:
        stdout = exc.stdout or ""
        stderr = (exc.stderr or "") + f"\nstep timed out after {timeout}s"
        exit_code = 124
        failure = "timeout"
    stdout_path.write_text(stdout, encoding="utf-8")
    stderr_path.write_text(stderr, encoding="utf-8")
    result = {"name": name, "command": command, "exit_code": exit_code, "ok": exit_code == 0, "stdout": str(stdout_path.resolve()), "stderr": str(stderr_path.resolve()), "timeout_seconds": timeout, "cache_key": None, "cache_hit": False, "deps": [], "duration_ms": None}
    if failure:
        result["failure"] = failure
    return result


def load_report(path: Path):
    if not path.is_file():
        return None
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception as exc:
        return {"valid": False, "status": "invalid", "issues": [{"code": "invalid_json", "message": f"{type(exc).__name__}: {exc}"}]}


def project_asset_requirements(project: Path) -> tuple[bool, bool]:
    """Read explicit icon/imagegen requirements without substring heuristics."""
    manifest = project / "slide-manifest.json"
    if not manifest.is_file():
        return False, False
    try:
        data = json.loads(manifest.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return False, False
    if not isinstance(data, dict):
        return False, False
    slides = [slide for slide in data.get("slides", []) if isinstance(slide, dict)]
    requires_icon = data.get("requires_icon_assets") is True or any(slide.get("requires_icon_assets") is True for slide in slides)
    requires_imagegen = data.get("requires_imagegen_assets") is True or any(slide.get("requires_imagegen_assets") is True for slide in slides)
    return requires_icon, requires_imagegen


def summarize_report(name: str, path: Path, report: dict):
    summary = normalize_child(name, path, report, required=True, stage=None, deck_sha256=report.get("deck_sha256"))
    summary["report"] = str(path.resolve())
    if name == "render_visual_gate":
        summary.update({"expected_pages": report.get("expected_pages"), "observed_pages": len(report.get("pages", []))})
    elif name == "visual_comparison":
        summary.update({"reference": report.get("reference"), "reference_dir": report.get("reference_dir"), "metrics": report.get("metrics", {}), "aggregate": report.get("aggregate", {}), "compared_pages": len(report.get("pages", [])), "human_visual_review_required": report.get("human_visual_review_required", True)})
    elif name == "ocr_text_check":
        summary.update({"language": report.get("language"), "slide_count": len(report.get("slides", [])), "human_visual_review_required": report.get("human_visual_review_required", True)})
    elif name == "route_validation":
        summary.update({"route": report.get("route"), "visual_authority": report.get("visual_authority"), "formal_content_authority": report.get("formal_content_authority")})
    elif name == "manifest_validation":
        summary.update({"warnings": report.get("warnings", []), "editability_protocol": report.get("editability_protocol"), "editability": report.get("editability", [])})
    elif name == "visual_compare_qa":
        summary.update({"native_status": report.get("status", "diagnostic"), "ok": report.get("ok"), "resized_for_comparison": report.get("resized_for_comparison"), "preview_size": report.get("preview_size")})
    elif name == "render":
        summary.update({"renderer": report.get("renderer"), "dpi": report.get("dpi"), "conversion": report.get("conversion", {}), "page_cache": report.get("page_cache", {}), "page_fingerprints": report.get("page_fingerprints", [])})
    return summary


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("project_dir")
    parser.add_argument("--deck", required=True)
    parser.add_argument("--expected-pages", type=int, required=True)
    parser.add_argument("--expected-ratio", type=float)
    parser.add_argument("--font-dir")
    parser.add_argument("--region", action="append", default=[])
    reference_group = parser.add_mutually_exclusive_group()
    reference_group.add_argument("--reference", help="single approved reference image; only valid for one-page decks")
    reference_group.add_argument("--reference-dir", help="directory containing slide-1.png, slide-2.png, ... for multi-page decks")
    parser.add_argument("--visual-threshold", type=float)
    parser.add_argument("--ocr-lang")
    parser.add_argument("--require-ocr", action="store_true")
    parser.add_argument("--revision-label")
    parser.add_argument("--require-cjk", action="store_true", help="block when the font report cannot support CJK delivery")
    parser.add_argument("--route-decision", help="route-decision.json declaring visual authority")
    parser.add_argument("--require-route", action="store_true", help="require and validate a route decision before downstream gates")
    parser.add_argument("--require-editability", action="store_true", help="require typed L0-L5 object records in the slide manifest")
    parser.add_argument("--require-icon-assets", action="store_true", help="require B4/B5 icon asset and layer audits")
    parser.add_argument("--require-imagegen-assets", action="store_true", help="require per-page imagegen asset provenance")
    parser.add_argument("--object-manifest", help="canonical slide-object-manifest.json")
    parser.add_argument("--require-object-manifest", action="store_true", help="require and validate the canonical object inventory")
    parser.add_argument("--require-independent-panels", action="store_true", help="reverse-audit independently movable semantic panels")
    parser.add_argument("--expected-panel-count", type=int, help="expected semantic panel count")
    parser.add_argument("--require-panel-approval", action="store_true", help="require explicit human approval metadata for panel assets")
    parser.add_argument("--require-text-style-map", action="store_true", help="validate rich text/style records when present")
    parser.add_argument("--text-manifest", help="canonical text-layout-manifest.json")
    parser.add_argument("--require-text-model", action="store_true", help="require and validate the canonical text layout manifest")
    parser.add_argument("--asset-manifest", action="append", default=[], help="asset manifest used for semantic object provenance checks")
    parser.add_argument("--manifest-registry", help="canonical cross-manifest registry.json")
    parser.add_argument("--require-manifest-registry", action="store_true", help="require and validate the cross-manifest registry")
    parser.add_argument("--release", action="store_true", help="run the strict release gate after technical validation")
    parser.add_argument("--handoff", help="handoff.json; required by --release")
    parser.add_argument("--human-signoff", help="human-closeout.json; required by --release")
    parser.add_argument("--issue-log", help="issue-log.json passed to the release gate")
    parser.add_argument("--quality-score", type=float, help="human/automated quality score for --release")
    parser.add_argument("--quality-threshold", type=float, default=80, help="minimum quality score for --release")
    parser.add_argument("--require-embedded-fonts", action="store_true", help="require verified OOXML embedded fonts in strict release delivery")
    parser.add_argument("--dpi", type=int, default=96, help="render DPI; same-ratio reference comparisons are normalized when pixel sizes differ")
    parser.add_argument("--strict-layout", action="store_true", help="treat layout-audit warnings (such as missing source_bbox) as blockers")
    parser.add_argument("--execution-mode", choices=["dag", "linear"], default="dag", help="DAG execution with caching, or the compatibility linear runner")
    parser.add_argument("--cache-dir", help="content-addressed pipeline cache directory; defaults to PROJECT_DIR/.pipeline-cache in DAG mode")
    parser.add_argument("--no-cache", action="store_true", help="disable successful-task cache restores/writes")
    parser.add_argument("--parallel-workers", type=int, default=4, help="maximum independent DAG checks to run concurrently")
    parser.add_argument("--affected-pages", help="only render and compare selected pages, e.g. 1,3-4")
    parser.add_argument("--page-cache-dir", help="content-addressed validated page PNG cache; defaults to .pipeline-cache/render-pages in DAG mode")
    parser.add_argument("--affected-region", action="append", default=[], help="critical region affected by the change: name=x,y,w,h; checked by the render QA gate")
    parser.add_argument("--output-dir")
    args = parser.parse_args()
    project = Path(args.project_dir).resolve()
    deck = Path(args.deck).resolve()
    if args.release:
        # A release is a stronger profile than technical validation.  Make
        # the required evidence explicit instead of allowing a green run to
        # be mistaken for a delivered deck.
        args.require_route = True
        args.require_editability = True
        args.require_embedded_fonts = True
        args.require_cjk = True
        missing = []
        if not args.font_dir:
            missing.append("--font-dir")
        if not args.route_decision:
            missing.append("--route-decision")
        if not args.handoff:
            missing.append("--handoff")
        if not args.human_signoff:
            missing.append("--human-signoff")
        if args.quality_score is None:
            missing.append("--quality-score")
        if missing:
            result = {"schema": "ai-ppt-plus/pipeline-run/v2", "valid": False, "technical_valid": False, "release_eligible": False, "code": "release_evidence_missing", "missing": missing}
            print(json.dumps(result, ensure_ascii=False))
            return 2
    if not project.is_dir() or not deck.is_file():
        print(json.dumps({"valid": False, "code": "project_or_deck_missing"}, ensure_ascii=False))
        return 3
    if args.release and not (project / "slide-manifest.json").is_file():
        result = {"schema": "ai-ppt-plus/pipeline-run/v2", "valid": False, "technical_valid": False, "release_eligible": False, "code": "slide_manifest_missing", "message": "--release requires project/slide-manifest.json"}
        print(json.dumps(result, ensure_ascii=False))
        return 2
    if args.reference and args.expected_pages != 1:
        result = {"schema": "ai-ppt-plus/pipeline-run/v1", "valid": False, "code": "single_reference_for_multipage", "message": "Use --reference-dir with slide-N.png files for multi-page decks"}
        print(json.dumps(result, ensure_ascii=False))
        return 2
    try:
        affected_pages = parse_page_selection(args.affected_pages, args.expected_pages)
    except ValueError as exc:
        print(json.dumps({"schema": "ai-ppt-plus/pipeline-run/v2", "valid": False, "code": "affected_pages_invalid", "message": str(exc)}, ensure_ascii=False))
        return 2
    if args.parallel_workers < 1:
        print(json.dumps({"schema": "ai-ppt-plus/pipeline-run/v2", "valid": False, "code": "parallel_workers_invalid"}, ensure_ascii=False))
        return 2
    if args.release and affected_pages:
        print(json.dumps({"schema": "ai-ppt-plus/pipeline-run/v2", "valid": False, "technical_valid": False, "release_eligible": False, "code": "release_requires_full_deck", "message": "--release requires a full-deck render; omit --affected-pages"}, ensure_ascii=False))
        return 2
    if args.requ…11825 tokens truncated… list):
        for child in value:
            records.extend(_asset_records(child))
    return records


def _record_ids(record: dict[str, Any]) -> set[str]:
    return {str(record[key]) for key in ("asset_id", "panel_id", "icon_id", "object_id") if record.get(key)}


def _record_paths(record: dict[str, Any]) -> set[str]:
    values = set()
    for key in ("path", "file", "source", "source_path", "copied_to"):
        value = record.get(key)
        if isinstance(value, str) and value:
            values.add(value.split("#", 1)[0])
    return values


def _values(value: Any) -> list[Any]:
    if isinstance(value, (list, tuple, set)):
        return list(value)
    return [value] if value else []


def _hash_candidates(obj: dict[str, Any], assets: list[dict[str, Any]], base: Path) -> tuple[set[str], list[Path]]:
    details = obj.get("details") if isinstance(obj.get("details"), dict) else {}
    hashes = set()
    paths: list[Path] = []
    for source in (obj, details):
        for key in ("source_hash", "path_sha256", "sha256"):
            value = source.get(key)
            if isinstance(value, str) and SHA256.fullmatch(value):
                hashes.add(value.lower())
        for key in ("source_path", "path", "file", "copied_to"):
            path = _resolve(base, source.get(key))
            if path is not None:
                paths.append(path)

    wanted_ids = {str(value) for value in _values(obj.get("asset_ids")) if value}
    if obj.get("embedded_asset"):
        wanted_ids.add(str(obj["embedded_asset"]))
    wanted_paths = {str(value).split("#", 1)[0] for value in _values(obj.get("source_paths")) if value}
    wanted_paths.update(str(value).split("#", 1)[0] for value in _values(obj.get("source_path")) if value)
    for record in assets:
        if wanted_ids and not wanted_ids.intersection(_record_ids(record)):
            continue
        record_paths = _record_paths(record)
        if wanted_paths and not wanted_paths.intersection(record_paths):
            continue
        if not wanted_ids and not wanted_paths:
            continue
        for key in ("source_hash", "path_sha256", "sha256"):
            value = record.get(key)
            if isinstance(value, str) and SHA256.fullmatch(value):
                hashes.add(value.lower())
        for key in ("path", "file", "source", "source_path", "copied_to"):
            path = _resolve(base, record.get(key))
            if path is not None:
                paths.append(path)
    unique_paths = []
    seen = set()
    for path in paths:
        resolved = path.resolve()
        if str(resolved) not in seen:
            seen.add(str(resolved))
            unique_paths.append(resolved)
    return hashes, unique_paths


def _shape_media(shape):
    """Return (package part name, bytes) for a picture shape, when available."""
    from pptx.oxml.ns import qn

    for element in shape._element.iter():
        if element.tag != qn("a:blip"):
            continue
        relationship_id = element.get(qn("r:embed"))
        relationship = shape.part.rels.get(relationship_id) if relationship_id else None
        if relationship is None or not hasattr(relationship, "target_part"):
            return None, None
        part = relationship.target_part
        return str(part.partname).lstrip("/"), part.blob
    return None, None


def _actual_kind(shape) -> str:
    if getattr(shape, "has_table", False):
        return "editable_table"
    if getattr(shape, "has_chart", False):
        return "editable_chart"
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return "native_group"
    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        return "native_shape"
    if getattr(shape, "has_text_frame", False):
        return "native_text"
    return "other"


def _text_from_shape(shape) -> str:
    return _normal_text(getattr(shape, "text", ""))


def _table_evidence(shape) -> dict[str, Any]:
    table = shape.table
    rows = [[_text_from_shape(cell) for cell in row.cells] for row in table.rows]
    return {
        "rows": len(table.rows),
        "columns": len(table.columns),
        "nonempty_cells": sum(bool(cell.strip()) for row in rows for cell in row),
        "values": rows,
    }


def _chart_evidence(shape) -> dict[str, Any]:
    chart = shape.chart
    series_count = len(chart.series)
    nonempty_series = 0
    for series in chart.series:
        try:
            values = list(series.values)
        except Exception:
            values = []
        if values:
            nonempty_series += 1
    chart_part = None
    for relationship in shape.part.rels.values():
        if not hasattr(relationship, "target_part"):
            continue
        part_name = str(relationship.target_part.partname).lstrip("/")
        if part_name.startswith("ppt/charts/"):
            chart_part = relationship.target_part
            break
    chart_xml = chart_part.blob if chart_part is not None else b""
    embedded_workbook = False
    if chart_part is not None:
        embedded_workbook = any(
            hasattr(rel, "target_part") and "embeddings/" in str(rel.target_part.partname)
            for rel in chart_part.rels.values()
        )
    return {
        "series": series_count,
        "nonempty_series": nonempty_series,
        "has_cached_values": b"<c:numCache" in chart_xml or b"<c:strCache" in chart_xml,
        "has_embedded_workbook": embedded_workbook,
        "chart_part": str(chart_part.partname).lstrip("/") if chart_part is not None else None,
    }


def _expected_text(obj: dict[str, Any], text_specs: dict[str, dict[str, Any]]) -> str | None:
    spec = obj.get("text_spec") if isinstance(obj.get("text_spec"), dict) else None
    if spec is None:
        spec = text_specs.get(str(obj.get("text_id") or obj.get("object_id")))
    if spec is None or "content" not in spec:
        return None
    return _normal_text(spec.get("content"))


def _type_matches(expected: str, actual: str) -> bool:
    if expected == "editable_text":
        return actual == "native_text"
    if expected == "editable_table":
        return actual == "editable_table"
    if expected == "editable_chart":
        return actual == "editable_chart"
    if expected in {"native_shape", "native_group"}:
        return actual == expected
    if expected == "editable_vector":
        return actual == "picture"
    if expected in {"independent_image", "extracted_icon", "traceable_static_graphic"}:
        return actual == "picture"
    return True


def audit(deck_path: Path, object_manifest_path: Path, text_manifest_path: Path | None = None, asset_manifest_paths: list[Path] | None = None) -> dict[str, Any]:
    from pptx import Presentation

    object_manifest = _read(object_manifest_path)
    text_specs: dict[str, dict[str, Any]] = {}
    if text_manifest_path and text_manifest_path.is_file():
        text_manifest = _read(text_manifest_path)
        for slide in text_manifest.get("slides", []):
            for spec in slide.get("text_specs", []) if isinstance(slide, dict) else []:
                if isinstance(spec, dict) and spec.get("text_id"):
                    text_specs[str(spec["text_id"])] = spec
    asset_records: list[dict[str, Any]] = []
    for path in asset_manifest_paths or []:
        if path.is_file():
            asset_records.extend(_asset_records(_read(path)))

    prs = Presentation(str(deck_path))
    errors: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []
    audited: list[dict[str, Any]] = []
    expected_count = 0
    seen_expected = set()
    for slide_index, slide_spec in enumerate(object_manifest.get("slides", []), 1):
        if not isinstance(slide_spec, dict):
            continue
        slide_no = int(slide_spec.get("slide_no", slide_index))
        if slide_no < 1 or slide_no > len(prs.slides):
            errors.append({"code": "manifest_slide_missing", "slide_no": slide_no})
            continue
        shapes = list(prs.slides[slide_no - 1].shapes)
        by_name: dict[str, list[Any]] = {}
        for shape in shapes:
            by_name.setdefault(shape.name, []).append(shape)
        for obj in slide_spec.get("objects", []):
            if not isinstance(obj, dict):
                continue
            expected_count += 1
            object_id = str(obj.get("object_id", ""))
            seen_expected.add((slide_no, object_id))
            matches = by_name.get(object_id, [])
            if not matches:
                errors.append({"code": "semantic_object_missing", "slide_no": slide_no, "object_id": object_id})
                continue
            if len(matches) != 1:
                errors.append({"code": "semantic_object_name_not_unique", "slide_no": slide_no, "object_id": object_id, "observed": len(matches)})
                continue
            shape = matches[0]
            actual = _actual_kind(shape)
            expected_type = str(obj.get("object_type", ""))
            record: dict[str, Any] = {
                "slide_no": slide_no,
                "object_id": object_id,
                "role": obj.get("role"),
                "expected_type": expected_type,
                "actual_type": actual,
                "shape_type": str(shape.shape_type),
                "text": _text_from_shape(shape),
                "semantic_checks": {},
            }
            if not _type_matches(expected_type, actual):
                errors.append({"code": "semantic_type_mismatch", "slide_no": slide_no, "object_id": object_id, "expected": expected_type, "actual": actual})
            record["semantic_checks"]["type"] = actual == "native_text" if expected_type == "editable_text" else _type_matches(expected_type, actual)

            expected_text = _expected_text(obj, text_specs)
            if expected_text is not None:
                observed_text = _text_from_shape(shape)
                matches_text = expected_text == observed_text
                record["semantic_checks"]["text_exact"] = matches_text
                if not matches_text:
                    errors.append({"code": "pptx_text_manifest_mismatch", "slide_no": slide_no, "object_id": object_id, "expected": expected_text, "observed": observed_text})

            if actual == "editable_table":
                table = _table_evidence(shape)
                record["table"] = table
                valid_table = table["rows"] > 0 and table["columns"] > 0
                record["semantic_checks"]["native_table_data"] = valid_table
                if not valid_table:
                    errors.append({"code": "native_table_empty", "slide_no": slide_no, "object_id": object_id})
            if actual == "editable_chart":
                chart = _chart_evidence(shape)
                record["chart"] = chart
                valid_chart = chart["series"] > 0 and chart["nonempty_series"] > 0 and chart["has_cached_values"] and chart["has_embedded_workbook"]
                record["semantic_checks"]["native_chart_data"] = valid_chart
                if not valid_chart:
                    errors.append({"code": "native_chart_data_missing", "slide_no": slide_no, "object_id": object_id, "evidence": chart})

            is_brand_lockup = obj.get("role") in {"brand_lockup", "logo", "brand-logo"} or obj.get("asset_policy") == "brand_lockup" or (isinstance(obj.get("details"), dict) and obj["details"].get("asset_policy") == "brand_lockup")
            if is_brand_lockup:
                brand_ok = actual == "picture" and not _text_from_shape(shape)
                record["semantic_checks"]["brand_lockup_whole_asset"] = brand_ok
                if not brand_ok:
                    errors.append({"code": "brand_lockup_not_whole_picture", "slide_no": slide_no, "object_id": object_id, "actual": actual})

            if actual == "picture":
                part_name, blob = _shape_media(shape)
                record["media_part"] = part_name
                record["embedded_sha256"] = hashlib.sha256(blob).hexdigest() if blob is not None else None
                expected_hashes, source_paths = _hash_candidates(obj, asset_records, object_manifest_path.parent)
                for source_path in source_paths:
                    if source_path.is_file():
                        source_hash = _digest(source_path)
                        expected_hashes.add(source_hash)
                    else:
                        warnings.append({"code": "source_file_unavailable", "slide_no": slide_no, "object_id": object_id, "path": str(source_path)})
                if expected_hashes and blob is not None:
                    hash_ok = hashlib.sha256(blob).hexdigest().lower() in expected_hashes
                    record["semantic_checks"]["source_hash"] = hash_ok
                    if not hash_ok:
                        errors.append({"code": "embedded_asset_hash_mismatch", "slide_no": slide_no, "object_id": object_id, "expected": sorted(expected_hashes), "observed": record["embedded_sha256"]})
                elif expected_hashes and blob is None:
                    errors.append({"code": "embedded_asset_unreadable", "slide_no": slide_no, "object_id": object_id})

            audited.append(record)

    result = {
        "schema": "ai-ppt-plus/semantic-object-audit/v1",
        "valid": not errors,
        "status": "passed" if not errors else "blocked",
        "deck": str(deck_path.resolve()),
        "deck_sha256": _digest(deck_path),
        "object_manifest": str(object_manifest_path.resolve()),
        "text_manifest": str(text_manifest_path.resolve()) if text_manifest_path else None,
        "expected_object_count": expected_count,
        "audited_object_count": len(audited),
        "objects": audited,
        "errors": errors,
        "warnings": warnings,
        "human_visual_review_required": True,
    }
    return result


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("deck")
    parser.add_argument("--object-manifest", required=True)
    parser.add_argument("--text-manifest")
    parser.add_argument("--asset-manifest", action="append", default=[])
    parser.add_argument("--report")
    args = parser.parse_args()
    try:
        result = audit(
            Path(args.deck).resolve(),
            Path(args.object_manifest).resolve(),
            Path(args.text_manifest).resolve() if args.text_manifest else None,
            [Path(path).resolve() for path in args.asset_manifest],
        )
    except Exception as exc:
        result = {"schema": "ai-ppt-plus/semantic-object-audit/v1", "valid": False, "status": "invalid", "errors": [{"code": "runtime_error", "message": f"{type(exc).__name__}: {exc}"}]}
    if args.report:
        report = Path(args.report).resolve()
        report.parent.mkdir(parents=True, exist_ok=True)
        report.write_text(json.dumps(result, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(json.dumps(result, ensure_ascii=False))
    return 0 if result.get("valid") is True else 2


if __name__ == "__main__":
    raise SystemExit(main())