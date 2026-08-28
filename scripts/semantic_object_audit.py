#!/usr/bin/env python3
"""Audit the semantic meaning and provenance of final PPTX objects.

The identity audit answers "does a shape with this name exist?".  This gate
answers the stronger questions required for delivery: is the object the
declared native type, is its data still present and source-consistent, does
the manifest cover the visible slide, and can embedded media be traced back to
the declared source bytes?
"""
from __future__ import annotations

import argparse
import csv
import hashlib
import io
import json
import math
import re
from pathlib import Path
from typing import Any


SCHEMA = "ai-ppt-plus/semantic-object-audit/v2"
SHA256 = re.compile(r"^[0-9a-fA-F]{64}$")
SOURCE_HASH_KEYS = ("source_hash", "source_sha256", "path_sha256", "asset_sha256", "file_sha256", "sha256")
SOURCE_PATH_KEYS = ("source_path", "path", "file", "source", "copied_to", "asset_path", "source_file")


def _read(path: Path) -> dict[str, Any]:
    value = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(value, dict):
        raise ValueError(f"manifest must be an object: {path}")
    return value


def _digest(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _json_digest(value: Any) -> str:
    return hashlib.sha256(
        json.dumps(value, ensure_ascii=False, sort_keys=True, separators=(",", ":")).encode("utf-8")
    ).hexdigest()


def _normal_text(value: Any) -> str:
    return str(value or "").replace("\r\n", "\n").replace("\r", "\n").replace("\v", "\n")


def _normal_table_cell(value: Any) -> str:
    if value is None:
        return ""
    return _normal_text(value)


def _normal_number(value: Any) -> float | str | None:
    if value is None:
        return None
    if isinstance(value, bool):
        return float(value)
    if isinstance(value, (int, float)):
        number = float(value)
        return number if math.isfinite(number) else None
    try:
        number = float(str(value).strip())
    except (TypeError, ValueError):
        return _normal_text(value)
    return number if math.isfinite(number) else None


def _resolve(base: Path, value: Any) -> Path | None:
    if not isinstance(value, str) or not value:
        return None
    candidate = value.split("#", 1)[0].strip()
    if not candidate or "://" in candidate or candidate.startswith("asset:"):
        return None
    path = Path(candidate)
    return path if path.is_absolute() else base / path


def _asset_records(value: Any, *, base: Path) -> list[dict[str, Any]]:
    records: list[dict[str, Any]] = []
    if isinstance(value, dict):
        if any(key in value for key in ("asset_id", "panel_id", "icon_id", "object_id", "path", "file", "copied_to")):
            record = dict(value)
            record["__manifest_base"] = str(base)
            records.append(record)
        for child in value.values():
            records.extend(_asset_records(child, base=base))
    elif isinstance(value, list):
        for child in value:
            records.extend(_asset_records(child, base=base))
    return records


def _record_ids(record: dict[str, Any]) -> set[str]:
    return {str(record[key]) for key in ("asset_id", "panel_id", "icon_id", "object_id") if record.get(key)}


def _record_paths(record: dict[str, Any]) -> set[str]:
    values = set()
    for key in SOURCE_PATH_KEYS:
        value = record.get(key)
        if isinstance(value, str) and value:
            values.add(value.split("#", 1)[0])
    return values


def _values(value: Any) -> list[Any]:
    if isinstance(value, (list, tuple, set)):
        return list(value)
    return [value] if value else []


def _hash_candidates(obj: dict[str, Any], assets: list[dict[str, Any]], base: Path) -> tuple[set[str], list[Path], list[str]]:
    """Collect declared hashes and source files for one object.

    Asset-manifest records are selected by object/asset ID or source path.  We
    intentionally never collect every hash in a manifest: an unrelated asset
    must not accidentally make a wrong embedded image pass.
    """
    details = obj.get("details") if isinstance(obj.get("details"), dict) else {}
    hashes: set[str] = set()
    paths: list[Path] = []
    invalid: list[str] = []

    def collect(source: dict[str, Any], source_base: Path) -> None:
        for key in SOURCE_HASH_KEYS:
            value = source.get(key)
            if isinstance(value, str) and SHA256.fullmatch(value):
                hashes.add(value.lower())
            elif value not in (None, ""):
                invalid.append(f"{key}={value}")
        for key in SOURCE_PATH_KEYS:
            path = _resolve(source_base, source.get(key))
            if path is not None:
                paths.append(path)

    collect(obj, base)
    collect(details, base)
    wanted_ids = {str(value) for value in _values(obj.get("asset_ids")) if value}
    if obj.get("embedded_asset"):
        wanted_ids.add(str(obj["embedded_asset"]))
    if obj.get("object_id"):
        wanted_ids.add(str(obj["object_id"]))
    wanted_paths = {str(value).split("#", 1)[0] for value in _values(obj.get("source_paths")) if value}
    wanted_paths.update(str(value).split("#", 1)[0] for value in _values(obj.get("source_path")) if value)

    for record in assets:
        record_ids = _record_ids(record)
        record_paths = _record_paths(record)
        if wanted_ids and not wanted_ids.intersection(record_ids):
            if wanted_paths and wanted_paths.intersection(record_paths):
                pass
            else:
                continue
        elif wanted_paths and not wanted_paths.intersection(record_paths):
            continue
        elif not wanted_ids and not wanted_paths:
            continue
        source_base = Path(str(record.get("__manifest_base") or base))
        collect(record, source_base)

    unique_paths: list[Path] = []
    seen: set[str] = set()
    for path in paths:
        resolved = path.resolve()
        if str(resolved) not in seen:
            seen.add(str(resolved))
            unique_paths.append(resolved)
    return hashes, unique_paths, invalid


def _shape_media(shape):
    """Return (package part name, bytes) for a picture shape, when available."""
    from pptx.oxml.ns import qn

    for element in shape._element.iter():
        if element.tag != qn("a:blip"):
            continue
        relationship_id = element.get(qn("r:embed"))
        relationship = shape.part.rels.get(relationship_id) if relationship_id else None
        if relationship is None or not hasattr(relationship, "target_part"):
            return None, None
        part = relationship.target_part
        return str(part.partname).lstrip("/"), part.blob
    return None, None


def _shape_type(shape):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    return getattr(shape, "shape_type", None), MSO_SHAPE_TYPE


def _actual_kind(shape) -> str:
    shape_type, types = _shape_type(shape)
    if getattr(shape, "has_table", False):
        return "editable_table"
    if getattr(shape, "has_chart", False):
        return "editable_chart"
    if shape_type == types.PICTURE:
        return "picture"
    if shape_type == types.GROUP:
        return "native_group"
    if shape_type == types.AUTO_SHAPE:
        return "native_shape"
    if shape_type == types.TEXT_BOX and getattr(shape, "has_text_frame", False):
        return "native_text"
    return "native_text" if getattr(shape, "has_text_frame", False) else "other"


def _is_native_textbox(shape) -> bool:
    _, types = _shape_type(shape)
    return bool(
        getattr(shape, "has_text_frame", False)
        and getattr(shape, "shape_type", None) == types.TEXT_BOX
        and not bool(getattr(shape, "is_placeholder", False))
    )


def _text_from_shape(shape) -> str:
    return _normal_text(getattr(shape, "text", ""))


def _text_run_evidence(shape) -> list[str]:
    if not getattr(shape, "has_text_frame", False):
        return []
    runs: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            runs.append(_normal_text(run.text))
    return runs


def _table_evidence(shape) -> dict[str, Any]:
    table = shape.table
    values = [[_normal_table_cell(cell.text) for cell in row.cells] for row in table.rows]
    columns = len(table.columns)
    rows = len(table.rows)
    return {
        "rows": rows,
        "columns": columns,
        "rectangular": all(len(row) == columns for row in values),
        "nonempty_cells": sum(bool(cell.strip()) for row in values for cell in row),
        "values": values,
        "values_sha256": _json_digest(values),
    }


def _rectangular_rows(value: Any) -> list[list[Any]] | None:
    if not isinstance(value, list) or not value or any(not isinstance(row, (list, tuple)) for row in value):
        return None
    width = max((len(row) for row in value), default=0)
    if width <= 0:
        return None
    return [[row[index] if index < len(row) else "" for index in range(width)] for row in value]


def _trim_empty_edges(rows: list[list[Any]] | None) -> list[list[Any]] | None:
    if not rows:
        return rows
    trimmed = [list(row) for row in rows]
    while trimmed and all(value is None or str(value) == "" for value in trimmed[-1]):
        trimmed.pop()
    if not trimmed:
        return None
    width = max((len(row) for row in trimmed), default=0)
    while width and all(len(row) < width or row[width - 1] is None or str(row[width - 1]) == "" for row in trimmed):
        width -= 1
    if width <= 0:
        return None
    return [row[:width] for row in trimmed]


def _table_matrix(value: Any, merges: Any = None) -> list[list[str]] | None:
    if isinstance(value, dict):
        merges = value.get("merges", merges)
        value = value.get("values", value.get("rows"))
    rows = _rectangular_rows(value)
    if rows is None:
        return None
    output = [[_normal_table_cell(cell) for cell in row] for row in rows]
    for merge in merges or []:
        if not isinstance(merge, list) or len(merge) != 4:
            continue
        r1, c1, r2, c2 = [int(part) for part in merge]
        for row in range(max(0, r1), min(len(output), r2 + 1)):
            for column in range(max(0, c1), min(len(output[row]), c2 + 1)):
                if row != r1 or column != c1:
                    output[row][column] = ""
    return output


def _workbook_rows(blob: bytes) -> tuple[list[list[Any]] | None, str | None]:
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(io.BytesIO(blob), data_only=False, read_only=True)
        if not workbook.sheetnames:
            return None, "workbook_has_no_sheets"
        sheet = workbook[workbook.sheetnames[0]]
        rows = [list(row) for row in sheet.iter_rows(values_only=True)]
        workbook.close()
        return _trim_empty_edges(_rectangular_rows(rows)), None
    except ImportError:
        return None, "openpyxl_unavailable"
    except Exception as exc:
        return None, f"workbook_parse_failed: {type(exc).__name__}: {exc}"


def _matrix_from_source(path: Path) -> tuple[list[list[Any]] | None, str | None]:
    suffix = path.suffix.casefold()
    try:
        if suffix in {".xlsx", ".xlsm", ".xltx", ".xltm"}:
            return _workbook_rows(path.read_bytes())
        if suffix in {".csv", ".tsv"}:
            delimiter = "\t" if suffix == ".tsv" else ","
            with path.open("r", encoding="utf-8-sig", newline="") as stream:
                return _trim_empty_edges(_rectangular_rows(list(csv.reader(stream, delimiter=delimiter)))), None
        if suffix == ".json":
            data = json.loads(path.read_text(encoding="utf-8"))
            if isinstance(data, dict):
                data = data.get("rows", data.get("values"))
            return _trim_empty_edges(_rectangular_rows(data)), None
    except Exception as exc:
        return None, f"data_source_parse_failed: {type(exc).__name__}: {exc}"
    return None, "unsupported_data_source_format"


def _chart_from_matrix(rows: list[list[Any]] | None) -> dict[str, Any] | None:
    if not rows or len(rows) < 2 or len(rows[0]) < 2:
        return None
    width = len(rows[0])
    categories = [_normal_text(row[0]) for row in rows[1:] if row]
    series: list[dict[str, Any]] = []
    for column in range(1, width):
        name = _normal_text(rows[0][column]) or "Series"
        values = []
        for row in rows[1:]:
            value = row[column] if column < len(row) else None
            number = _normal_number(value)
            if number is None:
                return None
            values.append(number)
        series.append({"name": name, "values": values})
    return {"kind": "category_chart", "categories": categories, "series": series}


def _chart_snapshot(categories: Any, series: Any) -> dict[str, Any] | None:
    if not isinstance(categories, (list, tuple)) or not categories:
        return None
    if not isinstance(series, (list, tuple)) or not series:
        return None
    output_series = []
    for item in series:
        if not isinstance(item, dict):
            return None
        raw_values = item.get("values")
        if not isinstance(raw_values, (list, tuple)) or len(raw_values) != len(categories):
            return None
        values = []
        for value in raw_values:
            number = _normal_number(value)
            if number is None:
                return None
            values.append(number)
        output_series.append({"name": _normal_text(item.get("name", "Series")), "values": values})
    return {
        "kind": "category_chart",
        "categories": [_normal_text(value) for value in categories],
        "series": output_series,
    }


def _chart_from_shape(shape) -> dict[str, Any] | None:
    chart = shape.chart
    try:
        categories = [getattr(category, "label", category) for category in chart.plots[0].categories]
        series = [{"name": item.name, "values": list(item.values)} for item in chart.series]
    except Exception:
        return None
    return _chart_snapshot(categories, series)


def _chart_evidence(shape) -> dict[str, Any]:
    chart = shape.chart
    chart_part = getattr(chart, "part", None)
    chart_xml = getattr(chart_part, "blob", b"") if chart_part is not None else b""
    workbook_part = None
    workbook_blob = None
    if chart_part is not None:
        for relationship in chart_part.rels.values():
            if not hasattr(relationship, "target_part"):
                continue
            target = relationship.target_part
            part_name = str(target.partname).lstrip("/")
            if "embeddings/" in part_name and part_name.casefold().endswith((".xlsx", ".xlsm")):
                workbook_part = part_name
                workbook_blob = target.blob
                break
    workbook_data = None
    workbook_error = None
    if workbook_blob is not None:
        rows, workbook_error = _workbook_rows(workbook_blob)
        workbook_data = _chart_from_matrix(rows)
    chart_data = _chart_from_shape(shape)
    nonempty_series = 0
    for item in chart.series:
        try:
            if list(item.values):
                nonempty_series += 1
        except Exception:
            continue
    return {
        "series": len(chart.series),
        "nonempty_series": nonempty_series,
        "chart_data": chart_data,
        "chart_data_sha256": _json_digest(chart_data) if chart_data is not None else None,
        "has_cached_values": b"<c:numCache" in chart_xml or b"<c:strCache" in chart_xml,
        "has_embedded_workbook": workbook_blob is not None,
        "workbook_part": workbook_part,
        "workbook_sha256": hashlib.sha256(workbook_blob).hexdigest() if workbook_blob is not None else None,
        "workbook_data": workbook_data,
        "workbook_data_sha256": _json_digest(workbook_data) if workbook_data is not None else None,
        "workbook_error": workbook_error,
        "chart_workbook_match": chart_data is not None and workbook_data is not None and chart_data == workbook_data,
    }


def _expected_text(obj: dict[str, Any], text_specs: dict[str, dict[str, Any]]) -> str | None:
    spec = obj.get("text_spec") if isinstance(obj.get("text_spec"), dict) else None
    if spec is None:
        spec = text_specs.get(str(obj.get("text_id") or obj.get("object_id")))
    if spec is not None and "content" in spec:
        return _normal_text(spec.get("content"))
    for key in ("content", "text"):
        if key in obj:
            return _normal_text(obj.get(key))
    return None


def _expected_data_source(obj: dict[str, Any], base: Path, *, chart: bool) -> Any:
    source = _resolve(base, obj.get("data_source_path") or obj.get("data_source"))
    if source is None or not source.is_file():
        return None
    if source.suffix.casefold() == ".json":
        try:
            value = json.loads(source.read_text(encoding="utf-8"))
            if chart and isinstance(value, dict) and "categories" in value and "series" in value:
                return _chart_snapshot(value["categories"], value["series"])
        except Exception:
            return None
    rows, _ = _matrix_from_source(source)
    if chart:
        return _chart_from_matrix(rows)
    return _table_matrix(rows)


def _expected_table_values(obj: dict[str, Any], base: Path) -> list[list[str]] | None:
    snapshot = obj.get("data_snapshot") or obj.get("table_data") or obj.get("data")
    values = _table_matrix(snapshot, obj.get("merges")) if snapshot is not None else None
    if values is None and isinstance(obj.get("rows"), list):
        values = _table_matrix(obj.get("rows"), obj.get("merges"))
    if values is None:
        values = _expected_data_source(obj, base, chart=False)
    return values


def _expected_chart_values(obj: dict[str, Any], base: Path) -> dict[str, Any] | None:
    snapshot = obj.get("data_snapshot") or obj.get("chart_data") or obj.get("data")
    if isinstance(snapshot, dict) and "categories" in snapshot and "series" in snapshot:
        return _chart_snapshot(snapshot["categories"], snapshot["series"])
    if "categories" in obj and "series" in obj:
        return _chart_snapshot(obj.get("categories"), obj.get("series"))
    source = _expected_data_source(obj, base, chart=True)
    return source if isinstance(source, dict) else None


def _data_source_evidence(obj: dict[str, Any], base: Path) -> dict[str, Any]:
    explicit_value = obj.get("data_source_path") or obj.get("data_source")
    source = _resolve(base, explicit_value)
    declared = obj.get("data_source_sha256")
    evidence = {
        "declared_sha256": declared.lower() if isinstance(declared, str) and SHA256.fullmatch(declared) else None,
        "path": str(source.resolve()) if source is not None else None,
        "available": bool(source and source.is_file()),
        "explicit_path": bool(obj.get("data_source_path")),
        "observed_sha256": _digest(source) if source and source.is_file() else None,
    }
    evidence["hash_match"] = (
        evidence["declared_sha256"] is None
        or evidence["observed_sha256"] is None
        or evidence["declared_sha256"] == evidence["observed_sha256"].lower()
    )
    return evidence


def _type_matches(expected: str, actual: str) -> bool:
    if expected == "editable_text":
        return actual == "native_text"
    if expected == "editable_table":
        return actual == "editable_table"
    if expected == "editable_chart":
        return actual == "editable_chart"
    if expected in {"native_shape", "native_group"}:
        return actual == expected
    if expected == "editable_vector":
        return actual == "picture"
    if expected in {"independent_image", "extracted_icon", "traceable_static_graphic", "decorative_art"}:
        return actual == "picture"
    return True


def _walk_shapes(shapes, parent: str | None = None):
    for shape in shapes:
        yield shape, parent
        if _actual_kind(shape) == "native_group":
            yield from _walk_shapes(shape.shapes, shape.name)


def _shape_box(shape) -> tuple[float, float, float, float]:
    return (float(shape.left), float(shape.top), float(shape.width), float(shape.height))


def _overlap_ratio(first, second) -> float:
    ax, ay, aw, ah = _shape_box(first)
    bx, by, bw, bh = _shape_box(second)
    left = max(ax, bx)
    top = max(ay, by)
    right = min(ax + aw, bx + bw)
    bottom = min(ay + ah, by + bh)
    intersection = max(0.0, right - left) * max(0.0, bottom - top)
    smaller = min(max(0.0, aw * ah), max(0.0, bw * bh))
    return intersection / smaller if smaller else 0.0


def _picture_crop_evidence(shape) -> dict[str, Any]:
    from pptx.oxml.ns import qn

    crop = {key: 0 for key in ("left", "top", "right", "bottom")}
    for element in shape._element.iter(qn("a:srcRect")):
        for key in crop:
            try:
                crop[key] = int(element.get(key, "0"))
            except (TypeError, ValueError):
                crop[key] = -1
        break
    return {"values": crop, "cropped": any(value != 0 for value in crop.values())}


def _is_brand_lockup(obj: dict[str, Any]) -> bool:
    details = obj.get("details") if isinstance(obj.get("details"), dict) else {}
    return bool(
        obj.get("role") in {"brand_lockup", "logo", "brand-logo"}
        or obj.get("asset_policy") == "brand_lockup"
        or details.get("asset_policy") == "brand_lockup"
    )


def _requires_source_hash(obj: dict[str, Any], brand: bool) -> bool:
    expected = str(obj.get("object_type") or "")
    return brand or expected in {"independent_image", "extracted_icon", "editable_vector", "traceable_static_graphic", "decorative_art"}


def _append_error(errors: list[dict[str, Any]], code: str, slide_no: int, object_id: str, **extra: Any) -> None:
    errors.append({"code": code, "slide_no": slide_no, "object_id": object_id, **extra})


def audit(
    deck_path: Path,
    object_manifest_path: Path,
    text_manifest_path: Path | None = None,
    asset_manifest_paths: list[Path] | None = None,
) -> dict[str, Any]:
    from pptx import Presentation

    object_manifest = _read(object_manifest_path)
    text_specs: dict[str, dict[str, Any]] = {}
    text_manifest_entries: list[dict[str, Any]] = []
    if text_manifest_path and text_manifest_path.is_file():
        text_manifest = _read(text_manifest_path)
        for slide_index, slide in enumerate(text_manifest.get("slides", []), 1):
            for spec in slide.get("text_specs", []) if isinstance(slide, dict) else []:
                if isinstance(spec, dict):
                    text_manifest_entries.append({"slide_no": int(slide.get("slide_no", slide_index)), **spec})
                    for key in (spec.get("text_id"), spec.get("object_id")):
                        if key:
                            text_specs[str(key)] = spec
    asset_records: list[dict[str, Any]] = []
    for path in asset_manifest_paths or []:
        if path.is_file():
            asset_records.extend(_asset_records(_read(path), base=path.parent.resolve()))

    prs = Presentation(str(deck_path))
    errors: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []
    audited: list[dict[str, Any]] = []
    coverage: list[dict[str, Any]] = []
    expected_count = 0

    for slide_index, slide_spec in enumerate(object_manifest.get("slides", []), 1):
        if not isinstance(slide_spec, dict):
            continue
        slide_no = int(slide_spec.get("slide_no", slide_index))
        if slide_no < 1 or slide_no > len(prs.slides):
            errors.append({"code": "manifest_slide_missing", "slide_no": slide_no})
            continue
        slide = prs.slides[slide_no - 1]
        entries = list(_walk_shapes(slide.shapes))
        by_name: dict[str, list[tuple[Any, str | None]]] = {}
        for shape, parent in entries:
            by_name.setdefault(shape.name, []).append((shape, parent))
        raw_objects = slide_spec.get("objects", [])
        objects = [obj for obj in raw_objects if isinstance(obj, dict)] if isinstance(raw_objects, list) else []
        declared_ids = {str(obj.get("object_id")) for obj in objects if obj.get("object_id")}
        allowed_names = set(str(value) for value in (object_manifest.get("allowed_shape_names") or []))
        allowed_names.update(str(value) for value in (slide_spec.get("allowed_shape_names") or []))
        ignored_names = set(str(value) for value in (object_manifest.get("ignored_shape_names") or []))
        ignored_names.update(str(value) for value in (slide_spec.get("ignored_shape_names") or []))

        unmanaged: list[dict[str, Any]] = []
        observed_top_level = 0
        for shape, parent in entries:
            if parent is not None:
                continue
            observed_top_level += 1
            if shape.name in declared_ids or shape.name in allowed_names or shape.name in ignored_names:
                continue
            if bool(getattr(shape, "is_placeholder", False)) and not _text_from_shape(shape).strip():
                warnings.append({"code": "template_placeholder_ignored", "slide_no": slide_no, "object_id": shape.name})
                continue
            unmanaged.append({
                "object_id": shape.name,
                "actual_type": _actual_kind(shape),
                "shape_type": str(shape.shape_type),
                "text": _text_from_shape(shape),
            })
            errors.append({"code": "undeclared_visible_shape", "slide_no": slide_no, "object_id": shape.name, "actual_type": _actual_kind(shape)})

        expected_count += len(objects)
        for obj in objects:
            object_id = str(obj.get("object_id", ""))
            matches = by_name.get(object_id, [])
            if not matches:
                _append_error(errors, "semantic_object_missing", slide_no, object_id)
                continue
            if len(matches) != 1:
                _append_error(errors, "semantic_object_name_not_unique", slide_no, object_id, observed=len(matches))
                continue
            shape, parent = matches[0]
            actual = _actual_kind(shape)
            expected_type = str(obj.get("object_type", ""))
            type_ok = _type_matches(expected_type, actual)
            record: dict[str, Any] = {
                "slide_no": slide_no,
                "object_id": object_id,
                "role": obj.get("role"),
                "expected_type": expected_type,
                "actual_type": actual,
                "shape_type": str(shape.shape_type),
                "parent_group": parent,
                "top_level": parent is None,
                "text": _text_from_shape(shape),
                "semantic_checks": {"type": type_ok},
            }
            if not type_ok:
                _append_error(errors, "semantic_type_mismatch", slide_no, object_id, expected=expected_type, actual=actual)

            if expected_type == "editable_text":
                native_textbox = _is_native_textbox(shape)
                record["text_box"] = {"native": native_textbox, "placeholder": bool(getattr(shape, "is_placeholder", False))}
                record["semantic_checks"]["native_textbox"] = native_textbox
                if not native_textbox:
                    _append_error(errors, "editable_text_not_native_textbox", slide_no, object_id, actual=actual)
                expected_text = _expected_text(obj, text_specs)
                if expected_text is None:
                    record["semantic_checks"]["text_exact"] = False
                    _append_error(errors, "text_manifest_evidence_missing", slide_no, object_id)
                else:
                    observed_text = _text_from_shape(shape)
                    text_ok = expected_text == observed_text
                    record["semantic_checks"]["text_exact"] = text_ok
                    record["expected_text"] = expected_text
                    record["observed_text_sha256"] = _json_digest(observed_text)
                    if not text_ok:
                        _append_error(errors, "pptx_text_manifest_mismatch", slide_no, object_id, expected=expected_text, observed=observed_text)
                observed_runs = _text_run_evidence(shape)
                record["text_runs"] = {"count": len(observed_runs), "values": observed_runs}

            if actual == "editable_table":
                table = _table_evidence(shape)
                expected_values = _expected_table_values(obj, object_manifest_path.parent)
                data_match = expected_values is not None and expected_values == table["values"]
                source_evidence = _data_source_evidence(obj, object_manifest_path.parent)
                record["table"] = {"observed": table, "expected_values": expected_values, "expected_values_sha256": _json_digest(expected_values) if expected_values is not None else None, "data_source": source_evidence}
                record["semantic_checks"]["native_table"] = True
                record["semantic_checks"]["native_table_data"] = bool(table["rows"] and table["columns"] and table["rectangular"] and table["nonempty_cells"])
                record["semantic_checks"]["table_data_present"] = bool(table["nonempty_cells"])
                record["semantic_checks"]["table_data_matches_source"] = data_match
                if expected_values is None:
                    _append_error(errors, "table_source_data_missing", slide_no, object_id)
                elif not data_match:
                    _append_error(errors, "table_data_mismatch", slide_no, object_id, expected=expected_values, observed=table["values"])
                if source_evidence["available"] and source_evidence["declared_sha256"] is None:
                    _append_error(errors, "table_data_source_hash_missing", slide_no, object_id)
                elif source_evidence["available"] and not source_evidence["hash_match"]:
                    _append_error(errors, "table_data_source_hash_mismatch", slide_no, object_id, evidence=source_evidence)
                elif source_evidence["explicit_path"] and not source_evidence["available"]:
                    _append_error(errors, "table_data_source_unavailable", slide_no, object_id, evidence=source_evidence)
                if not table["nonempty_cells"]:
                    _append_error(errors, "native_table_empty", slide_no, object_id)

            if actual == "editable_chart":
                chart = _chart_evidence(shape)
                expected_values = _expected_chart_values(obj, object_manifest_path.parent)
                chart_match = expected_values is not None and expected_values == chart["chart_data"]
                workbook_match = expected_values is not None and expected_values == chart["workbook_data"]
                source_evidence = _data_source_evidence(obj, object_manifest_path.parent)
                chart_ok = bool(
                    chart["series"]
                    and chart["nonempty_series"]
                    and chart["has_cached_values"]
                    and chart["has_embedded_workbook"]
                    and chart["chart_workbook_match"]
                    and chart_match
                    and workbook_match
                )
                record["chart"] = {"observed": chart, "expected_data": expected_values, "expected_data_sha256": _json_digest(expected_values) if expected_values is not None else None, "data_source": source_evidence}
                record["semantic_checks"]["native_chart"] = True
                record["semantic_checks"]["native_chart_data"] = chart_ok
                record["semantic_checks"]["chart_data_present"] = bool(chart["chart_data"])
                record["semantic_checks"]["embedded_workbook"] = bool(chart["has_embedded_workbook"])
                record["semantic_checks"]["chart_workbook_consistent"] = bool(chart["chart_workbook_match"])
                record["semantic_checks"]["chart_data_matches_source"] = bool(chart_match and workbook_match)
                if expected_values is None:
                    _append_error(errors, "chart_source_data_missing", slide_no, object_id)
                if source_evidence["available"] and source_evidence["declared_sha256"] is None:
                    _append_error(errors, "chart_data_source_hash_missing", slide_no, object_id)
                elif source_evidence["available"] and not source_evidence["hash_match"]:
                    _append_error(errors, "chart_data_source_hash_mismatch", slide_no, object_id, evidence=source_evidence)
                elif source_evidence["explicit_path"] and not source_evidence["available"]:
                    _append_error(errors, "chart_data_source_unavailable", slide_no, object_id, evidence=source_evidence)
                if not chart_ok:
                    _append_error(errors, "native_chart_data_invalid", slide_no, object_id, evidence=chart)

            if actual == "native_group":
                expected_children = [str(child) for child in obj.get("children", []) if child]
                observed_children = [child.name for child in shape.shapes]
                children_ok = expected_children == observed_children if expected_children else not observed_children
                record["group"] = {"expected_children": expected_children, "observed_children": observed_children}
                record["semantic_checks"]["group_children_declared"] = children_ok
                if not children_ok:
                    _append_error(errors, "group_children_manifest_mismatch", slide_no, object_id, expected=expected_children, observed=observed_children)

            brand = _is_brand_lockup(obj)
            if brand:
                crop = _picture_crop_evidence(shape) if actual == "picture" else {"values": {}, "cropped": False}
                duplicate_text = [
                    other.name
                    for other, other_parent in entries
                    if other is not shape and _actual_kind(other) == "native_text" and _text_from_shape(other).strip() and _overlap_ratio(shape, other) >= 0.25
                ]
                brand_ok = actual == "picture" and parent is None and not _text_from_shape(shape) and not crop["cropped"] and not duplicate_text
                record["brand"] = {"crop": crop, "duplicate_text_shapes": duplicate_text, "whole_asset_contract": obj.get("brand_asset_contract")}
                record["semantic_checks"]["brand_lockup_whole_asset"] = brand_ok
                if not brand_ok:
                    _append_error(errors, "brand_lockup_not_whole_independent_asset", slide_no, object_id, actual=actual, parent_group=parent, crop=crop, duplicate_text_shapes=duplicate_text)

            if actual == "picture":
                part_name, blob = _shape_media(shape)
                record["media_part"] = part_name
                record["embedded_sha256"] = hashlib.sha256(blob).hexdigest() if blob is not None else None
                declared_hashes, source_paths, invalid_hashes = _hash_candidates(obj, asset_records, object_manifest_path.parent)
                source_hashes: set[str] = set()
                missing_sources = []
                for source_path in source_paths:
                    if source_path.is_file():
                        source_hashes.add(_digest(source_path).lower())
                    else:
                        missing_sources.append(str(source_path))
                expected_hashes = declared_hashes | source_hashes
                record["source_paths"] = [str(path) for path in source_paths]
                record["declared_source_sha256"] = sorted(declared_hashes)
                record["source_file_sha256"] = sorted(source_hashes)
                record["expected_source_sha256"] = sorted(expected_hashes)
                source_required = _requires_source_hash(obj, brand)
                record["semantic_checks"]["source_hash_format"] = not invalid_hashes
                if invalid_hashes:
                    _append_error(errors, "source_hash_invalid", slide_no, object_id, values=invalid_hashes)
                declaration_ok = not declared_hashes or source_hashes <= declared_hashes
                record["semantic_checks"]["source_declaration_hash"] = declaration_ok
                if not declaration_ok:
                    _append_error(errors, "source_manifest_hash_mismatch", slide_no, object_id, declared=sorted(declared_hashes), observed=sorted(source_hashes))
                if missing_sources:
                    record["semantic_checks"]["source_files_available"] = False
                    if source_required:
                        _append_error(errors, "source_file_unavailable", slide_no, object_id, paths=missing_sources)
                    else:
                        warnings.append({"code": "source_file_unavailable", "slide_no": slide_no, "object_id": object_id, "paths": missing_sources})
                else:
                    record["semantic_checks"]["source_files_available"] = True
                if blob is None:
                    _append_error(errors, "embedded_asset_unreadable", slide_no, object_id)
                    record["semantic_checks"]["source_hash"] = False
                elif not expected_hashes:
                    record["semantic_checks"]["source_hash"] = False
                    if source_required:
                        _append_error(errors, "source_hash_evidence_missing", slide_no, object_id)
                    else:
                        warnings.append({"code": "source_hash_evidence_missing", "slide_no": slide_no, "object_id": object_id})
                else:
                    hash_ok = record["embedded_sha256"].lower() in expected_hashes and declaration_ok
                    record["semantic_checks"]["source_hash"] = hash_ok
                    if not hash_ok:
                        _append_error(errors, "embedded_asset_hash_mismatch", slide_no, object_id, expected=sorted(expected_hashes), observed=record["embedded_sha256"])

            audited.append(record)

        coverage_record = {
            "slide_no": slide_no,
            "declared_object_count": len(objects),
            "observed_top_level_shape_count": observed_top_level,
            "audited_object_count": sum(1 for item in audited if item["slide_no"] == slide_no),
            "undeclared_shape_count": len(unmanaged),
            "undeclared_shapes": unmanaged,
            "complete": not unmanaged,
        }
        coverage.append(coverage_record)

    text_refs_by_slide: dict[int, set[str]] = {}
    for slide_index, slide_spec in enumerate(object_manifest.get("slides", []), 1):
        if not isinstance(slide_spec, dict):
            continue
        slide_no = int(slide_spec.get("slide_no", slide_index))
        raw_objects = slide_spec.get("objects", [])
        text_refs_by_slide[slide_no] = {
            str(value)
            for obj in raw_objects if isinstance(raw_objects, list) and isinstance(obj, dict) and obj.get("object_type") == "editable_text"
            for value in (obj.get("object_id"), obj.get("text_id")) if value
        }
    orphan_text_specs = []
    for spec in text_manifest_entries:
        text_id = str(spec.get("object_id") or spec.get("text_id") or "")
        if text_id not in text_refs_by_slide.get(int(spec.get("slide_no", 0)), set()):
            orphan_text_specs.append({"slide_no": spec.get("slide_no"), "text_id": text_id})
            errors.append({"code": "text_manifest_orphan", "slide_no": spec.get("slide_no"), "object_id": text_id})

    result = {
        "schema": SCHEMA,
        "valid": not errors,
        "status": "passed" if not errors else "blocked",
        "deck": str(deck_path.resolve()),
        "deck_sha256": _digest(deck_path),
        "object_manifest": str(object_manifest_path.resolve()),
        "object_manifest_sha256": _digest(object_manifest_path),
        "text_manifest": str(text_manifest_path.resolve()) if text_manifest_path else None,
        "text_manifest_sha256": _digest(text_manifest_path) if text_manifest_path and text_manifest_path.is_file() else None,
        "expected_object_count": expected_count,
        "audited_object_count": len(audited),
        "observed_top_level_shape_count": sum(item["observed_top_level_shape_count"] for item in coverage),
        "undeclared_shape_count": sum(item["undeclared_shape_count"] for item in coverage),
        "text_manifest_coverage": {"declared": len(text_manifest_entries), "orphan_count": len(orphan_text_specs), "orphans": orphan_text_specs},
        "coverage": coverage,
        "objects": audited,
        "errors": errors,
        "warnings": warnings,
        "human_visual_review_required": True,
    }
    return result


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("deck")
    parser.add_argument("--object-manifest", required=True)
    parser.add_argument("--text-manifest")
    parser.add_argument("--asset-manifest", action="append", default=[])
    parser.add_argument("--report")
    args = parser.parse_args()
    try:
        result = audit(
            Path(args.deck).resolve(),
            Path(args.object_manifest).resolve(),
            Path(args.text_manifest).resolve() if args.text_manifest else None,
            [Path(path).resolve() for path in args.asset_manifest],
        )
    except Exception as exc:
        result = {"schema": SCHEMA, "valid": False, "status": "invalid", "errors": [{"code": "runtime_error", "message": f"{type(exc).__name__}: {exc}"}]}
    if args.report:
        from atomic_output import atomic_write_json

        atomic_write_json(Path(args.report).resolve(), result)
    print(json.dumps(result, ensure_ascii=False))
    return 0 if result.get("valid") is True else 2


if __name__ == "__main__":
    raise SystemExit(main())
