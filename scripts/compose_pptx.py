#!/usr/bin/env python3
"""Compose an editable .pptx from background images, cut-out icons, and text.

Final step of Part 2 (and the merge step of Part 1). Reads one layout file
describing every slide as stacked layers (back to front):
  background (full-bleed) -> whole-frame PNG -> icon/decoration PNGs (positioned) ->
  text boxes
and writes a real PowerPoint file where text stays editable and the frame/icons
stay as movable pictures. The Image2PPTX workflow defaults to the full-slide
`frame` field. If the user explicitly asks to split the frame into movable
parts, the script also accepts `icons[]` entries with `role:"frame_part"`
generated from `frame_parts/`.

Also handles Part 1 "image deck" output: a slide with only `background` and no
icons/texts becomes a full-bleed image slide.

Usage:
    python3 scripts/compose_pptx.py deck.json out.pptx
    python3 scripts/compose_pptx.py deck.json out.pptx --preview-dir out/preview
    python3 scripts/compose_pptx.py deck.json out.pptx --font-dir project-fonts --embed-fonts

Layout schema: see references/image-to-pptx.md. Quick form:
{
  "slide_width_in": 13.333, "slide_height_in": 7.5,
  "units": "fraction",            // "fraction" (of slide W/H) or "px"
  "ref_width": 2048, "ref_height": 1152,   // source image size; required for source_bbox/size_px QA
  "assets_dir": ".",              // base dir for relative file paths
  "slides": [
    {
      "background": "01-bg.png",
      "frame": "frame.png",
      "icons": [
        {"file": "icons/i_r1c1.png", "x":0.1,"y":0.2,"w":0.08,"h":0.08}
      ],
      "texts": [{"text":"标题","x":0.08,"y":0.05,"w":0.6,"h":0.12,
                 "size_ratio":0.074,"color":"#1A1A1A","bold":true,"align":"left","valign":"top",
                 "font":"Microsoft YaHei"}]
    }
  ]
}
A single-slide dict (without "slides") is also accepted.
"""
from __future__ import annotations

import argparse
import json
import os
import posixpath
import subprocess
import sys
import tempfile
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

EMU_PER_INCH = 914400


def _die(msg: str, code: int = 2):
    print(f"Error: {msg}", file=sys.stderr)
    raise SystemExit(code)


def _load_deck(path: Path):
    data = json.loads(path.read_text(encoding="utf-8"))
    if "slides" not in data:
        # Treat the whole object as a single slide; lift deck-level keys out.
        slide_keys = {"background", "frame", "panels", "shapes", "groups", "tables", "charts", "components", "speaker_notes", "notes", "icons", "texts"}
        slide = {k: data[k] for k in slide_keys if k in data}
        deck = {k: v for k, v in data.items() if k not in slide_keys}
        deck["slides"] = [slide]
        data = deck
    data.setdefault("slide_width_in", 13.333)
    data.setdefault("slide_height_in", 7.5)
    data.setdefault("units", "fraction")
    data.setdefault("assets_dir", str(path.parent))
    return data


def _expand_components(deck: dict):
    """Expand validated component instances into the existing object arrays."""
    if not isinstance(deck.get("slides"), list):
        keys = {"background", "frame", "panels", "shapes", "groups", "tables", "charts", "components", "speaker_notes", "notes", "icons", "texts"}
        slide = {key: deck[key] for key in keys if key in deck}
        deck = dict(deck)
        deck["slides"] = [slide]
    has_instances = any(isinstance(sl, dict) and sl.get("components") for sl in deck.get("slides", []))
    if not has_instances:
        return deck
    source = deck.get("component_library") or deck.get("component_library_path")
    if not source:
        _die("component instances require component_library")
    if isinstance(source, dict):
        library = source
    else:
        path = _resolve(Path(deck["assets_dir"]), str(source))
        if not path.exists():
            _die(f"component library not found: {path}")
        try:
            library = json.loads(path.read_text(encoding="utf-8"))
        except Exception as exc:
            _die(f"component library is not valid JSON: {exc}")
    if library.get("schema") != "ai-ppt-plus/component-library/v1":
        _die("component library schema is invalid")
    definitions = {item.get("component_id"): item for item in library.get("components", []) if isinstance(item, dict)}
    theme = dict(library.get("tokens", {}))
    theme.update(deck.get("theme", {}) if isinstance(deck.get("theme", {}), dict) else {})
    deck["theme"] = theme
    target_arrays = {"text": "texts", "shape": "shapes", "group": "groups", "table": "tables", "chart": "charts", "image": "icons", "vector": "icons"}
    for slide_no, slide in enumerate(deck["slides"], 1):
        layout = slide.get("layout_name", theme.get("layout_name", "Blank"))
        for instance in slide.get("components", []):
            if not isinstance(instance, dict):
                _die(f"slide {slide_no}: component instance must be an object")
            component_id = str(instance.get("component_id", ""))
            definition = definitions.get(component_id)
            if definition is None:
                _die(f"slide {slide_no}: component not found: {component_id}")
            if layout not in definition.get("allowed_layouts", []):
                _die(f"slide {slide_no}: component {component_id} is not allowed on layout {layout}")
            primitive = dict(definition.get("defaults", {}))
            primitive.update(instance.get("object", {}))
            primitive["object_id"] = str(instance.get("object_id") or primitive.get("object_id") or component_id)
            primitive["component_id"] = component_id
            target = target_arrays[definition["type"]]
            slide.setdefault(target, []).append(primitive)
        slide.pop("components", None)
    return deck


def _resolve(assets_dir: Path, file: str) -> Path:
    p = Path(file)
    return p if p.is_absolute() else (assets_dir / p)


def _frac(deck, item, key_xy, axis, ref):
    """Return a 0..1 fraction for x/y/w/h given the deck's unit system."""
    val = item[key_xy]
    if deck["units"] == "px":
        return val / ref
    return val


def _text_size_pt(item, sh_pt: float, ref_h: float, default=None) -> float:
    """Return text size in points.

    Preferred Image2PPTX field is size_ratio/size_pct: source text height as a
    fraction/percent of the source image height. This scales cleanly to any PPT
    slide height. Legacy size_px and absolute size(pt) remain supported.
    """
    if item.get("size") is not None:
        return float(item["size"])
    if item.get("size_ratio") is not None:
        return float(item["size_ratio"]) * sh_pt
    if item.get("size_pct") is not None:
        return float(item["size_pct"]) / 100.0 * sh_pt
    if item.get("size_px") is not None and ref_h:
        return float(item["size_px"]) * sh_pt / ref_h
    return float(default) if default is not None else 18.0


def _set_run_fonts(run, name: str):
    """Set latin + east-asian + complex-script typeface so CJK renders correctly."""
    from pptx.oxml.ns import qn
    run.font.name = name
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rpr.find(qn(tag))
        if el is None:
            el = rpr.makeelement(qn(tag), {})
            rpr.append(el)
        el.set("typeface", name)


def _set_run_alpha(run, opacity: float):
    """Set text color alpha for lightly visible GPT-vision text overlays."""
    from pptx.oxml.ns import qn
    opacity = max(0.0, min(1.0, float(opacity)))
    rpr = run._r.get_or_add_rPr()
    solid = rpr.find(qn("a:solidFill"))
    if solid is None:
        return
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:
        return
    for old in srgb.findall(qn("a:alpha")):
        srgb.remove(old)
    alpha = srgb.makeelement(qn("a:alpha"), {"val": str(int(opacity * 100000))})
    srgb.append(alpha)


def _hex_to_rgb(value: str):
    from pptx.dml.color import RGBColor
    v = value.strip().lstrip("#")
    if len(v) == 3:
        v = "".join(ch * 2 for ch in v)
    if len(v) != 6:
        return RGBColor(0x11, 0x11, 0x11)
    return RGBColor(int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def _hex_to_tuple(value: str):
    v = value.strip().lstrip("#")
    if len(v) == 3:
        v = "".join(ch * 2 for ch in v)
    if len(v) != 6:
        return (17, 17, 17)
    return (int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def _hex_to_rgba(value: str, opacity: float = 1.0):
    r, g, b = _hex_to_tuple(value)
    return (r, g, b, int(max(0.0, min(1.0, float(opacity))) * 255))


def _set_fill_alpha(shape, opacity: float):
    """Add an <a:alpha> to a shape's solid fill so cards can be translucent."""
    from pptx.oxml.ns import qn
    spPr = shape._element.spPr
    solid = spPr.find(qn("a:solidFill"))
    if solid is None:
        return
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:
        return
    for old in srgb.findall(qn("a:alpha")):
        srgb.remove(old)
    alpha = srgb.makeelement(qn("a:alpha"), {"val": str(int(max(0.0, min(1.0, opacity)) * 100000))})
    srgb.append(alpha)


def _add_outer_shadow(shape, blur_pt=6.0, dist_pt=3.0, alpha=0.35):
    """Best-effort soft drop shadow for card-like shapes."""
    from pptx.oxml.ns import qn
    from pptx.util import Pt
    spPr = shape._element.spPr
    for old in spPr.findall(qn("a:effectLst")):
        spPr.remove(old)
    eff = spPr.makeelement(qn("a:effectLst"), {})
    sh = eff.makeelement(qn("a:outerShdw"), {
        "blurRad": str(int(Pt(blur_pt))), "dist": str(int(Pt(dist_pt))),
        "dir": "5400000", "rotWithShape": "0"})
    clr = sh.makeelement(qn("a:srgbClr"), {"val": "000000"})
    a = clr.makeelement(qn("a:alpha"), {"val": str(int(alpha * 100000))})
    clr.append(a)
    sh.append(clr)
    eff.append(sh)
    spPr.append(eff)


def _set_gradient_fill(shape, gradient: dict):
    """Write a deterministic PresentationML linear gradient fill."""
    from pptx.oxml.ns import qn
    sp_pr = shape._element.spPr
    for tag in ("a:solidFill", "a:noFill", "a:gradFill"):
        old = sp_pr.find(qn(tag))
        if old is not None:
            sp_pr.remove(old)
    stops = gradient.get("stops", []) if isinstance(gradient, dict) else []
    if len(stops) < 2:
        _die("gradient fill requires at least two color stops")
    grad = sp_pr.makeelement(qn("a:gradFill"), {"rotWithShape": "1"})
    gs_lst = grad.makeelement(qn("a:gsLst"), {})
    for stop in stops:
        if not isinstance(stop, dict) or not isinstance(stop.get("color"), str):
            _die("gradient stops require color fields")
        raw_pos = float(stop.get("position", stop.get("pos", 0)))
        pos = max(0, min(100000, int(raw_pos * (100000 if 0 <= raw_pos <= 1 else 1000))))
        gs = gs_lst.makeelement(qn("a:gs"), {"pos": str(pos)})
        color = gs.makeelement(qn("a:srgbClr"), {"val": stop["color"].lstrip("#")[:6]})
        if stop.get("opacity") is not None:
            color.append(color.makeelement(qn("a:alpha"), {"val": str(int(max(0, min(1, float(stop["opacity"]))) * 100000))}))
        gs.append(color)
        gs_lst.append(gs)
    grad.append(gs_lst)
    angle = float(gradient.get("angle", 0))
    grad.append(grad.makeelement(qn("a:lin"), {"ang": str(int(angle * 60000) % 21600000), "scaled": "1"}))
    line = sp_pr.find(qn("a:ln"))
    if line is not None:
        sp_pr.insert(list(sp_pr).index(line), grad)
    else:
        sp_pr.append(grad)


def _apply_shape_fill(shape, spec: dict):
    gradient = spec.get("gradient") or spec.get("fill_gradient")
    if gradient:
        _set_gradient_fill(shape, gradient)
    elif spec.get("fill"):
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(spec["fill"])
        if spec.get("opacity") is not None:
            _set_fill_alpha(shape, float(spec["opacity"]))
    else:
        shape.fill.background()


def _set_alt_text(shape, text: str | None):
    if not text:
        return
    element = shape._element
    for attr in ("nvSpPr", "nvPicPr", "nvGrpSpPr", "nvGraphicFramePr"):
        container = getattr(element, attr, None)
        if container is not None and getattr(container, "cNvPr", None) is not None:
            container.cNvPr.set("descr", str(text))
            return


def _replace_svg_media(pptx_path: Path, svg_assets: list[tuple[int, str, Path]]):
    """Replace temporary PNG media with native SVG package parts."""
    if not svg_assets:
        return
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    with zipfile.ZipFile(pptx_path, "r") as zin:
        entries = {name: zin.read(name) for name in zin.namelist()}
    replacements = []
    for slide_no, object_name, svg_path in svg_assets:
        slide_name = f"ppt/slides/slide{slide_no}.xml"
        rels_name = f"ppt/slides/_rels/slide{slide_no}.xml.rels"
        root = ET.fromstring(entries[slide_name])
        rel_root = ET.fromstring(entries[rels_name])
        target = None
        for pic in root.findall(f".//{{{p_ns}}}pic"):
            node = pic.find(f".//{{{p_ns}}}cNvPr")
            if node is None or node.get("name") != object_name:
                continue
            blip = pic.find(f".//{{{a_ns}}}blip")
            rid = blip.get(f"{{{r_ns}}}embed") if blip is not None else None
            for rel in rel_root:
                if rel.get("Id") == rid:
                    target = posixpath.normpath(posixpath.join("ppt/slides", rel.get("Target")))
                    rel.set("Target", "../media/" + Path(target).with_suffix(".svg").name)
                    break
            break
        if target is None or target not in entries:
            _die(f"could not resolve SVG media relationship for {object_name}")
        new_target = str(Path(target).with_suffix(".svg")).replace("\\", "/")
        entries[rels_name] = ET.tostring(rel_root, encoding="utf-8", xml_declaration=True)
        entries[new_target] = svg_path.read_bytes()
        del entries[target]
        replacements.append((target, new_target))
    content = ET.fromstring(entries["[Content_Types].xml"])
    for old, new in replacements:
        for override in content:
            if override.get("PartName") == "/" + old:
                override.set("PartName", "/" + new)
                override.set("ContentType", "image/svg+xml")
    entries["[Content_Types].xml"] = ET.tostring(content, encoding="utf-8", xml_declaration=True)
    tmp = pptx_path.with_suffix(".svg-rewrite.tmp")
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in entries.items():
            zout.writestr(name, data)
    tmp.replace(pptx_path)


def _svg_to_png(svg_path: Path) -> Path:
    fd, temp_name = tempfile.mkstemp(suffix=".png")
    os.close(fd)
    output = Path(temp_name)
    try:
        from cairosvg import svg2png
        svg2png(url=str(svg_path), write_to=str(output))
    except ImportError:
        for command in (("inkscape", str(svg_path), "--export-filename", str(output)), ("convert", str(svg_path), str(output))):
            try:
                result = subprocess.run(command, capture_output=True, text=True, timeout=30, check=False)
            except (FileNotFoundError, subprocess.TimeoutExpired):
                continue
            if result.returncode == 0 and output.is_file():
                return output
        output.unlink(missing_ok=True)
        _die(f"SVG asset requires cairosvg, inkscape, or ImageMagick: {svg_path}")
    return output


def _choose_slide_layout(prs, slide_spec: dict, theme: dict):
    """Select an existing template layout without inventing a new master."""
    requested = slide_spec.get("layout_name", theme.get("layout_name"))
    if requested:
        for layout in prs.slide_layouts:
            if layout.name == str(requested):
                return layout
        _die(f"slide layout not found: {requested}")
    if slide_spec.get("layout_index") is not None:
        index = int(slide_spec["layout_index"])
        if index < 0 or index >= len(prs.slide_layouts):
            _die(f"slide layout index out of range: {index}")
        return prs.slide_layouts[index]
    return prs.slide_layouts[6]


def build_pptx(deck, out_path: Path):
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    shape_map = {"rect": MSO_SHAPE.RECTANGLE, "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
                 "oval": MSO_SHAPE.OVAL, "ellipse": MSO_SHAPE.OVAL,
                 "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE, "chevron": MSO_SHAPE.CHEVRON,
                 "right_arrow": MSO_SHAPE.RIGHT_ARROW, "left_arrow": MSO_SHAPE.LEFT_ARROW,
                 "up_arrow": MSO_SHAPE.UP_ARROW, "down_arrow": MSO_SHAPE.DOWN_ARROW,
                 "pentagon": MSO_SHAPE.REGULAR_PENTAGON, "hexagon": MSO_SHAPE.HEXAGON,
                 "parallelogram": MSO_SHAPE.PARALLELOGRAM, "trapezoid": MSO_SHAPE.TRAPEZOID,
                 "diamond": MSO_SHAPE.DIAMOND}
    chart_map = {"bar": XL_CHART_TYPE.BAR_CLUSTERED, "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                 "line": XL_CHART_TYPE.LINE_MARKERS, "pie": XL_CHART_TYPE.PIE,
                 "doughnut": XL_CHART_TYPE.DOUGHNUT}

    assets_dir = Path(deck["assets_dir"])
    sw_in = float(deck["slide_width_in"])
    sh_in = float(deck["slide_height_in"])
    sw_emu = int(round(sw_in * EMU_PER_INCH))
    sh_emu = int(round(sh_in * EMU_PER_INCH))
    ref_w = float(deck.get("ref_width") or 0)
    ref_h = float(deck.get("ref_height") or 0)
    sh_pt = sh_in * 72.0

    prs = Presentation()
    prs.slide_width = Emu(sw_emu)
    prs.slide_height = Emu(sh_emu)
    blank = prs.slide_layouts[6]

    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                 "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
    anchor_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                  "center": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}

    svg_assets = []
    theme = deck.get("theme", {}) if isinstance(deck.get("theme", {}), dict) else {}
    for idx, sl in enumerate(deck["slides"], 1):
        slide = prs.slides.add_slide(_choose_slide_layout(prs, sl, theme))

        bg = sl.get("background")
        if bg:
            bg_path = _resolve(assets_dir, bg)
            if not bg_path.exists():
                _die(f"slide {idx}: background not found: {bg_path}")
            picture = slide.shapes.add_picture(str(bg_path), 0, 0, width=Emu(sw_emu), height=Emu(sh_emu))
            picture.name = str(sl.get("background_object_id", "background"))
            _set_alt_text(picture, sl.get("background_alt_text"))

        frame = sl.get("frame")
        if frame:
            fr_path = _resolve(assets_dir, frame)
            if not fr_path.exists():
                _die(f"slide {idx}: frame not found: {fr_path}")
            picture = slide.shapes.add_picture(str(fr_path), 0, 0, width=Emu(sw_emu), height=Emu(sh_emu))
            picture.name = str(sl.get("frame_object_id", "frame"))
            _set_alt_text(picture, sl.get("frame_alt_text"))

        for shape_index, shp in enumerate(sl.get("shapes", []), 1):
            fx = _frac(deck, shp, "x", "x", ref_w)
            fy = _frac(deck, shp, "y", "y", ref_h)
            fw = _frac(deck, shp, "w", "w", ref_w)
            fh = _frac(deck, shp, "h", "h", ref_h)
            left, top = Emu(int(fx * sw_emu)), Emu(int(fy * sh_emu))
            wid, hei = Emu(int(fw * sw_emu)), Emu(int(fh * sh_emu))
            stype = shp.get("type", "rounded_rect")

            if stype == "line":
                conn = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT, left, top, Emu(int((fx + fw) * sw_emu)),
                    Emu(int((fy + fh) * sh_emu)))
                conn.line.color.rgb = _hex_to_rgb(shp.get("line", shp.get("fill", "#FFFFFF")))
                conn.line.width = Pt(float(shp.get("line_width", 1.5)))
                conn.name = str(shp.get("object_id") or shp.get("name") or f"shape-{shape_index:02d}")
                continue

            shape = slide.shapes.add_shape(shape_map.get(stype, MSO_SHAPE.ROUNDED_RECTANGLE),
                                           left, top, wid, hei)
            shape.name = str(shp.get("object_id") or shp.get("name") or f"shape-{shape_index:02d}")
            if stype == "rounded_rect" and "radius" in shp:
                try:
                    shape.adjustments[0] = float(shp["radius"])
                except Exception:
                    pass

            _apply_shape_fill(shape, shp)

            if shp.get("line"):
                shape.line.color.rgb = _hex_to_rgb(shp["line"])
                shape.line.width = Pt(float(shp.get("line_width", 1.0)))
            else:
                shape.line.fill.background()

            if shp.get("shadow"):
                _add_outer_shadow(shape)
            else:
                shape.shadow.inherit = False

            if shp.get("rotation"):
                shape.rotation = float(shp["rotation"])
            _set_alt_text(shape, shp.get("alt_text"))

        # Groups preserve a semantic component boundary while retaining
        # independently editable child shapes. Children use slide coordinates
        # by default; set children_coordinate_space=local to position them
        # inside the group's x/y/w/h box using fractional local coordinates.
        for group_index, group_spec in enumerate(sl.get("groups", []), 1):
            group = slide.shapes.add_group_shape()
            group.name = str(group_spec.get("object_id") or group_spec.get("name") or f"group-{group_index:02d}")
            local = group_spec.get("children_coordinate_space") == "local"
            gx = _frac(deck, group_spec, "x", "x", ref_w) if all(k in group_spec for k in ("x", "y", "w", "h")) else 0
            gy = _frac(deck, group_spec, "y", "y", ref_h) if all(k in group_spec for k in ("x", "y", "w", "h")) else 0
            gw = _frac(deck, group_spec, "w", "w", ref_w) if all(k in group_spec for k in ("x", "y", "w", "h")) else 1
            gh = _frac(deck, group_spec, "h", "h", ref_h) if all(k in group_spec for k in ("x", "y", "w", "h")) else 1
            for child_index, child in enumerate(group_spec.get("children", []), 1):
                if not isinstance(child, dict):
                    continue
                child = dict(child)
                if local:
                    child["x"] = gx + _frac(deck, child, "x", "x", ref_w) * gw
                    child["y"] = gy + _frac(deck, child, "y", "y", ref_h) * gh
                    child["w"] = _frac(deck, child, "w", "w", ref_w) * gw
                    child["h"] = _frac(deck, child, "h", "h", ref_h) * gh
                    child_deck = dict(deck, units="fraction")
                else:
                    child_deck = deck
                cfx = _frac(child_deck, child, "x", "x", ref_w)
                cfy = _frac(child_deck, child, "y", "y", ref_h)
                cfw = _frac(child_deck, child, "w", "w", ref_w)
                cfh = _frac(child_deck, child, "h", "h", ref_h)
                if child.get("type", "rounded_rect") == "line":
                    conn = group.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(int(cfx * sw_emu)), Emu(int(cfy * sh_emu)), Emu(int((cfx + cfw) * sw_emu)), Emu(int((cfy + cfh) * sh_emu)))
                    conn.line.color.rgb = _hex_to_rgb(child.get("line", child.get("fill", "#FFFFFF")))
                    conn.line.width = Pt(float(child.get("line_width", 1.5)))
                    conn.name = str(child.get("object_id") or child.get("name") or f"{group.name}-child-{child_index:02d}")
                    continue
                child_shape = group.shapes.add_shape(shape_map.get(child.get("type", "rounded_rect"), MSO_SHAPE.ROUNDED_RECTANGLE), Emu(int(cfx * sw_emu)), Emu(int(cfy * sh_emu)), Emu(int(cfw * sw_emu)), Emu(int(cfh * sh_emu)))
                child_shape.name = str(child.get("object_id") or child.get("name") or f"{group.name}-child-{child_index:02d}")
                _apply_shape_fill(child_shape, child)
                if child.get("line"):
                    child_shape.line.color.rgb = _hex_to_rgb(child["line"])
                    child_shape.line.width = Pt(float(child.get("line_width", 1.0)))
                else:
                    child_shape.line.fill.background()
                if child.get("rotation"):
                    child_shape.rotation = float(child["rotation"])
                _set_alt_text(child_shape, child.get("alt_text"))
            group.shapes._recalculate_extents()
            _set_alt_text(group, group_spec.get("alt_text"))

        for table_index, table_spec in enumerate(sl.get("tables", []), 1):
            rows = table_spec.get("rows", [])
            columns = int(table_spec.get("columns") or (len(rows[0]) if rows and isinstance(rows[0], list) else 0))
            if not rows or columns <= 0:
                _die(f"slide {idx}: table requires non-empty rows and columns")
            table = slide.shapes.add_table(len(rows), columns,
                Emu(int(_frac(deck, table_spec, "x", "x", ref_w) * sw_emu)),
                Emu(int(_frac(deck, table_spec, "y", "y", ref_h) * sh_emu)),
                Emu(int(_frac(deck, table_spec, "w", "w", ref_w) * sw_emu)),
                Emu(int(_frac(deck, table_spec, "h", "h", ref_h) * sh_emu))).table
            frame = table._graphic_frame
            frame.name = str(table_spec.get("object_id") or table_spec.get("name") or f"table-{table_index:02d}")
            font = str(table_spec.get("font") or theme.get("font") or "Microsoft YaHei")
            header_fill = table_spec.get("header_fill") or theme.get("table_header_fill")
            body_fill = table_spec.get("fill") or theme.get("table_fill")
            widths = table_spec.get("column_widths") or []
            for ci, width in enumerate(widths[:columns]):
                table.columns[ci].width = Emu(int(float(width) * sw_emu)) if deck["units"] == "fraction" else int(width)
            for ri, row in enumerate(rows):
                for ci in range(columns):
                    cell = table.cell(ri, ci)
                    cell.text = str(row[ci]) if ci < len(row) else ""
                    fill = header_fill if ri == 0 else body_fill
                    if fill:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = _hex_to_rgb(fill)
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.name = font
                            run.font.size = Pt(float(table_spec.get("size", theme.get("size", 12))))
                            if table_spec.get("color") or theme.get("text_color"):
                                run.font.color.rgb = _hex_to_rgb(table_spec.get("color") or theme.get("text_color"))
                            run.font.bold = bool(ri == 0 and table_spec.get("header_bold", True))
            for merge in table_spec.get("merges", []):
                if isinstance(merge, list) and len(merge) == 4:
                    r1, c1, r2, c2 = [int(value) for value in merge]
                    table.cell(r1, c1).merge(table.cell(r2, c2))
            _set_alt_text(frame, table_spec.get("alt_text"))

        for chart_index, chart_spec in enumerate(sl.get("charts", []), 1):
            chart_type = str(chart_spec.get("type", "column")).casefold()
            if chart_type not in chart_map:
                _die(f"slide {idx}: unsupported chart type: {chart_type}")
            data = CategoryChartData()
            data.categories = [str(value) for value in chart_spec.get("categories", [])]
            series = chart_spec.get("series", [])
            if not data.categories or not series:
                _die(f"slide {idx}: chart requires categories and series")
            for item in series:
                data.add_series(str(item.get("name", "Series")), [float(value) for value in item.get("values", [])])
            frame = slide.shapes.add_chart(chart_map[chart_type],
                Emu(int(_frac(deck, chart_spec, "x", "x", ref_w) * sw_emu)),
                Emu(int(_frac(deck, chart_spec, "y", "y", ref_h) * sh_emu)),
                Emu(int(_frac(deck, chart_spec, "w", "w", ref_w) * sw_emu)),
                Emu(int(_frac(deck, chart_spec, "h", "h", ref_h) * sh_emu)), data)
            frame.name = str(chart_spec.get("object_id") or chart_spec.get("name") or f"chart-{chart_index:02d}")
            chart = frame.chart
            chart.has_title = bool(chart_spec.get("title"))
            if chart.has_title:
                chart.chart_title.text_frame.text = str(chart_spec["title"])
            chart.has_legend = bool(chart_spec.get("legend", len(series) > 1))
            if chart.has_legend:
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
            if chart_spec.get("data_labels"):
                for series_item in chart.series:
                    series_item.has_data_labels = True
                    series_item.data_labels.show_value = True
            palette = chart_spec.get("colors") or theme.get("chart_colors") or []
            for series_item, color in zip(chart.series, palette):
                series_item.format.fill.solid()
                series_item.format.fill.fore_color.rgb = _hex_to_rgb(color)
            _set_alt_text(frame, chart_spec.get("alt_text"))

        notes = sl.get("speaker_notes") or sl.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)

        # Semantic panels are independent movable assets. They intentionally
        # render before icons and text, and are not folded into the full-slide
        # frame layer.
        for panel in sl.get("panels", []):
            pp = _resolve(assets_dir, panel["file"])
            if not pp.exists():
                _die(f"slide {idx}: panel not found: {pp}")
            fx = _frac(deck, panel, "x", "x", ref_w)
            fy = _frac(deck, panel, "y", "y", ref_h)
            fw = _frac(deck, panel, "w", "w", ref_w)
            fh = _frac(deck, panel, "h", "h", ref_h)
            picture = slide.shapes.add_picture(
                str(pp), Emu(int(fx * sw_emu)), Emu(int(fy * sh_emu)),
                width=Emu(int(fw * sw_emu)), height=Emu(int(fh * sh_emu)))
            picture.name = str(panel.get("object_id") or panel.get("panel_id") or f"panel-{idx}")
            _set_alt_text(picture, panel.get("alt_text"))

        for icon_index, ic in enumerate(sl.get("icons", []), 1):
            ip = _resolve(assets_dir, ic["file"])
            if not ip.exists():
                _die(f"slide {idx}: icon not found: {ip}")
            fx = _frac(deck, ic, "x", "x", ref_w)
            fy = _frac(deck, ic, "y", "y", ref_h)
            fw = _frac(deck, ic, "w", "w", ref_w)
            fh = _frac(deck, ic, "h", "h", ref_h)
            source_path = ip
            if ip.suffix.casefold() == ".svg":
                source_path = _svg_to_png(ip)
            picture = slide.shapes.add_picture(
                str(source_path), Emu(int(fx * sw_emu)), Emu(int(fy * sh_emu)),
                width=Emu(int(fw * sw_emu)), height=Emu(int(fh * sh_emu)))
            picture.name = str(ic.get("name") or ic.get("object_id") or f"icon-{icon_index:02d}")
            _set_alt_text(picture, ic.get("alt_text"))
            if ip.suffix.casefold() == ".svg":
                svg_assets.append((idx, picture.name, ip))

        for tx in sl.get("texts", []):
            fx = _frac(deck, tx, "x", "x", ref_w)
            fy = _frac(deck, tx, "y", "y", ref_h)
            fw = _frac(deck, tx, "w", "w", ref_w)
            fh = _frac(deck, tx, "h", "h", ref_h)
            box = slide.shapes.add_textbox(
                Emu(int(fx * sw_emu)), Emu(int(fy * sh_emu)),
                Emu(int(fw * sw_emu)), Emu(int(fh * sh_emu)))
            tf = box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = anchor_map.get(tx.get("valign", "top"), MSO_ANCHOR.TOP)
            # PowerPoint text boxes have non-zero default internal margins. The
            # layout bbox already represents the visual text box measured on the
            # source image, so default margins create systematic placement drift.
            tf.margin_left = Pt(float(tx.get("margin_left", 0)))
            tf.margin_right = Pt(float(tx.get("margin_right", 0)))
            tf.margin_top = Pt(float(tx.get("margin_top", 0)))
            tf.margin_bottom = Pt(float(tx.get("margin_bottom", 0)))

            size_pt = _text_size_pt(tx, sh_pt, ref_h)

            color = _hex_to_rgb(tx.get("color", theme.get("text_color", "#111111")))
            font = tx.get("font", theme.get("font", "Microsoft YaHei"))
            bold = bool(tx.get("bold", False))
            italic = bool(tx.get("italic", False))
            align = align_map.get(tx.get("align", "left"), PP_ALIGN.LEFT)
            line_spacing = tx.get("line_spacing")
            opacity = float(tx.get("opacity", 1.0))
            if tx.get("name") or tx.get("object_id"):
                box.name = str(tx.get("name") or tx["object_id"])

            runs = tx.get("runs")
            if runs:
                para = tf.paragraphs[0]
                para.alignment = align
                if line_spacing:
                    para.line_spacing = float(line_spacing)
                for rinfo in runs:
                    run = para.add_run()
                    run.text = str(rinfo.get("text", ""))
                    run.font.size = Pt(_text_size_pt(rinfo, sh_pt, ref_h, default=size_pt))
                    run.font.bold = bool(rinfo.get("bold", bold))
                    run.font.italic = bool(rinfo.get("italic", italic))
                    run.font.color.rgb = _hex_to_rgb(rinfo["color"]) if rinfo.get("color") else color
                    _set_run_fonts(run, str(rinfo.get("font", font)))
                    _set_run_alpha(run, float(rinfo.get("opacity", opacity)))
            else:
                lines = str(tx.get("text", "")).split("\n")
                for li, line in enumerate(lines):
                    para = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                    para.alignment = align
                    if line_spacing:
                        para.line_spacing = float(line_spacing)
                    run = para.add_run()
                    run.text = line
                    run.font.size = Pt(size_pt)
                    run.font.bold = bold
                    run.font.italic = italic
                    run.font.color.rgb = color
                    _set_run_fonts(run, font)
                    _set_run_alpha(run, opacity)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    _replace_svg_media(out_path, svg_assets)
    print(f"Wrote {out_path}  ({len(deck['slides'])} slides)")


# --------------------------- optional PNG preview ---------------------------

def _find_cjk_font(font_dir=None, bold=False):
    if font_dir:
        local = sorted(Path(font_dir).glob('*.ttf')) + sorted(Path(font_dir).glob('*.otf')) + sorted(Path(font_dir).glob('*.ttc'))
        if local:
            return str(local[0])
    candidates_bold = [
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/STHeiti Medium.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Bold.ttc",
    ]
    candidates = [
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/STHeiti Light.ttc",
        "/Library/Fonts/Arial Unicode.ttf",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    pool = candidates_bold + candidates if bold else candidates
    for c in pool:
        if Path(c).exists():
            return c
    return None


def _poly_points(stype, x0, y0, x1, y1):
    """Approximate polygon vertices for preview rendering of non-rect shapes."""
    w, h = x1 - x0, y1 - y0
    if stype == "triangle":
        return [(x0 + w / 2, y0), (x1, y1), (x0, y1)]
    if stype == "diamond":
        return [(x0 + w / 2, y0), (x1, y0 + h / 2), (x0 + w / 2, y1), (x0, y0 + h / 2)]
    if stype == "right_arrow":
        m = h * 0.30
        return [(x0, y0 + m), (x1 - w * 0.4, y0 + m), (x1 - w * 0.4, y0), (x1, y0 + h / 2),
                (x1 - w * 0.4, y1), (x1 - w * 0.4, y1 - m), (x0, y1 - m)]
    if stype == "left_arrow":
        m = h * 0.30
        return [(x1, y0 + m), (x0 + w * 0.4, y0 + m), (x0 + w * 0.4, y0), (x0, y0 + h / 2),
                (x0 + w * 0.4, y1), (x0 + w * 0.4, y1 - m), (x1, y1 - m)]
    if stype == "up_arrow":
        m = w * 0.30
        return [(x0 + m, y1), (x0 + m, y0 + h * 0.4), (x0, y0 + h * 0.4), (x0 + w / 2, y0),
                (x1, y0 + h * 0.4), (x1 - m, y0 + h * 0.4), (x1 - m, y1)]
    if stype == "chevron":
        notch = w * 0.25
        return [(x0, y0), (x1 - notch, y0), (x1, y0 + h / 2), (x1 - notch, y1),
                (x0, y1), (x0 + notch, y0 + h / 2)]
    if stype == "trapezoid":
        inset = w * 0.22
        return [(x0 + inset, y0), (x1 - inset, y0), (x1, y1), (x0, y1)]
    if stype == "parallelogram":
        sk = w * 0.22
        return [(x0 + sk, y0), (x1, y0), (x1 - sk, y1), (x0, y1)]
    if stype == "pentagon":
        import math
        cx, cy, rx, ry = x0 + w / 2, y0 + h / 2, w / 2, h / 2
        return [(cx + rx * math.sin(2 * math.pi * i / 5),
                 cy - ry * math.cos(2 * math.pi * i / 5)) for i in range(5)]
    if stype == "hexagon":
        return [(x0 + w * 0.25, y0), (x0 + w * 0.75, y0), (x1, y0 + h / 2),
                (x0 + w * 0.75, y1), (x0 + w * 0.25, y1), (x0, y0 + h / 2)]
    return None


def _wrap_text(draw, text, font, max_w):
    out = []
    for raw in text.split("\n"):
        if not raw:
            out.append("")
            continue
        line = ""
        for ch in raw:
            trial = line + ch
            w = draw.textlength(trial, font=font)
            if w > max_w and line:
                out.append(line)
                line = ch
            else:
                line = trial
        out.append(line)
    return out


def render_previews(deck, preview_dir: Path):
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        print("Preview skipped: Pillow not available.", file=sys.stderr)
        return

    assets_dir = Path(deck["assets_dir"])
    sw_in = float(deck["slide_width_in"])
    sh_in = float(deck["slide_height_in"])
    ratio = sw_in / sh_in
    if deck["units"] == "px" and deck.get("ref_width") and deck.get("ref_height"):
        CW, CH = int(deck["ref_width"]), int(deck["ref_height"])
    else:
        CW = 1600
        CH = int(round(CW / ratio))
    sh_pt = sh_in * 72.0
    ref_w = float(deck.get("ref_width") or CW)
    ref_h = float(deck.get("ref_height") or CH)

    preview_dir.mkdir(parents=True, exist_ok=True)
    regular_path = _find_cjk_font(deck.get("font_dir"), bold=False)
    bold_path = _find_cjk_font(deck.get("font_dir"), bold=True)

    for idx, sl in enumerate(deck["slides"], 1):
        canvas = Image.new("RGBA", (CW, CH), (255, 255, 255, 255))
        bg = sl.get("background")
        if bg:
            bp = _resolve(assets_dir, bg)
            if bp.exists():
                with Image.open(bp) as im:
                    canvas.paste(im.convert("RGBA").resize((CW, CH)), (0, 0))

        frame = sl.get("frame")
        if frame:
            fp = _resolve(assets_dir, frame)
            if fp.exists():
                with Image.open(fp) as im:
                    canvas.alpha_composite(im.convert("RGBA").resize((CW, CH)), (0, 0))

        for shp in sl.get("shapes", []):
            fx = _frac(deck, shp, "x", "x", ref_w)
            fy = _frac(deck, shp, "y", "y", ref_h)
            fw = _frac(deck, shp, "w", "w", ref_w)
            fh = _frac(deck, shp, "h", "h", ref_h)
            x0, y0 = int(fx * CW), int(fy * CH)
            x1, y1 = int((fx + fw) * CW), int((fy + fh) * CH)
            overlay = Image.new("RGBA", (CW, CH), (0, 0, 0, 0))
            od = ImageDraw.Draw(overlay)
            stype = shp.get("type", "rounded_rect")
            a = int(float(shp.get("opacity", 1.0)) * 255) if shp.get("fill") else 0
            fill = _hex_to_tuple(shp["fill"]) + (a,) if shp.get("fill") else None
            line = _hex_to_tuple(shp["line"]) + (255,) if shp.get("line") else None
            lw = int(round(float(shp.get("line_width", 1.0)) * CH / (sh_in * 72.0)))
            lw = max(1, lw) if line else 0
            poly = _poly_points(stype, x0, y0, x1, y1)
            if stype == "line":
                od.line([x0, y0, x1, y1], fill=line or (255, 255, 255, 255),
                        width=max(1, lw or 2))
            elif stype in ("oval", "ellipse"):
                od.ellipse([x0, y0, x1, y1], fill=fill, outline=line, width=lw)
            elif stype == "rounded_rect":
                rad = int(float(shp.get("radius", 0.12)) * min(x1 - x0, y1 - y0))
                od.rounded_rectangle([x0, y0, x1, y1], radius=max(1, rad), fill=fill,
                                     outline=line, width=lw)
            elif poly is not None:
                od.polygon(poly, fill=fill, outline=line, width=lw)
            else:
                od.rectangle([x0, y0, x1, y1], fill=fill, outline=line, width=lw)
            canvas.alpha_composite(overlay)

        for panel in sl.get("panels", []):
            pp = _resolve(assets_dir, panel["file"])
            if not pp.exists():
                continue
            fx = _frac(deck, panel, "x", "x", ref_w)
            fy = _frac(deck, panel, "y", "y", ref_h)
            fw = _frac(deck, panel, "w", "w", ref_w)
            fh = _frac(deck, panel, "h", "h", ref_h)
            with Image.open(pp) as im:
                panel_im = im.convert("RGBA").resize((max(1, int(fw * CW)), max(1, int(fh * CH))))
            canvas.alpha_composite(panel_im, (int(fx * CW), int(fy * CH)))

        for ic in sl.get("icons", []):
            ip = _resolve(assets_dir, ic["file"])
            if not ip.exists():
                continue
            fx = _frac(deck, ic, "x", "x", ref_w)
            fy = _frac(deck, ic, "y", "y", ref_h)
            fw = _frac(deck, ic, "w", "w", ref_w)
            fh = _frac(deck, ic, "h", "h", ref_h)
            tw, th = max(1, int(fw * CW)), max(1, int(fh * CH))
            with Image.open(ip) as im:
                icon = im.convert("RGBA").resize((tw, th))
            canvas.alpha_composite(icon, (int(fx * CW), int(fy * CH)))

        draw = ImageDraw.Draw(canvas)
        for tx in sl.get("texts", []):
            fx = _frac(deck, tx, "x", "x", ref_w)
            fy = _frac(deck, tx, "y", "y", ref_h)
            fw = _frac(deck, tx, "w", "w", ref_w)
            size_pt = _text_size_pt(tx, sh_pt, ref_h)
            px = max(8, int(round(size_pt * CH / sh_pt)))
            bold = bool(tx.get("bold", False))
            fpath = (bold_path if bold else regular_path) or regular_path
            try:
                font = ImageFont.truetype(fpath, px) if fpath else ImageFont.load_default()
            except Exception:
                font = ImageFont.load_default()
            color = tx.get("color", "#111111")
            align = tx.get("align", "left")
            valign = tx.get("valign", "top")
            opacity = float(tx.get("opacity", 1.0))
            bx0, by0 = int(fx * CW), int(fy * CH)
            bw, bh = int(fw * CW), int(_frac(deck, tx, "h", "h", ref_h) * CH)
            line_h = int(px * 1.3)
            runs = tx.get("runs")
            if runs:
                # Lay out rich runs character-by-character so wrapping follows
                # the same visual box as plain text. The previous preview path
                # drew each run as one unwrapped string, which made a valid
                # PPTX preview look broken whenever runs crossed a line break.
                chars = []
                for rinfo in runs:
                    rb = bool(rinfo.get("bold", bold))
                    rpx = max(8, int(round(_text_size_pt(rinfo, sh_pt, ref_h, default=size_pt) * CH / sh_pt)))
                    rfp = (bold_path if rb else regular_path) or regular_path
                    try:
                        rfont = ImageFont.truetype(rfp, rpx) if rfp else ImageFont.load_default()
                    except Exception:
                        rfont = ImageFont.load_default()
                    for ch in str(rinfo.get("text", "")):
                        chars.append((ch, rfont, rinfo.get("color", color),
                                      float(rinfo.get("opacity", opacity)), rb))
                lines, line = [], []
                line_width = 0.0
                for ch, cfont, ccolor, copacity, cbold in chars:
                    if ch == "\n":
                        lines.append(line); line = []; line_width = 0.0
                        continue
                    cw = draw.textlength(ch, font=cfont)
                    if line and line_width + cw > bw:
                        lines.append(line); line = []; line_width = 0.0
                    line.append((ch, cfont, ccolor, copacity, cbold, cw))
                    line_width += cw
                if line or not lines:
                    lines.append(line)
                total_h = line_h * len(lines)
                ty = by0 + (max(0, bh - total_h) if valign == "bottom" else
                            max(0, (bh - total_h) // 2) if valign in ("middle", "center") else 0)
                for li, chars_line in enumerate(lines):
                    total_w = sum(item[5] for item in chars_line)
                    cx = bx0 + (bw - total_w if align == "right" else
                                (bw - total_w) / 2 if align == "center" else 0)
                    for ch, cfont, ccolor, copacity, cbold, cw in chars_line:
                        stroke = 1 if cbold and not bold_path else 0
                        fill = _hex_to_rgba(ccolor, copacity)
                        draw.text((max(bx0, int(cx)), ty + li * line_h), ch, fill=fill,
                                  font=cfont, stroke_width=stroke, stroke_fill=fill)
                        cx += cw
            else:
                lines = _wrap_text(draw, str(tx.get("text", "")), font, max(1, bw))
                total_h = line_h * max(1, len(lines))
                ty = by0 + (max(0, bh - total_h) if valign == "bottom" else
                            max(0, (bh - total_h) // 2) if valign in ("middle", "center") else 0)
                stroke = 1 if bold and not bold_path else 0
                for li, line in enumerate(lines):
                    lw_ = draw.textlength(line, font=font)
                    lx = bx0 + (bw - lw_ if align == "right" else
                                (bw - lw_) / 2 if align == "center" else 0)
                    fill = _hex_to_rgba(color, opacity)
                    draw.text((max(bx0, int(lx)), ty + li * line_h), line, fill=fill,
                              font=font, stroke_width=stroke, stroke_fill=fill)

        out = preview_dir / f"slide_{idx:02d}.png"
        canvas.convert("RGB").save(out)
        print(f"Preview: {out}")


def main() -> None:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("layout", help="deck.json / layout.json path.")
    ap.add_argument("out", help="Output .pptx path.")
    ap.add_argument("--preview-dir", help="If set, also render PNG previews here for QA.")
    ap.add_argument("--font-dir", help="Task-local licensed font directory; also used by previews.")
    ap.add_argument("--font-manifest", help="Font manifest; defaults to FONT_DIR/font-manifest.json.")
    ap.add_argument("--embed-fonts", action="store_true", help="Post-process the generated PPTX with OOXML font parts.")
    ap.add_argument("--embedding-report", help="JSON report for the OOXML font embedding step.")
    args = ap.parse_args()

    lp = Path(args.layout)
    if not lp.exists():
        _die(f"layout file not found: {lp}")
    deck = _expand_components(_load_deck(lp))
    out_path = Path(args.out).resolve()
    if args.font_dir:
        deck["font_dir"] = str(Path(args.font_dir).resolve())
    if args.embed_fonts:
        font_dir = args.font_dir or deck.get("font_dir")
        manifest = args.font_manifest or deck.get("font_manifest")
        if not font_dir and not manifest:
            _die("--embed-fonts requires --font-dir or --font-manifest")
        from embed_fonts import embed_pptx_fonts, load_specs
        try:
            specs, _ = load_specs(font_dir, manifest, [])
        except Exception as exc:
            _die(f"font embedding input invalid: {exc}")
        out_path.parent.mkdir(parents=True, exist_ok=True)
        with tempfile.NamedTemporaryFile(prefix=f".{out_path.stem}-", suffix=".pptx", dir=out_path.parent, delete=False) as temporary:
            staging_path = Path(temporary.name)
        try:
            build_pptx(deck, staging_path)
            report_path = Path(args.embedding_report).resolve() if args.embedding_report else out_path.with_name(f"{out_path.stem}.font-embedding.json")
            result = embed_pptx_fonts(staging_path, out_path, specs, report_path, overwrite=True)
            if not result.get("valid"):
                _die(f"font embedding failed; see {report_path}")
        finally:
            if staging_path.exists():
                os.unlink(staging_path)
    else:
        build_pptx(deck, out_path)
    if args.preview_dir:
        render_previews(deck, Path(args.preview_dir))


if __name__ == "__main__":
    main()
