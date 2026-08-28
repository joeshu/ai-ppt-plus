#!/usr/bin/env python3
"""Native PowerPoint primitive writers.

This module owns text, shape, group, table and chart creation. It deliberately
does not know about the command line, component libraries, previews or font
embedding; the authoring backend supplies the slide and coordinate context.
"""
from __future__ import annotations

import sys


def _die(message: str, code: int = 2):
    print(f"Error: {message}", file=sys.stderr)
    raise SystemExit(code)


def text_size_pt(item: dict, slide_height_pt: float, reference_height: float, default=None) -> float:
    """Convert ratio, pixel or absolute text sizes to points."""
    if item.get("size") is not None:
        return float(item["size"])
    if item.get("size_ratio") is not None:
        return float(item["size_ratio"]) * slide_height_pt
    if item.get("size_pct") is not None:
        return float(item["size_pct"]) / 100.0 * slide_height_pt
    if item.get("size_px") is not None and reference_height:
        return float(item["size_px"]) * slide_height_pt / reference_height
    return float(default) if default is not None else 18.0


def _set_run_fonts(run, name: str):
    """Set latin, east-asian and complex-script typefaces."""
    from pptx.oxml.ns import qn

    run.font.name = name
    run_properties = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        element = run_properties.find(qn(tag))
        if element is None:
            element = run_properties.makeelement(qn(tag), {})
            run_properties.append(element)
        element.set("typeface", name)


def _set_run_alpha(run, opacity: float):
    from pptx.oxml.ns import qn

    opacity = max(0.0, min(1.0, float(opacity)))
    run_properties = run._r.get_or_add_rPr()
    solid = run_properties.find(qn("a:solidFill"))
    if solid is None:
        return
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:
        return
    for old in srgb.findall(qn("a:alpha")):
        srgb.remove(old)
    srgb.append(srgb.makeelement(qn("a:alpha"), {"val": str(int(opacity * 100000))}))


def _hex_to_rgb(value: str):
    from pptx.dml.color import RGBColor

    raw = str(value).strip().lstrip("#")
    if len(raw) == 3:
        raw = "".join(char * 2 for char in raw)
    if len(raw) != 6:
        return RGBColor(0x11, 0x11, 0x11)
    return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))


def _hex_to_tuple(value: str):
    raw = str(value).strip().lstrip("#")
    if len(raw) == 3:
        raw = "".join(char * 2 for char in raw)
    if len(raw) != 6:
        return (17, 17, 17)
    return (int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))


def _hex_to_rgba(value: str, opacity: float = 1.0):
    red, green, blue = _hex_to_tuple(value)
    return red, green, blue, int(max(0.0, min(1.0, float(opacity))) * 255)


def _set_fill_alpha(shape, opacity: float):
    """Add an alpha transform to a solid shape fill."""
    from pptx.oxml.ns import qn

    shape_properties = shape._element.spPr
    solid = shape_properties.find(qn("a:solidFill"))
    if solid is None:
        return
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:
        return
    for old in srgb.findall(qn("a:alpha")):
        srgb.remove(old)
    alpha = int(max(0.0, min(1.0, opacity)) * 100000)
    srgb.append(srgb.makeelement(qn("a:alpha"), {"val": str(alpha)}))


def _add_outer_shadow(shape, blur_pt=6.0, dist_pt=3.0, alpha=0.35):
    """Best-effort soft drop shadow for card-like shapes."""
    from pptx.oxml.ns import qn
    from pptx.util import Pt

    shape_properties = shape._element.spPr
    for old in shape_properties.findall(qn("a:effectLst")):
        shape_properties.remove(old)
    effects = shape_properties.makeelement(qn("a:effectLst"), {})
    shadow = effects.makeelement(qn("a:outerShdw"), {
        "blurRad": str(int(Pt(blur_pt))),
        "dist": str(int(Pt(dist_pt))),
        "dir": "5400000",
        "rotWithShape": "0",
    })
    color = shadow.makeelement(qn("a:srgbClr"), {"val": "000000"})
    color.append(color.makeelement(qn("a:alpha"), {"val": str(int(alpha * 100000))}))
    shadow.append(color)
    effects.append(shadow)
    shape_properties.append(effects)


def _set_gradient_fill(shape, gradient: dict):
    """Write a deterministic PresentationML linear gradient fill."""
    from pptx.oxml.ns import qn

    shape_properties = shape._element.spPr
    for tag in ("a:solidFill", "a:noFill", "a:gradFill"):
        old = shape_properties.find(qn(tag))
        if old is not None:
            shape_properties.remove(old)
    stops = gradient.get("stops", []) if isinstance(gradient, dict) else []
    if len(stops) < 2:
        _die("gradient fill requires at least two color stops")
    gradient_fill = shape_properties.makeelement(qn("a:gradFill"), {"rotWithShape": "1"})
    stop_list = gradient_fill.makeelement(qn("a:gsLst"), {})
    for stop in stops:
        if not isinstance(stop, dict) or not isinstance(stop.get("color"), str):
            _die("gradient stops require color fields")
        raw_position = float(stop.get("position", stop.get("pos", 0)))
        position = max(0, min(100000, int(raw_position * (100000 if 0 <= raw_position <= 1 else 1000))))
        gradient_stop = stop_list.makeelement(qn("a:gs"), {"pos": str(position)})
        color = gradient_stop.makeelement(qn("a:srgbClr"), {"val": stop["color"].lstrip("#")[:6]})
        if stop.get("opacity") is not None:
            alpha = int(max(0, min(1, float(stop["opacity"]))) * 100000)
            color.append(color.makeelement(qn("a:alpha"), {"val": str(alpha)}))
        gradient_stop.append(color)
        stop_list.append(gradient_stop)
    gradient_fill.append(stop_list)
    angle = float(gradient.get("angle", 0))
    gradient_fill.append(gradient_fill.makeelement(qn("a:lin"), {"ang": str(int(angle * 60000) % 21600000), "scaled": "1"}))
    line = shape_properties.find(qn("a:ln"))
    if line is not None:
        shape_properties.insert(list(shape_properties).index(line), gradient_fill)
    else:
        shape_properties.append(gradient_fill)


def apply_shape_fill(shape, spec: dict):
    gradient = spec.get("gradient") or spec.get("fill_gradient")
    if gradient:
        _set_gradient_fill(shape, gradient)
    elif spec.get("fill"):
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(spec["fill"])
        if spec.get("opacity") is not None:
            _set_fill_alpha(shape, float(spec["opacity"]))
    else:
        shape.fill.background()


def set_alt_text(shape, text: str | None):
    if not text:
        return
    element = shape._element
    for attribute in ("nvSpPr", "nvPicPr", "nvGrpSpPr", "nvGraphicFramePr"):
        container = getattr(element, attribute, None)
        if container is not None and getattr(container, "cNvPr", None) is not None:
            container.cNvPr.set("descr", str(text))
            return


def shape_map():
    from pptx.enum.shapes import MSO_SHAPE

    return {
        "rect": MSO_SHAPE.RECTANGLE,
        "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "oval": MSO_SHAPE.OVAL,
        "ellipse": MSO_SHAPE.OVAL,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "chevron": MSO_SHAPE.CHEVRON,
        "right_arrow": MSO_SHAPE.RIGHT_ARROW,
        "left_arrow": MSO_SHAPE.LEFT_ARROW,
        "up_arrow": MSO_SHAPE.UP_ARROW,
        "down_arrow": MSO_SHAPE.DOWN_ARROW,
        "pentagon": MSO_SHAPE.REGULAR_PENTAGON,
        "hexagon": MSO_SHAPE.HEXAGON,
        "parallelogram": MSO_SHAPE.PARALLELOGRAM,
        "trapezoid": MSO_SHAPE.TRAPEZOID,
        "diamond": MSO_SHAPE.DIAMOND,
    }


def chart_map():
    from pptx.enum.chart import XL_CHART_TYPE

    return {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
    }


def add_shapes(slide, specs: list[dict], deck: dict, ref_w: float, ref_h: float, sw_emu: int, sh_emu: int):
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.util import Emu, Pt

    shape_types = shape_map()
    for index, spec in enumerate(specs, 1):
        fx = _frac(deck, spec, "x", ref_w)
        fy = _frac(deck, spec, "y", ref_h)
        fw = _frac(deck, spec, "w", ref_w)
        fh = _frac(deck, spec, "h", ref_h)
        left, top = Emu(int(fx * sw_emu)), Emu(int(fy * sh_emu))
        width, height = Emu(int(fw * sw_emu)), Emu(int(fh * sh_emu))
        shape_type = spec.get("type", "rounded_rect")
        if shape_type == "line":
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                left,
                top,
                Emu(int((fx + fw) * sw_emu)),
                Emu(int((fy + fh) * sh_emu)),
            )
            connector.line.color.rgb = _hex_to_rgb(spec.get("line", spec.get("fill", "#FFFFFF")))
            connector.line.width = Pt(float(spec.get("line_width", 1.5)))
            connector.name = str(spec.get("object_id") or spec.get("name") or f"shape-{index:02d}")
            continue

        shape = slide.shapes.add_shape(shape_types.get(shape_type, MSO_SHAPE.ROUNDED_RECTANGLE), left, top, width, height)
        shape.name = str(spec.get("object_id") or spec.get("name") or f"shape-{index:02d}")
        if shape_type == "rounded_rect" and "radius" in spec:
            try:
                shape.adjustments[0] = float(spec["radius"])
            except Exception:
                pass
        apply_shape_fill(shape, spec)
        if spec.get("line"):
            shape.line.color.rgb = _hex_to_rgb(spec["line"])
            shape.line.width = Pt(float(spec.get("line_width", 1.0)))
        else:
            shape.line.fill.background()
        if spec.get("shadow"):
            _add_outer_shadow(shape)
        else:
            shape.shadow.inherit = False
        if spec.get("rotation"):
            shape.rotation = float(spec["rotation"])
        set_alt_text(shape, spec.get("alt_text"))


def add_groups(slide, specs: list[dict], deck: dict, ref_w: float, ref_h: float, sw_emu: int, sh_emu: int):
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.util import Emu, Pt

    shape_types = shape_map()
    for group_index, group_spec in enumerate(specs, 1):
        group = slide.shapes.add_group_shape()
        group.name = str(group_spec.get("object_id") or group_spec.get("name") or f"group-{group_index:02d}")
        local = group_spec.get("children_coordinate_space") == "local"
        has_box = all(key in group_spec for key in ("x", "y", "w", "h"))
        gx = _frac(deck, group_spec, "x", ref_w) if has_box else 0
        gy = _frac(deck, group_spec, "y", ref_h) if has_box else 0
        gw = _frac(deck, group_spec, "w", ref_w) if has_box else 1
        gh = _frac(deck, group_spec, "h", ref_h) if has_box else 1
        for child_index, raw_child in enumerate(group_spec.get("children", []), 1):
            if not isinstance(raw_child, dict):
                continue
            child = dict(raw_child)
            if local:
                child["x"] = gx + _frac(deck, child, "x", ref_w) * gw
                child["y"] = gy + _frac(deck, child, "y", ref_h) * gh
                child["w"] = _frac(deck, child, "w", ref_w) * gw
                child["h"] = _frac(deck, child, "h", ref_h) * gh
                child_deck = dict(deck, units="fraction")
            else:
                child_deck = deck
            cfx = _frac(child_deck, child, "x", ref_w)
            cfy = _frac(child_deck, child, "y", ref_h)
            cfw = _frac(child_deck, child, "w", ref_w)
            cfh = _frac(child_deck, child, "h", ref_h)
            child_name = str(child.get("object_id") or child.get("name") or f"{group.name}-child-{child_index:02d}")
            if child.get("type", "rounded_rect") == "line":
                connector = group.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Emu(int(cfx * sw_emu)),
                    Emu(int(cfy * sh_emu)),
                    Emu(int((cfx + cfw) * sw_emu)),
                    Emu(int((cfy + cfh) * sh_emu)),
                )
                connector.line.color.rgb = _hex_to_rgb(child.get("line", child.get("fill", "#FFFFFF")))
                connector.line.width = Pt(float(child.get("line_width", 1.5)))
                connector.name = child_name
                continue
            shape = group.shapes.add_shape(
                shape_types.get(child.get("type", "rounded_rect"), MSO_SHAPE.ROUNDED_RECTANGLE),
                Emu(int(cfx * sw_emu)),
                Emu(int(cfy * sh_emu)),
                Emu(int(cfw * sw_emu)),
                Emu(int(cfh * sh_emu)),
            )
            shape.name = child_name
            apply_shape_fill(shape, child)
            if child.get("line"):
                shape.line.color.rgb = _hex_to_rgb(child["line"])
                shape.line.width = Pt(float(child.get("line_width", 1.0)))
            else:
                shape.line.fill.background()
            if child.get("rotation"):
                shape.rotation = float(child["rotation"])
            set_alt_text(shape, child.get("alt_text"))
        group.shapes._recalculate_extents()
        set_alt_text(group, group_spec.get("alt_text"))


def add_tables(slide, specs: list[dict], deck: dict, theme: dict, ref_w: float, ref_h: float, sw_emu: int, sh_emu: int):
    from pptx.util import Emu, Pt

    for table_index, spec in enumerate(specs, 1):
        rows = spec.get("rows", [])
        columns = int(spec.get("columns") or (len(rows[0]) if rows and isinstance(rows[0], list) else 0))
        if not rows or columns <= 0:
            _die(f"table {table_index} requires non-empty rows and columns")
        table = slide.shapes.add_table(
            len(rows),
            columns,
            Emu(int(_frac(deck, spec, "x", ref_w) * sw_emu)),
            Emu(int(_frac(deck, spec, "y", ref_h) * sh_emu)),
            Emu(int(_frac(deck, spec, "w", ref_w) * sw_emu)),
            Emu(int(_frac(deck, spec, "h", ref_h) * sh_emu)),
        ).table
        graphic_frame = table._graphic_frame
        graphic_frame.name = str(spec.get("object_id") or spec.get("name") or f"table-{table_index:02d}")
        font = str(spec.get("font") or theme.get("font") or "Noto Sans CJK SC")
        header_fill = spec.get("header_fill") or theme.get("table_header_fill")
        body_fill = spec.get("fill") or theme.get("table_fill")
        for column_index, width in enumerate((spec.get("column_widths") or [])[:columns]):
            table.columns[column_index].width = Emu(int(float(width) * sw_emu)) if deck["units"] == "fraction" else int(width)
        for row_index, row in enumerate(rows):
            for column_index in range(columns):
                cell = table.cell(row_index, column_index)
                cell.text = str(row[column_index]) if column_index < len(row) else ""
                fill = header_fill if row_index == 0 else body_fill
                if fill:
                    cell.fill.solid()
                 …1866 tokens truncated…cation."""
from __future__ import annotations

import sys
from io import BytesIO
from pathlib import Path

from asset_placement import cleanup_temporary_files, svg_to_png
from atomic_output import atomic_write_bytes
from component_expander import _frac, _resolve
from pptx_primitives import _hex_to_rgba, _hex_to_tuple, text_size_pt


def find_cjk_font(font_dir=None, bold=False):
    if font_dir:
        local = sorted(Path(font_dir).glob("*.ttf")) + sorted(Path(font_dir).glob("*.otf")) + sorted(Path(font_dir).glob("*.ttc"))
        if local:
            if bold:
                bold_local = [path for path in local if any(token in path.stem.casefold() for token in ("bold", "medium", "semibold"))]
                if bold_local:
                    return str(bold_local[0])
            return str(local[0])
    candidates_bold = [
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/STHeiti Medium.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Bold.ttc",
    ]
    candidates = [
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/STHeiti Light.ttc",
        "/Library/Fonts/Arial Unicode.ttf",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    for candidate in candidates_bold + candidates if bold else candidates:
        if Path(candidate).exists():
            return candidate
    return None


def _poly_points(shape_type, x0, y0, x1, y1):
    """Approximate polygon vertices for preview rendering of non-rect shapes."""
    width, height = x1 - x0, y1 - y0
    if shape_type == "triangle":
        return [(x0 + width / 2, y0), (x1, y1), (x0, y1)]
    if shape_type == "diamond":
        return [(x0 + width / 2, y0), (x1, y0 + height / 2), (x0 + width / 2, y1), (x0, y0 + height / 2)]
    if shape_type == "right_arrow":
        margin = height * 0.30
        return [(x0, y0 + margin), (x1 - width * 0.4, y0 + margin), (x1 - width * 0.4, y0), (x1, y0 + height / 2), (x1 - width * 0.4, y1), (x1 - width * 0.4, y1 - margin), (x0, y1 - margin)]
    if shape_type == "left_arrow":
        margin = height * 0.30
        return [(x1, y0 + margin), (x0 + width * 0.4, y0 + margin), (x0 + width * 0.4, y0), (x0, y0 + height / 2), (x0 + width * 0.4, y1), (x0 + width * 0.4, y1 - margin), (x1, y1 - margin)]
    if shape_type == "up_arrow":
        margin = width * 0.30
        return [(x0 + margin, y1), (x0 + margin, y0 + height * 0.4), (x0, y0 + height * 0.4), (x0 + width / 2, y0), (x1, y0 + height * 0.4), (x1 - margin, y0 + height * 0.4), (x1 - margin, y1)]
    if shape_type == "chevron":
        notch = width * 0.25
        return [(x0, y0), (x1 - notch, y0), (x1, y0 + height / 2), (x1 - notch, y1), (x0, y1), (x0 + notch, y0 + height / 2)]
    if shape_type == "trapezoid":
        inset = width * 0.22
        return [(x0 + inset, y0), (x1 - inset, y0), (x1, y1), (x0, y1)]
    if shape_type == "parallelogram":
        skew = width * 0.22
        return [(x0 + skew, y0), (x1, y0), (x1 - skew, y1), (x0, y1)]
    if shape_type == "pentagon":
        import math

        center_x, center_y = x0 + width / 2, y0 + height / 2
        radius_x, radius_y = width / 2, height / 2
        return [(center_x + radius_x * math.sin(2 * math.pi * i / 5), center_y - radius_y * math.cos(2 * math.pi * i / 5)) for i in range(5)]
    if shape_type == "hexagon":
        return [(x0 + width * 0.25, y0), (x0 + width * 0.75, y0), (x1, y0 + height / 2), (x0 + width * 0.75, y1), (x0 + width * 0.25, y1), (x0, y0 + height / 2)]
    return None


def _wrap_text(draw, text, font, max_width):
    output = []
    for raw in text.split("\n"):
        if not raw:
            output.append("")
            continue
        line = ""
        for character in raw:
            trial = line + character
            if draw.textlength(trial, font=font) > max_width and line:
                output.append(line)
                line = character
            else:
                line = trial
        output.append(line)
    return output


def _open_asset(path: Path, temporary_files: list[Path]) -> Path:
    if path.suffix.casefold() == ".svg":
        raster = svg_to_png(path)
        temporary_files.append(raster)
        return raster
    return path


def render_previews(deck: dict, preview_dir: Path) -> None:
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        print("Preview skipped: Pillow not available.", file=sys.stderr)
        return

    assets_dir = Path(deck["assets_dir"])
    slide_width = float(deck["slide_width_in"])
    slide_height = float(deck["slide_height_in"])
    ratio = slide_width / slide_height
    if deck["units"] == "px" and deck.get("ref_width") and deck.get("ref_height"):
        canvas_width, canvas_height = int(deck["ref_width"]), int(deck["ref_height"])
    else:
        canvas_width = 1600
        canvas_height = int(round(canvas_width / ratio))
    slide_height_pt = slide_height * 72.0
    ref_width = float(deck.get("ref_width") or canvas_width)
    ref_height = float(deck.get("ref_height") or canvas_height)
    preview_dir.mkdir(parents=True, exist_ok=True)
    regular_path = find_cjk_font(deck.get("font_dir"), bold=False)
    bold_path = find_cjk_font(deck.get("font_dir"), bold=True)
    temporary_files: list[Path] = []

    try:
        for index, slide_spec in enumerate(deck["slides"], 1):
            canvas = Image.new("RGBA", (canvas_width, canvas_height), (255, 255, 255, 255))
            background = slide_spec.get("background")
            if background:
                path = _resolve(assets_dir, background)
                if path.exists():
                    with Image.open(path) as image:
                        canvas.paste(image.convert("RGBA").resize((canvas_width, canvas_height)), (0, 0))

            frame = slide_spec.get("frame")
            if frame:
                path = _resolve(assets_dir, frame)
                if path.exists():
                    with Image.open(path) as image:
                        canvas.alpha_composite(image.convert("RGBA").resize((canvas_width, canvas_height)), (0, 0))

            for shape_spec in slide_spec.get("shapes", []):
                fx = _frac(deck, shape_spec, "x", "x", ref_width)
                fy = _frac(deck, shape_spec, "y", "y", ref_height)
                fw = _frac(deck, shape_spec, "w", "x", ref_width)
                fh = _frac(deck, shape_spec, "h", "y", ref_height)
                x0, y0 = int(fx * canvas_width), int(fy * canvas_height)
                x1, y1 = int((fx + fw) * canvas_width), int((fy + fh) * canvas_height)
                overlay = Image.new("RGBA", (canvas_width, canvas_height), (0, 0, 0, 0))
                draw = ImageDraw.Draw(overlay)
                shape_type = shape_spec.get("type", "rounded_rect")
                alpha = int(float(shape_spec.get("opacity", 1.0)) * 255) if shape_spec.get("fill") else 0
                fill = _hex_to_tuple(shape_spec["fill"]) + (alpha,) if shape_spec.get("fill") else None
                line = _hex_to_tuple(shape_spec["line"]) + (255,) if shape_spec.get("line") else None
                line_width = int(round(float(shape_spec.get("line_width", 1.0)) * canvas_height / (slide_height * 72.0)))
                line_width = max(1, line_width) if line else 0
                polygon = _poly_points(shape_type, x0, y0, x1, y1)
                if shape_type == "line":
                    draw.line([x0, y0, x1, y1], fill=line or (255, 255, 255, 255), width=max(1, line_width or 2))
                elif shape_type in ("oval", "ellipse"):
                    draw.ellipse([x0, y0, x1, y1], fill=fill, outline=line, width=line_width)
                elif shape_type == "rounded_rect":
                    radius = int(float(shape_spec.get("radius", 0.12)) * min(x1 - x0, y1 - y0))
                    draw.rounded_rectangle([x0, y0, x1, y1], radius=max(1, radius), fill=fill, outline=line, width=line_width)
                elif polygon is not None:
                    draw.polygon(polygon, fill=fill, outline=line, width=line_width)
                else:
                    draw.rectangle([x0, y0, x1, y1], fill=fill, outline=line, width=line_width)
                canvas.alpha_composite(overlay)

            for panel in slide_spec.get("panels", []):
                path = _resolve(assets_dir, panel["file"])
                if not path.exists():
                    continue
                fx = _frac(deck, panel, "x", "x", ref_width)
                fy = _frac(deck, panel, "y", "y", ref_height)
                fw = _frac(deck, panel, "w", "x", ref_width)
                fh = _frac(deck, panel, "h", "y", ref_height)
                with Image.open(path) as image:
                    panel_image = image.convert("RGBA").resize((max(1, int(fw * canvas_width)), max(1, int(fh * canvas_height))))
                canvas.alpha_composite(panel_image, (int(fx * canvas_width), int(fy * canvas_height)))

            for icon in slide_spec.get("icons", []):
                path = _resolve(assets_dir, icon["file"])
                if not path.exists():
                    continue
                path = _open_asset(path, temporary_files)
                fx = _frac(deck, icon, "x", "x", ref_width)
                fy = _frac(deck, icon, "y", "y", ref_height)
                fw = _frac(deck, icon, "w", "x", ref_width)
                fh = _frac(deck, icon, "h", "y", ref_height)
                with Image.open(path) as image:
                    icon_image = image.convert("RGBA").resize((max(1, int(fw * canvas_width)), max(1, int(fh * canvas_height))))
                canvas.alpha_composite(icon_image, (int(fx * canvas_width), int(fy * canvas_height)))

            draw = ImageDraw.Draw(canvas)
            for text_spec in slide_spec.get("texts", []):
                fx = _frac(deck, text_spec, "x", "x", ref_width)
                fy = _frac(deck, text_spec, "y", "y", ref_height)
                fw = _frac(deck, text_spec, "w", "x", ref_width)
                size_pt = text_size_pt(text_spec, slide_height_pt, ref_height)
                pixel_size = max(8, int(round(size_pt * canvas_height / slide_height_pt)))
                bold = bool(text_spec.get("bold", False))
                font_path = (bold_path if bold else regular_path) or regular_path
                try:
                    font = ImageFont.truetype(font_path, pixel_size) if font_path else ImageFont.load_default()
                except Exception:
                    font = ImageFont.load_default()
                color = text_spec.get("color", "#111111")
                align = text_spec.get("align", "left")
                valign = text_spec.get("valign", "top")
                opacity = float(text_spec.get("opacity", 1.0))
                box_x, box_y = int(fx * canvas_width), int(fy * canvas_height)
                box_width = int(fw * canvas_width)
                box_height = int(_frac(deck, text_spec, "h", "y", ref_height) * canvas_height)
                line_height = int(pixel_size * 1.3)
                runs = text_spec.get("runs")
                if runs:
                    characters = []
                    for run_spec in runs:
                        run_bold = bool(run_spec.get("bold", bold))
                        run_size = max(8, int(round(text_size_pt(run_spec, slide_height_pt, ref_height, default=size_pt) * canvas_height / slide_height_pt)))
                        run_path = (bold_path if run_bold else regular_path) or regular_path
                        try:
                            run_font = ImageFont.truetype(run_path, run_size) if run_path else ImageFont.load_default()
                        except Exception:
                            run_font = ImageFont.load_default()
                        for character in str(run_spec.get("text", "")):
                            characters.append((character, run_font, run_spec.get("color", color), float(run_spec.get("opacity", opacity)), run_bold))
                    lines, line, line_width = [], [], 0.0
                    for character, character_font, character_color, character_opacity, character_bold in characters:
                        if character == "\n":
                            lines.append(line)
                            line, line_width = [], 0.0
                            continue
                        character_width = draw.textlength(character, font=character_font)
                        if line and line_width + character_width > box_width:
                            lines.append(line)
                            line, line_width = [], 0.0
                        line.append((character, character_font, character_color, character_opacity, character_bold, character_width))
                        line_width += character_width
                    if line or not lines:
                        lines.append(line)
                    total_height = line_height * len(lines)
                    text_y = box_y + (max(0, box_height - total_height) if valign == "bottom" else max(0, (box_height - total_height) // 2) if valign in ("middle", "center") else 0)
                    for line_index, line_items in enumerate(lines):
                        total_width = sum(item[5] for item in line_items)
                        cursor_x = box_x + (box_width - total_width if align == "right" else (box_width - total_width) / 2 if align == "center" else 0)
                        for character, character_font, character_color, character_opacity, character_bold, character_width in line_items:
                            stroke = 1 if character_bold and not bold_path else 0
                            fill = _hex_to_rgba(character_color, character_opacity)
                            draw.text((max(box_x, int(cursor_x)), text_y + line_index * line_height), character, fill=fill, font=character_font, stroke_width=stroke, stroke_fill=fill)
                            cursor_x += character_width
                else:
                    lines = _wrap_text(draw, str(text_spec.get("text", "")), font, max(1, box_width))
                    total_height = line_height * max(1, len(lines))
                    text_y = box_y + (max(0, box_height - total_height) if valign == "bottom" else max(0, (box_height - total_height) // 2) if valign in ("middle", "center") else 0)
                    stroke = 1 if bold and not bold_path else 0
                    for line_index, line in enumerate(lines):
                        line_width = draw.textlength(line, font=font)
                        text_x = box_x + (box_width - line_width if align == "right" else (box_width - line_width) / 2 if align == "center" else 0)
                        fill = _hex_to_rgba(color, opacity)
                        draw.text((max(box_x, int(text_x)), text_y + line_index * line_height), line, fill=fill, font=font, stroke_width=stroke, stroke_fill=fill)

            output = preview_dir / f"slide_{index:02d}.png"
            buffer = BytesIO()
            canvas.convert("RGB").save(buffer, format="PNG")
            atomic_write_bytes(output, buffer.getvalue(), suffix=".tmp.png")
            print(f"Preview: {output}")
    finally:
        cleanup_temporary_files(temporary_files)