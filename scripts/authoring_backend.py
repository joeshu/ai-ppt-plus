#!/usr/bin/env python3
"""Python-pptx authoring backend for the editable deck contract."""
from __future__ import annotations

import sys
import tempfile
from pathlib import Path

from asset_placement import (
    add_background,
    add_frame,
    add_icons,
    add_panels,
    cleanup_temporary_files,
    replace_svg_media,
)
from atomic_output import atomic_replace, atomic_save_presentation
from component_expander import _choose_slide_layout
from pptx_primitives import add_charts, add_groups, add_shapes, add_tables, add_texts


# Use the family name exposed by the bundled redistributable font. Keep the
# declared family and the asset name table aligned so renderers do not
# silently substitute a thinner host face.
DEFAULT_FONT_FAMILY = "Noto Sans CJK SC"
BACKEND_ID = "python-pptx"


def _die(message: str, code: int = 2):
    print(f"Error: {message}", file=sys.stderr)
    raise SystemExit(code)


def _normalized_theme(deck: dict) -> dict:
    theme = deck.get("theme", {}) if isinstance(deck.get("theme", {}), dict) else {}
    normalized = dict(theme)
    # A missing task font must not silently select the proprietary Microsoft
    # YaHei family. The font asset gate remains responsible for proving that
    # the chosen family is available and licensed for delivery.
    normalized.setdefault("font", deck.get("font_family", DEFAULT_FONT_FAMILY))
    return normalized


def build_pptx(deck: dict, out_path: Path) -> None:
    """Create a PPTX using the stable backend contract."""
    from pptx import Presentation
    from pptx.util import Emu

    assets_dir = Path(deck["assets_dir"])
    slide_width = float(deck["slide_width_in"])
    slide_height = float(deck["slide_height_in"])
    slide_width_emu = int(round(slide_width * 914400))
    slide_height_emu = int(round(slide_height * 914400))
    reference_width = float(deck.get("ref_width") or 0)
    reference_height = float(deck.get("ref_height") or 0)
    presentation = Presentation()
    presentation.slide_width = Emu(slide_width_emu)
    presentation.slide_height = Emu(slide_height_emu)
    theme = _normalized_theme(deck)
    svg_assets: list[tuple[int, str, Path]] = []
    temporary_files: list[Path] = []

    try:
        for slide_no, slide_spec in enumerate(deck["slides"], 1):
            slide = presentation.slides.add_slide(_choose_slide_layout(presentation, slide_spec, theme, deck))
            add_background(slide, slide_spec, assets_dir, slide_width_emu, slide_height_emu)
            add_frame(slide, slide_spec, assets_dir, slide_width_emu, slide_height_emu)
            # Semantic panel images are structural substrates.  Place them
            # before native overlays so badges, legend keys and bullet marks
            # remain visible and independently editable.  The previous order
            # silently covered those overlays whenever a panel occupied the
            # same region.
            add_panels(slide, slide_spec.get("panels", []), assets_dir, deck, reference_width, reference_height, slide_width_emu, slide_height_emu, slide_no)
            add_shapes(slide, slide_spec.get("shapes", []), deck, reference_width, reference_height, slide_width_emu, slide_height_emu)
            add_groups(slide, slide_spec.get("groups", []), deck, reference_width, reference_height, slide_width_emu, slide_height_emu)
            add_tables(slide, slide_spec.get("tables", []), deck, theme, reference_width, reference_height, slide_width_emu, slide_height_emu)
            add_charts(slide, slide_spec.get("charts", []), deck, theme, reference_width, reference_height, slide_width_emu, slide_height_emu)
            notes = slide_spec.get("speaker_notes") or slide_spec.get("notes")
            if notes:
                slide.notes_slide.notes_text_frame.text = str(notes)
            add_icons(slide, slide_spec.get("icons", []), assets_dir, deck, reference_width, reference_height, slide_width_emu, slide_height_emu, slide_no, svg_assets, temporary_files)
            add_texts(slide, slide_spec.get("texts", []), deck, theme, reference_width, reference_height, slide_width_emu, slide_height_emu)

        if svg_assets:
            # Keep the final target untouched until both python-pptx output
            # and native SVG package rewriting have completed successfully.
            def write_package(path: Path) -> None:
                presentation.save(str(path))
                replace_svg_media(path, svg_assets)

            atomic_replace(out_path, write_package, suffix=".tmp.pptx")
        else:
            atomic_save_presentation(presentation, out_path)
        print(f"Wrote {out_path}  ({len(deck['slides'])} slides)")
    finally:
        cleanup_temporary_files(temporary_files)


def build_with_embedded_fonts(
    deck: dict,
    out_path: Path,
    *,
    font_dir: str | None,
    font_manifest: str | None,
    embedding_report: Path | None,
) -> None:
    """Compose a staging package, then atomically publish its font-embedded form."""
    from embed_fonts import embed_pptx_fonts, load_specs

    try:
        specs, _ = load_specs(font_dir, font_manifest, [])
    except Exception as exc:
        _die(f"font embedding input invalid: {exc}")
    out_path = out_path.resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.NamedTemporaryFile(prefix=f".{out_path.stem}-", suffix=".staging.pptx", dir=out_path.parent, delete=False) as temporary:
        staging_path = Path(temporary.name)
    try:
        build_pptx(deck, staging_path)
        report_path = embedding_report or out_path.with_name(f"{out_path.stem}.font-embedding.json")
        result = embed_pptx_fonts(staging_path, out_path, specs, report_path, overwrite=True)
        if not result.get("valid"):
            _die(f"font embedding failed; see {report_path}")
    finally:
        staging_path.unlink(missing_ok=True)
