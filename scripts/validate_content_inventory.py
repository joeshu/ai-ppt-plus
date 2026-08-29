#!/usr/bin/env python3
"""Validate the independent visible-content inventory for a slide deck.

The object manifest proves that declared objects exist.  This gate proves that
the visible text and chart annotations identified during source review were
not silently omitted and that the independent authority agrees with both the
text manifest and the final PPTX.
"""
from __future__ import annotations

import argparse
import hashlib
import json
import re
from pathlib import Path
from typing import Any

from atomic_output import atomic_write_json


SCHEMA = "ai-ppt-plus/content-inventory/v1"
AUTHORITIES = {"approved_outline", "user_transcription", "approved_outline_or_user_transcription"}
SHA256_RE = re.compile(r"^[0-9a-f]{64}$")


def _digest(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _read(path: Path) -> dict[str, Any]:
    value = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(value, dict):
        raise ValueError(f"JSON object required: {path}")
    return value


def _normal_text(value: Any) -> str:
    return str(value or "").replace("\r\n", "\n").replace("\r", "\n").replace("\v", "\n")


def _content(value: dict[str, Any]) -> str:
    if "content" in value:
        return _normal_text(value.get("content"))
    if "text" in value:
        return _normal_text(value.get("text"))
    runs = value.get("runs")
    if isinstance(runs, list):
        return _normal_text("".join(str(run.get("text", "")) for run in runs if isinstance(run, dict)))
    return ""


def _slide_entries(data: Any) -> list[tuple[int, dict[str, Any]]]:
    raw = data.get("slides") if isinstance(data, dict) else None
    if not isinstance(raw, list):
        return []
    entries: list[tuple[int, dict[str, Any]]] = []
    for index, item in enumerate(raw, 1):
        if not isinstance(item, dict):
            continue
        try:
            number = int(item.get("slide_no", index))
        except (TypeError, ValueError):
            number = 0
        entries.append((number, item))
    return entries


def _manifest_objects(data: dict[str, Any]) -> dict[tuple[int, str], dict[str, Any]]:
    result: dict[tuple[int, str], dict[str, Any]] = {}
    for slide_no, slide in _slide_entries(data):
        objects = slide.get("objects")
        if not isinstance(objects, list):
            continue
        for obj in objects:
            if isinstance(obj, dict) and obj.get("object_id"):
                result[(slide_no, str(obj["object_id"]))] = obj
    return result


def _manifest_text(data: dict[str, Any]) -> dict[tuple[int, str], dict[str, Any]]:
    result: dict[tuple[int, str], dict[str, Any]] = {}
    for slide_no, slide in _slide_entries(data):
        specs = slide.get("text_specs")
        if not isinstance(specs, list):
            continue
        for spec in specs:
            if not isinstance(spec, dict):
                continue
            ident = spec.get("text_id") or spec.get("object_id")
            if ident:
                result[(slide_no, str(ident))] = spec
    return result


def _shape_map(deck_path: Path) -> dict[tuple[int, str], list[Any]]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(deck_path))
    output: dict[tuple[int, str], list[Any]] = {}

    def walk(shapes):
        for shape in shapes:
            yield shape
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                yield from walk(shape.shapes)

    for slide_no, slide in enumerate(prs.slides, 1):
        for shape in walk(slide.shapes):
            output.setdefault((slide_no, str(shape.name)), []).append(shape)
    return output


def _is_native_textbox(shape) -> bool:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    return bool(
        getattr(shape, "has_text_frame", False)
        and getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.TEXT_BOX
        and not bool(getattr(shape, "is_placeholder", False))
    )


def _iter_inventory_text(slide_no: int, slide: dict[str, Any], errors: list[dict[str, Any]]):
    raw = slide.get("visible_text", slide.get("texts", []))
    if raw is None:
        raw = []
    if not isinstance(raw, list):
        errors.append({"severity": "blocker", "code": "visible_text_not_array", "slide_no": slide_no})
        raw = []
    for index, item in enumerate(raw, 1):
        if not isinstance(item, dict):
            errors.append({"severity": "blocker", "code": "visible_text_item_invalid", "slide_no": slide_no, "index": index})
            continue
        yield item, "visible_text", index

    charts = slide.get("charts", [])
    if charts is None:
        charts = []
    if not isinstance(charts, list):
        errors.append({"severity": "blocker", "code": "chart_inventory_not_array", "slide_no": slide_no})
        charts = []
    for chart_index, chart in enumerate(charts, 1):
        if not isinstance(chart, dict):
            errors.append({"severity": "blocker", "code": "chart_inventory_item_invalid", "slide_no": slide_no, "index": chart_index})
            continue
        chart_id = chart.get("chart_id") or chart.get("object_id")
        if not isinstance(chart_id, str) or not chart_id.strip():
            errors.append({"severity": "blocker", "code": "chart_id_missing", "slide_no": slide_no, "index": chart_index})
        representation = chart.get("representation")
        if representation not in {"native_chart", "static_line_primitives", "svg", "raster_fallback"}:
            errors.append({"severity": "blocker", "code": "chart_representation_invalid", "slide_no": slide_no, "chart_id": chart_id, "observed": representation})
        source_status = chart.get("source_data_status")
        if source_status not in {"verified", "unverified", "unavailable"}:
            errors.append({"severity": "blocker", "code": "chart_source_status_invalid", "slide_no": slide_no, "chart_id": chart_id, "observed": source_status})
        required_elements = chart.get("required_elements")
        if not isinstance(required_elements, list) or not required_elements:
            errors.append({"severity": "blocker", "code": "chart_required_elements_missing", "slide_no": slide_no, "chart_id": chart_id})
            required_elements = []
        visible_elements = chart.get("visible_elements")
        if not isinstance(visible_elements, dict):
            errors.append({"severity": "blocker", "code": "chart_visible_elements_missing", "slide_no": slide_no, "chart_id": chart_id})
            visible_elements = {}
        for element_name in required_elements:
            entries = visible_elements.get(element_name)
            if not isinstance(entries, list) or not entries:
                errors.append({"severity": "blocker", "code": "chart_visible_element_missing", "slide_no": slide_no, "chart_id": chart_id, "element": element_name})
                continue
            for item_index, item in enumerate(entries, 1):
                if not isinstance(item, dict):
                    errors.append({"severity": "blocker", "code": "chart_visible_element_invalid", "slide_no": slide_no, "chart_id": chart_id, "element": element_name, "index": item_index})
                    continue
                yield item, f"chart.{element_name}", item_index


def validate(
    inventory_path: Path,
    *,
    object_manifest_path: Path | None = None,
    text_manifest_path: Path | None = None,
    deck_path: Path | None = None,
    expected_pages: int | None = None,
) -> dict[str, Any]:
    inventory = _read(inventory_path)
    errors: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []
    if inventory.get("schema") != SCHEMA:
        errors.append({"severity": "blocker", "code": "schema_invalid", "observed": inventory.get("schema")})
    authority = inventory.get("authority")
    if authority not in AUTHORITIES:
        errors.append({"severity": "blocker", "code": "content_authority_invalid", "observed": authority, "allowed": sorted(AUTHORITIES)})
    source_reference = inventory.get("source_reference")
    source_sha256 = inventory.get("source_sha256")
    if source_reference:
        if not isinstance(source_reference, str):
            errors.append({"severity": "blocker", "code": "source_reference_invalid"})
        elif source_sha256 is None or not isinstance(source_sha256, str) or not SHA256_RE.fullmatch(source_sha256):
            errors.append({"severity": "blocker", "code": "source_hash_missing_or_invalid", "source_reference": source_reference, "observed": source_sha256})
        else:
            source_path = (inventory_path.parent / source_reference).resolve()
            observed_source_hash = _digest(source_path) if source_path.is_file() else None
            if source_path.is_file() and observed_source_hash != source_sha256:
                errors.append({"severity": "blocker", "code": "source_hash_mismatch", "source_reference": str(source_path), "expected": source_sha256, "observed": observed_source_hash})
            elif not source_path.is_file():
                warnings.append({"severity": "major", "code": "source_reference_not_local_file", "source_reference": source_reference})

    slides = _slide_entries(inventory)
    numbers = [number for number, _slide in slides]
    if not slides:
        errors.append({"severity": "blocker", "code": "slides_missing"})
    if len(numbers) != len(set(numbers)):
        errors.append({"severity": "blocker", "code": "duplicate_slide_number", "observed": numbers})
    expected = list(range(1, expected_pages + 1)) if expected_pages is not None else list(range(1, len(numbers) + 1))
    if sorted(numbers) != expected:
        errors.append({"severity": "blocker", "code": "slide_coverage_mismatch", "expected": expected, "observed": sorted(numbers)})

    object_data = _read(object_manifest_path) if object_manifest_path and object_manifest_path.is_file() else None
    text_data = _read(text_manifest_path) if text_manifest_path and text_manifest_path.is_file() else None
    object_map = _manifest_objects(object_data) if object_data else {}
    text_map = _manifest_text(text_data) if text_data else {}
    if object_manifest_path and not object_manifest_path.is_file():
        errors.append({"severity": "blocker", "code": "object_manifest_missing", "path": str(object_manifest_path)})
    if text_manifest_path and not text_manifest_path.is_file():
        errors.append({"severity": "blocker", "code": "text_manifest_missing", "path": str(text_manifest_path)})
    shape_map = _shape_map(deck_path) if deck_path and deck_path.is_file() else {}
    if deck_path and not deck_path.is_file():
        errors.append({"severity": "blocker", "code": "deck_missing", "path": str(deck_path)})

    text_count = 0
    chart_count = 0
    element_count = 0
    seen_ids: set[tuple[int, str]] = set()
    for slide_no, slide in slides:
        inventory_ids: set[str] = set()
        for item, kind, index in _iter_inventory_text(slide_no, slide, errors):
            text_count += 1
            ident = item.get("object_id") or item.get("text_id")
            content = _content(item)
            if not isinstance(ident, str) or not ident.strip():
                errors.append({"severity": "blocker", "code": "visible_text_id_missing", "slide_no": slide_no, "kind": kind, "index": index})
                continue
            key = (slide_no, ident)
            inventory_ids.add(ident)
            if key in seen_ids:
                errors.append({"severity": "blocker", "code": "visible_text_id_duplicate", "slide_no": slide_no, "object_id": ident})
            seen_ids.add(key)
            if not content:
                errors.append({"severity": "blocker", "code": "visible_text_content_missing", "slide_no": slide_no, "object_id": ident})
            if object_data is not None:
                obj = object_map.get(key)
                if obj is None:
                    errors.append({"severity": "blocker", "code": "visible_text_not_in_object_manifest", "slide_no": slide_no, "object_id": ident})
                else:
                    if obj.get("object_type") != "editable_text":
                        errors.append({"severity": "blocker", "code": "visible_text_object_not_editable_text", "slide_no": slide_no, "object_id": ident, "observed": obj.get("object_type")})
                    object_text = _content(obj.get("text_spec", {})) if isinstance(obj.get("text_spec"), dict) else _content(obj)
                    if object_text != content:
                        errors.append({"severity": "blocker", "code": "visible_text_object_mismatch", "slide_no": slide_no, "object_id": ident, "inventory": content, "object_manifest": object_text})
            if text_data is not None:
                spec = text_map.get(key)
                if spec is None:
                    errors.append({"severity": "blocker", "code": "visible_text_not_in_text_manifest", "slide_no": slide_no, "object_id": ident})
                elif _content(spec) != content:
                    errors.append({"severity": "blocker", "code": "visible_text_text_manifest_mismatch", "slide_no": slide_no, "object_id": ident, "inventory": content, "text_manifest": _content(spec)})
            if deck_path and deck_path.is_file():
                matches = shape_map.get(key, [])
                if len(matches) != 1:
                    errors.append({"severity": "blocker", "code": "visible_text_shape_count_mismatch", "slide_no": slide_no, "object_id": ident, "observed": len(matches)})
                elif not _is_native_textbox(matches[0]):
                    errors.append({"severity": "blocker", "code": "visible_text_shape_not_native_textbox", "slide_no": slide_no, "object_id": ident})
                elif _normal_text(matches[0].text) != content:
                    errors.append({"severity": "blocker", "code": "visible_text_pptx_mismatch", "slide_no": slide_no, "object_id": ident, "inventory": content, "pptx": _normal_text(matches[0].text)})
            if kind.startswith("chart."):
                element_count += 1

        charts = slide.get("charts", [])
        chart_count += len(charts) if isinstance(charts, list) else 0

        required_objects = slide.get("required_object_ids", [])
        if required_objects is not None and not isinstance(required_objects, list):
            errors.append({"severity": "blocker", "code": "required_object_ids_not_array", "slide_no": slide_no})
        elif object_data is not None:
            for ident in required_objects or []:
                if (slide_no, str(ident)) not in object_map:
                    errors.append({"severity": "blocker", "code": "required_object_missing", "slide_no": slide_no, "object_id": ident})
        if object_data is not None:
            for (object_slide, object_id), obj in object_map.items():
                if object_slide == slide_no and obj.get("object_type") == "editable_text" and object_id not in inventory_ids:
                    errors.append({"severity": "blocker", "code": "editable_text_not_in_content_inventory", "slide_no": slide_no, "object_id": object_id})
        if text_data is not None:
            for (text_slide, text_id) in text_map:
                if text_slide == slide_no and text_id not in inventory_ids:
                    errors.append({"severity": "blocker", "code": "text_spec_not_in_content_inventory", "slide_no": slide_no, "text_id": text_id})

    result = {
        "schema": "ai-ppt-plus/content-inventory-validation/v1",
        "valid": not errors,
        "status": "passed" if not errors else "blocked",
        "inventory": str(inventory_path.resolve()),
        "inventory_sha256": _digest(inventory_path),
        "authority": inventory.get("authority"),
        "source_reference": inventory.get("source_reference"),
        "source_sha256": inventory.get("source_sha256"),
        "object_manifest": str(object_manifest_path.resolve()) if object_manifest_path else None,
        "text_manifest": str(text_manifest_path.resolve()) if text_manifest_path else None,
        "deck": str(deck_path.resolve()) if deck_path else None,
        "slides": len(slides),
        "visible_text_count": text_count,
        "chart_count": chart_count,
        "chart_annotation_count": element_count,
        "errors": errors,
        "warnings": warnings,
        "human_visual_review_required": True,
    }
    return result


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("inventory")
    parser.add_argument("--object-manifest")
    parser.add_argument("--text-manifest")
    parser.add_argument("--deck")
    parser.add_argument("--expected-pages", type=int)
    parser.add_argument("--report")
    args = parser.parse_args()
    try:
        result = validate(
            Path(args.inventory).resolve(),
            object_manifest_path=Path(args.object_manifest).resolve() if args.object_manifest else None,
            text_manifest_path=Path(args.text_manifest).resolve() if args.text_manifest else None,
            deck_path=Path(args.deck).resolve() if args.deck else None,
            expected_pages=args.expected_pages,
        )
    except Exception as exc:
        result = {"schema": "ai-ppt-plus/content-inventory-validation/v1", "valid": False, "status": "invalid", "errors": [{"severity": "blocker", "code": "runtime_error", "message": f"{type(exc).__name__}: {exc}"}]}
    if args.report:
        atomic_write_json(Path(args.report).resolve(), result)
    print(json.dumps(result, ensure_ascii=False))
    return 0 if result.get("valid") is True else 2


if __name__ == "__main__":
    raise SystemExit(main())
