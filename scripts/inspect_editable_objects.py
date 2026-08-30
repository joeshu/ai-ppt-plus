#!/usr/bin/env python3
"""Reverse-audit final PPTX objects against a slide-object manifest."""
from __future__ import annotations

import argparse
import json
from pathlib import Path
from atomic_output import atomic_write_json


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("deck")
    ap.add_argument("--object-manifest")
    ap.add_argument("--require-independent-panels", action="store_true")
    ap.add_argument("--report")
    args = ap.parse_args()
    from pptx import Presentation
    prs = Presentation(args.deck)
    expected_by_slide: dict[int, dict[str, dict]] = {}
    if args.object_manifest:
        data = json.loads(Path(args.object_manifest).read_text(encoding="utf-8"))
        for slide_index, manifest_slide in enumerate(data.get("slides", []), 1):
            slide_no = int(manifest_slide.get("slide_no", slide_index))
            page_objects = expected_by_slide.setdefault(slide_no, {})
            for obj in manifest_slide.get("objects", []):
                object_id = obj.get("object_id")
                if object_id in page_objects:
                    raise ValueError(f"duplicate object_id on slide {slide_no}: {object_id}")
                page_objects[object_id] = obj
    errors, warnings, slides = [], [], []
    for si, slide in enumerate(prs.slides, 1):
        shapes = []
        for shape in slide.shapes:
            kind = str(getattr(shape, "shape_type", "unknown"))
            text = getattr(shape, "text", "") if hasattr(shape, "text") else ""
            shapes.append({"name": shape.name, "type": kind, "has_text": bool(text.strip()), "text": text[:160], "left": shape.left, "top": shape.top, "width": shape.width, "height": shape.height})
        slides.append({"slide": si, "shape_count": len(shapes), "shapes": shapes})
        names = {s["name"] for s in shapes}
        expected = expected_by_slide.get(si, {})
        for oid, obj in expected.items():
            if oid not in names:
                errors.append({"code": "manifest_object_missing_in_pptx", "slide": si, "object_id": oid, "role": obj.get("role")})
        if args.require_independent_panels:
            panel_ids = [oid for oid, obj in expected.items() if obj.get("role") in {"semantic-panel", "panel", "frame-panel"}]
            for oid in panel_ids:
                matches = [s for s in shapes if s["name"] == oid]
                if len(matches) != 1:
                    errors.append({"code": "panel_shape_count_not_one", "slide": si, "object_id": oid, "observed": len(matches)})
    whole_slide_images = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if str(getattr(shape, "shape_type", "")) == "PICTURE" and shape.width >= prs.slide_width * .95 and shape.height >= prs.slide_height * .95:
                whole_slide_images.append(shape.name)
    expected_object_count = sum(len(objects) for objects in expected_by_slide.values())
    if whole_slide_images and expected_object_count:
        warnings.append({"code": "whole_slide_picture_present", "names": whole_slide_images})
    result = {"schema": "ai-ppt-plus/editable-object-audit/v1", "valid": not errors, "deck": str(Path(args.deck).resolve()), "slides": slides, "expected_object_count": expected_object_count, "observed_shape_count": sum(s["shape_count"] for s in slides), "whole_slide_pictures": whole_slide_images, "errors": errors, "warnings": warnings, "human_visual_review_required": True}
    if args.report:
        out = Path(args.report); out.parent.mkdir(parents=True, exist_ok=True)
        atomic_write_json(out.resolve(), result)
    print(json.dumps(result, ensure_ascii=False))
    return 0 if result["valid"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
