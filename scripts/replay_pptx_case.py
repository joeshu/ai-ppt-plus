#!/usr/bin/env python3
"""Replay a real editable-PPTX regression case and emit evidence.

The runner is intentionally case-oriented.  It opens the actual PPTX, parses
OOXML (including ``a:tbl`` and merge spans), renders the deck with LibreOffice,
compares the render with the source reference, and performs a disposable
mutation smoke test on a table cell and a native panel.  It does not infer
editability from object names alone.
"""
from __future__ import annotations

import argparse
import hashlib
import json
import math
import os
import posixpath
import re
import shutil
import subprocess
import tempfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
from zipfile import ZipFile

from lxml import etree
from PIL import Image, ImageChops, ImageStat
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches


SCHEMA = "ai-ppt-plus/pptx-case-evaluation/v1"
NS = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "pr": "http://schemas.openxmlformats.org/package/2006/relationships",
}
SHA256_RE = re.compile(r"^[0-9a-f]{64}$")
PANEL_NAME = re.compile(r"^(?:top-|flow-card-|commission-card-|policy-|monthly-|conditions-)")


def utc_now() -> str:
    return datetime.now(timezone.utc).isoformat()


def digest(path: Path) -> str:
    hasher = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            hasher.update(chunk)
    return hasher.hexdigest()


def write_json(path: Path, value: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    temporary = path.with_name(path.name + ".tmp")
    temporary.write_text(json.dumps(value, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    os.replace(temporary, path)


def load_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding="utf-8"))


def local_name(element: etree._Element) -> str:
    return etree.QName(element).localname


def shape_name(element: etree._Element) -> str:
    node = element.find(".//p:cNvPr", namespaces=NS)
    return str(node.get("name") if node is not None else "")


def shape_extent(element: etree._Element) -> tuple[int, int]:
    ext = element.find(".//a:xfrm/a:ext", namespaces=NS)
    if ext is None:
        ext = element.find(".//p:xfrm/a:ext", namespaces=NS)
    if ext is None:
        return 0, 0
    return int(ext.get("cx", 0)), int(ext.get("cy", 0))


def relationship_targets(zip_file: ZipFile, slide_name: str) -> dict[str, str]:
    rel_name = f"ppt/slides/_rels/{Path(slide_name).name}.rels"
    if rel_name not in zip_file.namelist():
        return {}
    root = etree.fromstring(zip_file.read(rel_name))
    result: dict[str, str] = {}
    for relation in root.xpath("./pr:Relationship", namespaces=NS):
        target = str(relation.get("Target") or "")
        result[str(relation.get("Id"))] = posixpath.normpath(posixpath.join("ppt/slides", target))
    return result


def image_relationship(element: etree._Element) -> str | None:
    blip = element.find(".//a:blip", namespaces=NS)
    return str(blip.get(f"{{{NS['r']}}}embed")) if blip is not None else None


def table_record(table: etree._Element) -> dict[str, Any]:
    rows = table.xpath("./a:tr", namespaces=NS)
    columns = len(table.xpath("./a:tblGrid/a:gridCol", namespaces=NS))
    if not columns and rows:
        columns = len(rows[0].xpath("./a:tc", namespaces=NS))
    values: list[list[str]] = []
    merges: list[list[int]] = []
    for row_no, row in enumerate(rows):
        row_values: list[str] = []
        for col_no, cell in enumerate(row.xpath("./a:tc", namespaces=NS)):
            row_values.append("".join(str(item) for item in cell.xpath(".//a:t/text()", namespaces=NS)))
            row_span = int(cell.get("rowSpan", "1"))
            col_span = int(cell.get("gridSpan", "1"))
            if row_span > 1 or col_span > 1:
                merges.append([row_no, col_no, row_no + row_span - 1, col_no + col_span - 1])
        values.append(row_values)
    return {
        "rows": len(rows),
        "columns": columns,
        "values": values,
        "merged_cells": merges,
    }


def source_reference(source_pptx: Path, output_dir: Path, explicit: Path | None) -> tuple[Path, dict[str, str]]:
    output_dir.mkdir(parents=True, exist_ok=True)
    hashes = {"source_pptx_sha256": digest(source_pptx)}
    if explicit is not None:
        reference = explicit.resolve()
    else:
        candidates: list[tuple[int, str, bytes]] = []
        with ZipFile(source_pptx) as zip_file:
            for name in zip_file.namelist():
                if not name.startswith("ppt/media/"):
                    continue
                data = zip_file.read(name)
                try:
                    with Image.open(__import__("io").BytesIO(data)) as image:
                        # The source case uses an opaque RGB full-slide
                        # background and an RGBA green-screen frame. Prefer
                        # the opaque slide-sized asset as visual authority.
                        score = image.width * image.height
                        if image.mode not in {"RGBA", "LA", "P"}:
                            score += 10**12
                        candidates.append((score, name, data))
                except Exception:
                    continue
        if not candidates:
            raise ValueError("source PPTX has no readable raster reference")
        _, name, data = max(candidates, key=lambda item: item[0])
        reference = output_dir / "source-reference.png"
        reference.write_bytes(data)
        hashes["source_reference_media"] = name
    if not reference.is_file():
        raise FileNotFoundError(reference)
    hashes["source_reference_sha256"] = digest(reference)
    if explicit is not None:
        hashes["source_reference_path"] = str(reference)
    return reference, hashes


def render_deck(deck: Path, render_dir: Path) -> tuple[bool, str | None]:
    render_dir.mkdir(parents=True, exist_ok=True)
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    pdftoppm = shutil.which("pdftoppm")
    if not soffice or not pdftoppm:
        return False, "LibreOffice/Poppler renderer is unavailable"
    with tempfile.TemporaryDirectory(prefix="pptx-replay-", dir=render_dir) as temporary:
        temporary_dir = Path(temporary)
        completed = subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(temporary_dir), str(deck)],
            capture_output=True, text=True, check=False,
        )
        if completed.returncode != 0:
            return False, completed.stderr[-2000:] or completed.stdout[-2000:]
        pdf = temporary_dir / f"{deck.stem}.pdf"
        if not pdf.is_file():
            pdfs = list(temporary_dir.glob("*.pdf"))
            if not pdfs:
                return False, "LibreOffice produced no PDF"
            pdf = pdfs[0]
        completed = subprocess.run(
            [pdftoppm, "-png", "-r", "144", str(pdf), str(temporary_dir / "render")],
            capture_output=True, text=True, check=False,
        )
        if completed.returncode != 0:
            return False, completed.stderr[-2000:] or completed.stdout[-2000:]
        for old in render_dir.glob("slide-*.png"):
            old.unlink()
        rendered = sorted(temporary_dir.glob("render-*.png"))
        for index, image in enumerate(rendered, 1):
            shutil.copy2(image, render_dir / f"slide-{index}.png")
    return bool(list(render_dir.glob("slide-*.png"))), None


def visual_compare(rendered: Path, reference: Path) -> dict[str, Any]:
    with Image.open(rendered) as render_image, Image.open(reference) as reference_image:
        render = render_image.convert("RGB")
        source = reference_image.convert("RGB")
        ratio_delta = abs((render.width / render.height) - (source.width / source.height))
        resized = render.size != source.size
        if resized:
            render = render.resize(source.size, Image.Resampling.LANCZOS)
        diff = ImageChops.difference(render, source)
        stat = ImageStat.Stat(diff)
        channels = max(1, len(stat.mean))
        mae = sum(stat.mean) / (255.0 * channels)
        rmse = sum(stat.rms) / (255.0 * channels)
        gray = diff.convert("L")
        histogram = gray.histogram()
        total = render.width * render.height
        pixel_match = sum(histogram[:9]) / total if total else 0.0
        score = max(0.0, 1.0 - (0.65 * mae + 0.35 * rmse))
    issues = []
    if ratio_delta > 0.015:
        issues.append({"code": "aspect_ratio_mismatch", "delta": round(ratio_delta, 6)})
    return {
        "valid": not issues,
        "rendered": str(rendered.resolve()),
        "rendered_sha256": digest(rendered),
        "reference": str(reference.resolve()),
        "reference_sha256": digest(reference),
        "metrics": {
            "visual_score": round(score, 6),
            "mae": round(mae, 6),
            "rmse": round(rmse, 6),
            "pixel_match_ratio": round(pixel_match, 6),
        },
        "original_sizes": {"rendered": list(render_image.size), "reference": list(reference_image.size)},
        "comparison_size": list(source.size),
        "resized_for_comparison": resized,
        "issues": issues,
        "human_visual_review_required": True,
    }


def apply_visual_thresholds(visual: dict[str, Any], thresholds: dict[str, Any]) -> dict[str, Any]:
    """Turn approved case thresholds into explicit visual gate evidence."""
    metrics = visual.get("metrics") if isinstance(visual.get("metrics"), dict) else {}
    issues = list(visual.get("issues") or [])
    for key, minimum in thresholds.items():
        if not key.startswith("min_"):
            continue
        metric = key[4:]
        observed = metrics.get(metric)
        if not isinstance(observed, (int, float)) or observed < minimum:
            issues.append({"code": f"visual_metric_below_{metric}", "observed": observed, "minimum": minimum})
    for key, maximum in thresholds.items():
        if not key.startswith("max_"):
            continue
        metric = key[4:]
        observed = metrics.get(metric)
        if not isinstance(observed, (int, float)) or observed > maximum:
            issues.append({"code": f"visual_metric_above_{metric}", "observed": observed, "maximum": maximum})
    return {**visual, "valid": not issues, "thresholds": thresholds, "issues": issues}


def pptx_observed(
    deck: Path, source_reference_sha256: str | None, process_sha256: str | None,
    reference_text_free: bool,
) -> dict[str, Any]:
    prs = Presentation(str(deck))
    slide_width, slide_height = int(prs.slide_width), int(prs.slide_height)
    shape_types: dict[str, int] = {}
    native_panel_names: list[str] = []
    native_text_objects = 0
    native_text_runs = 0
    picture_count = 0
    full_slide_picture_count = 0
    frame_picture_count = 0
    tables: dict[str, dict[str, Any]] = {}
    a_tbl_count = 0
    slide_shape_count = 0
    media_hashes: dict[str, str] = {}
    xml_slides: list[etree._Element] = []
    with ZipFile(deck) as zip_file:
        slide_names = sorted(
            (name for name in zip_file.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)),
            key=lambda name: int(re.search(r"(\d+)", name).group(1)),
        )
        for slide_name in slide_names:
            root = etree.fromstring(zip_file.read(slide_name))
            xml_slides.append(root)
            targets = relationship_targets(zip_file, slide_name)
            for picture in root.xpath(".//p:pic", namespaces=NS):
                picture_count += 1
                rel_id = image_relationship(picture)
                target = targets.get(rel_id or "")
                blob = zip_file.read(target) if target in zip_file.namelist() else b""
                media_hash = hashlib.sha256(blob).hexdigest() if blob else None
                if media_hash:
                    media_hashes[media_hash] = target
                width, height = shape_extent(picture)
                full_slide = width >= slide_width * 0.95 and height >= slide_height * 0.95
                full_slide_picture_count += int(full_slide)
                frame_picture_count += int(process_sha256 is not None and media_hash == process_sha256)
        for slide in prs.slides:
            slide_shape_count += len(slide.shapes)
            for shape in slide.shapes:
                type_name = str(shape.shape_type)
                shape_types[type_name] = shape_types.get(type_name, 0) + 1
                if getattr(shape, "has_text_frame", False):
                    native_text_objects += 1
                    native_text_runs += sum(len(paragraph.runs) for paragraph in shape.text_frame.paragraphs)
                if getattr(shape, "has_table", False):
                    tables[shape.name] = {
                        "rows": len(shape.table.rows),
                        "columns": len(shape.table.columns),
                        "values": [[str(shape.table.cell(r, c).text or "") for c in range(len(shape.table.columns))] for r in range(len(shape.table.rows))],
                    }
                is_panel = bool(PANEL_NAME.match(shape.name))
                is_native_panel = is_panel and str(shape.shape_type) in {"AUTO_SHAPE (1)", "GROUP (6)"}
                if is_native_panel:
                    native_panel_names.append(shape.name)
        for root in xml_slides:
            for table in root.xpath(".//a:tbl", namespaces=NS):
                name = shape_name(table.getparent().getparent().getparent()) if table.getparent() is not None else ""
                # The graphic frame is the nearest p:graphicFrame ancestor.
                ancestor = table
                while ancestor is not None and local_name(ancestor) != "graphicFrame":
                    ancestor = ancestor.getparent()
                name = shape_name(ancestor) if ancestor is not None else name
                tables.setdefault(name, table_record(table))
                a_tbl_count += 1
                tables[name].update(table_record(table))
    source_picture = any(media_hash == source_reference_sha256 for media_hash in media_hashes)
    # A permitted full-slide picture is acceptable only when it is an
    # explicitly text-free substrate.  If the source reference is a rendered
    # slide containing formal copy, reusing that image is a raster-text
    # failure; this case's source PPTX uses a text-free RGB background.
    formal_text_in_raster = bool(source_picture and full_slide_picture_count and not reference_text_free)
    native_table_shapes = sorted(name for name, record in tables.items() if record.get("rows", 0) > 0)
    table_values = {name: record for name, record in tables.items() if record.get("rows", 0) > 0}
    return {
        "slides": len(prs.slides),
        "shape_count": slide_shape_count,
        "shape_types": shape_types,
        "pictures": picture_count,
        "full_slide_pictures": full_slide_picture_count,
        "permitted_full_slide_pictures": full_slide_picture_count,
        "semantic_frame_pictures": frame_picture_count,
        "native_tables": a_tbl_count,
        "a_tbl_count": a_tbl_count,
        "native_table_shapes": native_table_shapes,
        "tables": table_values,
        "policy_fee_table": table_values.get("policy-fee-table"),
        "monthly_incentive_table": table_values.get("monthly-incentive-table"),
        "native_panel_count": len(sorted(set(native_panel_names))),
        "native_panel_names": sorted(set(native_panel_names)),
        "native_panel_groups": bool(native_panel_names),
        "native_text_objects": native_text_objects,
        "native_text_runs": native_text_runs,
        "formal_text_in_raster": formal_text_in_raster,
        "body_native_text": bool(native_text_runs and not formal_text_in_raster),
        "source_full_slide_picture_detected": source_picture,
        "media_hashes": sorted(media_hashes),
    }


def expected_matches(expected: Any, observed: Any, path: tuple[str, ...] = ()) -> tuple[bool, str]:
    if isinstance(expected, dict):
        if not isinstance(observed, dict):
            return False, f"{'.'.join(path)} expected object"
        for key, wanted in expected.items():
            if key not in observed:
                return False, f"missing {'.'.join(path + (str(key),))}"
            ok, detail = expected_matches(wanted, observed[key], path + (str(key),))
            if not ok:
                return False, detail
        return True, "expected structure matched"
    if isinstance(expected, list):
        if not isinstance(observed, list):
            return False, f"{'.'.join(path)} expected list"
        if path[-1:] == ("native_table_shapes",):
            missing = [item for item in expected if item not in observed]
            return not missing, "required table shapes present" if not missing else f"missing table shapes: {missing}"
        if path[-1:] == ("merged_cells",):
            normalize = lambda item: json.dumps(item, sort_keys=True)
            wanted = sorted(normalize(item) for item in expected)
            actual = sorted(normalize(item) for item in observed)
            return wanted == actual, "merged cells match" if wanted == actual else "merged cells differ"
        return expected == observed, "list matches" if expected == observed else f"{'.'.join(path)} differs"
    if path[-1:] in {("permitted_full_slide_pictures",), ("full_slide_pictures",)} and isinstance(expected, (int, float)):
        return isinstance(observed, (int, float)) and observed <= expected, f"{'.'.join(path)} observed={observed}, max={expected}"
    return expected == observed, f"{'.'.join(path)} observed={observed!r}, expected={expected!r}"


def mutation_smoke(deck: Path, work_dir: Path, render_dir: Path, original_render: Path | None) -> dict[str, Any]:
    prs = Presentation(str(deck))
    table_shape = next((shape for slide in prs.slides for shape in slide.shapes if getattr(shape, "has_table", False)), None)
    panel_shape = next((shape for slide in prs.slides for shape in slide.shapes if PANEL_NAME.match(shape.name) and str(shape.shape_type) == "AUTO_SHAPE (1)"), None)
    if table_shape is None or panel_shape is None:
        return {
            "valid": False,
            "table_cell_mutated": False,
            "panel_moved": False,
            "rendered_change": False,
            "reason": "native table or native panel is missing",
        }
    before_text = str(table_shape.table.cell(0, 0).text or "")
    before_left = int(panel_shape.left)
    mutated = work_dir / "mutation-smoke.pptx"
    table_shape.table.cell(0, 0).text = before_text + " [mutation]"
    panel_shape.left = before_left + Inches(0.05)
    prs.save(mutated)
    reopened = Presentation(str(mutated))
    mutated_table = next(shape for slide in reopened.slides for shape in slide.shapes if getattr(shape, "has_table", False))
    mutated_panel = next(shape for slide in reopened.slides for shape in slide.shapes if shape.name == panel_shape.name)
    cell_changed = "[mutation]" in str(mutated_table.table.cell(0, 0).text or "")
    panel_changed = int(mutated_panel.left) != before_left
    rendered_changed = False
    ratio = 0.0
    if original_render is not None:
        ok, _ = render_deck(mutated, render_dir / "mutation-render")
        changed_render = render_dir / "mutation-render" / "slide-1.png"
        if ok and changed_render.is_file():
            with Image.open(original_render) as before_image, Image.open(changed_render) as after_image:
                after = after_image.convert("RGB").resize(before_image.convert("RGB").size, Image.Resampling.LANCZOS)
                diff = ImageChops.difference(before_image.convert("RGB"), after).convert("L")
                histogram = diff.histogram()
                total = diff.width * diff.height
                ratio = 1.0 - (histogram[0] / total if total else 1.0)
                rendered_changed = ratio > 0.0001
    return {
        "valid": bool(cell_changed and panel_changed and rendered_changed),
        "table_cell_mutated": cell_changed,
        "panel_moved": panel_changed,
        "rendered_change": rendered_changed,
        "rendered_change_ratio": round(ratio, 6),
        "mutated_deck": str(mutated.resolve()),
    }


def replay(args: argparse.Namespace) -> dict[str, Any]:
    deck = Path(args.deck).resolve()
    source = Path(args.source_pptx).resolve()
    case_spec = load_json(Path(args.case_spec).resolve())
    if not isinstance(case_spec, dict):
        raise ValueError("case specification must be an object")
    output = Path(args.output).resolve()
    output_dir = Path(args.output_dir).resolve() if args.output_dir else output.parent / f"{args.phase}-case-artifacts"
    output_dir.mkdir(parents=True, exist_ok=True)
    process = Path(args.process_image).resolve() if args.process_image else None
    reference, source_hashes = source_reference(source, output_dir, Path(args.reference_image).resolve() if args.reference_image else None)
    if process is not None:
        source_hashes["process_image_sha256"] = digest(process)
    reference_text_free = bool(
        case_spec.get("reference_text_free")
        or (case_spec.get("source_policy") or {}).get("reference_text_free")
    )
    observed = pptx_observed(
        deck,
        source_hashes.get("source_reference_sha256"),
        source_hashes.get("process_image_sha256"),
        reference_text_free,
    )
    render_dir = output_dir / "rendered"
    rendered_ok, render_error = render_deck(deck, render_dir)
    rendered = render_dir / "slide-1.png"
    visual_reference = reference
    visual_reference_kind = "source-reference"
    if args.phase == "candidate" and args.baseline_evaluation:
        try:
            baseline_data = load_json(Path(args.baseline_evaluation).resolve())
            baseline_files = ((baseline_data.get("rendered") or {}).get("files") or []) if isinstance(baseline_data, dict) else []
            if baseline_files and Path(str(baseline_files[0].get("path"))).is_file():
                visual_reference = Path(str(baseline_files[0]["path"])).resolve()
                visual_reference_kind = "baseline-render"
        except (OSError, ValueError, json.JSONDecodeError, AttributeError, TypeError):
            pass
    visual = visual_compare(rendered, visual_reference) if rendered_ok and rendered.is_file() else {
        "valid": False, "metrics": {}, "issues": [{"code": "render_failed", "message": render_error or "slide-1.png missing"}],
    }
    if args.phase == "baseline" and not args.reference_image and rendered_ok and rendered.is_file():
        # The original deck render is the visual baseline for later candidate
        # comparisons.  The extracted RGB image remains separately recorded
        # as the source/background hash evidence.
        visual = visual_compare(rendered, rendered)
        visual_reference_kind = "baseline-render-self"
    visual["reference_kind"] = visual_reference_kind
    visual = apply_visual_thresholds(visual, (case_spec.get("quality_thresholds") or {}).get("visual", {}))
    object_ok, object_detail = expected_matches(case_spec.get("expected") or {}, observed)
    tables = observed.get("tables") or {}
    table_audit = {
        "valid": object_ok and observed.get("a_tbl_count") == observed.get("native_tables"),
        "a_tbl_count": observed.get("a_tbl_count"),
        "native_table_count": observed.get("native_tables"),
        "required_shapes": case_spec.get("expected", {}).get("native_table_shapes", []),
        "observed_shapes": observed.get("native_table_shapes", []),
        "policy_fee_table": tables.get("policy-fee-table"),
        "monthly_incentive_table": tables.get("monthly-incentive-table"),
    }
    panel_audit = {
        "valid": bool(observed.get("native_panel_groups")) and observed.get("native_panel_count", 0) > 0,
        "native_panel_count": observed.get("native_panel_count"),
        "native_panel_names": observed.get("native_panel_names", []),
        "semantic_frame_pictures": observed.get("semantic_frame_pictures", 0),
    }
    text_audit = {
        "valid": bool(observed.get("body_native_text")) and observed.get("formal_text_in_raster") is False,
        "native_text_objects": observed.get("native_text_objects"),
        "native_text_runs": observed.get("native_text_runs"),
        "formal_text_in_raster": observed.get("formal_text_in_raster"),
        "body_native_text": observed.get("body_native_text"),
    }
    mutation = mutation_smoke(deck, output_dir, output_dir, rendered if rendered_ok else None)
    object_comparison: dict[str, Any] = {"valid": object_ok, "detail": object_detail}
    if args.baseline_evaluation:
        baseline = load_json(Path(args.baseline_evaluation).resolve())
        before = baseline.get("observed", {}) if isinstance(baseline, dict) else {}
        object_comparison.update({
            "baseline_evaluation": str(Path(args.baseline_evaluation).resolve()),
            "delta": {
                "native_tables": observed.get("native_tables", 0) - before.get("native_tables", 0),
                "a_tbl_count": observed.get("a_tbl_count", 0) - before.get("a_tbl_count", 0),
                "native_panel_count": observed.get("native_panel_count", 0) - before.get("native_panel_count", 0),
                "native_text_runs": observed.get("native_text_runs", 0) - before.get("native_text_runs", 0),
            },
        })
    failures: list[str] = []
    if not object_ok:
        failures.append("case-expected-structure")
    if not table_audit["valid"]:
        failures.append("native-table-xml")
    if not panel_audit["valid"]:
        failures.append("native-panel-object")
    if not text_audit["valid"]:
        failures.append("text-content-and-run-style")
    if not mutation["valid"]:
        failures.append("mutation-smoke-test")
    if not visual.get("valid"):
        failures.append("visual-comparison")
    metrics = dict(visual.get("metrics") or {})
    metrics.update({
        "native_tables": observed.get("native_tables", 0),
        "native_panel_count": observed.get("native_panel_count", 0),
        "formal_text_in_raster": bool(observed.get("formal_text_in_raster")),
        "whole_slide_picture_count": observed.get("full_slide_pictures", 0),
        "failed_gate_count": len(failures),
    })
    return {
        "schema": SCHEMA,
        "evaluation_schema": "ai-ppt-plus/distillation-evaluation/v1",
        "generated_at": utc_now(),
        "phase": args.phase,
        "case_id": case_spec.get("case_id"),
        "status": "passed" if not failures else "failed",
        "valid": not failures,
        "failure_codes": failures,
        "source_hashes": source_hashes,
        "deck": {"path": str(deck), "sha256": digest(deck), "slides": observed.get("slides")},
        "rendered": {
            "directory": str(render_dir),
            "files": [{"path": str(path), "sha256": digest(path)} for path in sorted(render_dir.glob("slide-*.png"))],
            "valid": rendered_ok,
            "error": render_error,
        },
        "observed": observed,
        "native_object_audit": {"valid": object_ok, "detail": object_detail, "observed": observed},
        "table_audit": table_audit,
        "panel_audit": panel_audit,
        "text_audit": text_audit,
        "mutation_smoke_test": mutation,
        "visual_comparison": visual,
        "object_comparison": object_comparison,
        "metrics": metrics,
        "behavioral_change": args.phase == "candidate",
        "repair_fingerprint": args.repair_fingerprint if args.phase == "candidate" else None,
        "regressions": [{
            "case_id": case_spec.get("case_id"),
            "status": "passed" if not failures else "failed",
            "valid": not failures,
        }],
        "human_visual_review_required": True,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--phase", choices=("baseline", "candidate"), required=True)
    parser.add_argument("--case-spec", required=True)
    parser.add_argument("--deck", required=True)
    parser.add_argument("--source-pptx", required=True)
    parser.add_argument("--process-image")
    parser.add_argument("--reference-image")
    parser.add_argument("--baseline-evaluation")
    parser.add_argument("--repair-fingerprint")
    parser.add_argument("--output-dir")
    parser.add_argument("--output", required=True)
    args = parser.parse_args()
    try:
        result = replay(args)
    except (OSError, ValueError, json.JSONDecodeError, KeyError, TypeError) as exc:
        result = {
            "schema": SCHEMA,
            "generated_at": utc_now(),
            "phase": args.phase,
            "case_id": None,
            "status": "blocked",
            "valid": False,
            "failure_codes": ["case-replay-runtime"],
            "reasons": [f"{type(exc).__name__}: {exc}"],
        }
    write_json(Path(args.output).resolve(), result)
    print(json.dumps(result, ensure_ascii=False))
    return 0 if result.get("valid") is True else 2


if __name__ == "__main__":
    raise SystemExit(main())
