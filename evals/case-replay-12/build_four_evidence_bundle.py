#!/usr/bin/env python3
"""Build self-contained four-evidence bundles for reconstruction replay cases."""
from __future__ import annotations

import argparse
import hashlib
import json
import shutil
from pathlib import Path

from PIL import Image, ImageChops
from pptx import Presentation

REPO = Path(__file__).resolve().parents[2]


def digest(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def resolve_path(value: str, replay_root: Path) -> Path:
    candidate = Path(value)
    if candidate.is_absolute():
        return candidate
    for base in (REPO, replay_root, Path(__file__).resolve().parent):
        path = (base / candidate).resolve()
        if path.exists():
            return path
    return (REPO / candidate).resolve()


def shape_record(shape) -> dict:
    value = {
        "shape_id": int(shape.shape_id),
        "name": str(shape.name),
        "shape_type": str(shape.shape_type),
        "x": int(shape.left),
        "y": int(shape.top),
        "w": int(shape.width),
        "h": int(shape.height),
        "has_text_frame": bool(getattr(shape, "has_text_frame", False)),
        "has_table": bool(getattr(shape, "has_table", False)),
        "has_chart": bool(getattr(shape, "has_chart", False)),
    }
    if value["has_text_frame"]:
        value["text"] = str(shape.text)
    if value["has_table"]:
        value["table"] = [[cell.text for cell in row.cells] for row in shape.table.rows]
    if value["has_chart"]:
        value["chart_type"] = str(shape.chart.chart_type)
    if hasattr(shape, "shapes"):
        value["children"] = [shape_record(child) for child in shape.shapes]
    return value


def object_tree(deck: Path) -> dict:
    prs = Presentation(str(deck))
    return {
        "schema": "ai-ppt-plus/pptx-object-tree/v1",
        "deck_sha256": digest(deck),
        "slide_count": len(prs.slides),
        "slides": [
            {"slide_no": index, "objects": [shape_record(shape) for shape in slide.shapes]}
            for index, slide in enumerate(prs.slides, start=1)
        ],
    }


def save_diff(reference: Path, rendered: Path, output: Path) -> None:
    with Image.open(reference).convert("RGB") as ref, Image.open(rendered).convert("RGB") as out:
        if out.size != ref.size:
            out = out.resize(ref.size)
        ImageChops.difference(ref, out).save(output)


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--candidate-evaluation", default=str(Path(__file__).resolve().parent / "candidate-evaluation.json"))
    parser.add_argument("--replay-root", required=True)
    parser.add_argument("--output", required=True)
    parser.add_argument("--strict", action="store_true")
    args = parser.parse_args()

    evaluation = json.loads(Path(args.candidate_evaluation).read_text(encoding="utf-8"))
    replay_root = Path(args.replay_root).resolve()
    manifest_cases = []
    failures = []

    for case in evaluation.get("cases", []):
        case_id = str(case.get("case_id") or "")
        candidate = case.get("candidate") or {}
        reference = resolve_path(str(candidate.get("reference") or ""), replay_root)
        rendered = resolve_path(str(candidate.get("rendered") or ""), replay_root)
        deck = resolve_path(str(candidate.get("deck") or ""), replay_root)
        bundle = replay_root / "four-evidence" / case_id
        bundle.mkdir(parents=True, exist_ok=True)
        required = {"reference": reference, "rendered": rendered, "deck": deck}
        missing = [name for name, path in required.items() if not path.is_file()]
        if missing:
            failures.append({"case_id": case_id, "missing": missing})
            continue

        reference_copy = bundle / "reference.png"
        rendered_copy = bundle / "rendered.png"
        diff_path = bundle / "visual-diff.png"
        object_tree_path = bundle / "object-tree.json"
        shutil.copy2(reference, reference_copy)
        shutil.copy2(rendered, rendered_copy)
        save_diff(reference_copy, rendered_copy, diff_path)
        object_tree_path.write_text(json.dumps(object_tree(deck), ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

        evidence = {
            "schema": "ai-ppt-plus/reconstruction-four-evidence/v1",
            "case_id": case_id,
            "deck": str(deck),
            "deck_sha256": digest(deck),
            "reference": str(reference_copy),
            "reference_sha256": digest(reference_copy),
            "rendered": str(rendered_copy),
            "rendered_sha256": digest(rendered_copy),
            "visual_diff": str(diff_path),
            "visual_diff_sha256": digest(diff_path),
            "object_tree": str(object_tree_path),
            "object_tree_sha256": digest(object_tree_path),
            "technical_status": candidate.get("technical_status"),
            "human_visual_review_required": True,
            "release_eligible": False,
        }
        evidence_path = bundle / "evidence.json"
        evidence_path.write_text(json.dumps(evidence, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        manifest_cases.append({"case_id": case_id, "evidence": str(evidence_path), "evidence_sha256": digest(evidence_path)})

    result = {
        "schema": "ai-ppt-plus/reconstruction-four-evidence-suite/v1",
        "case_count": len(evaluation.get("cases", [])),
        "evidence_count": len(manifest_cases),
        "cases": manifest_cases,
        "failures": failures,
        "valid": not failures and len(manifest_cases) == len(evaluation.get("cases", [])),
        "human_visual_review_required": True,
        "release_eligible": False,
    }
    output = Path(args.output)
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(json.dumps(result, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(json.dumps(result, ensure_ascii=False))
    if args.strict and not result["valid"]:
        raise SystemExit("four-evidence suite incomplete")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
