#!/usr/bin/env python3
"""Run baseline and native case replays for the 12 distillation scenes.

The reference PNGs are the visual authority.  The approved case-suite JSON is
the formal-copy/data authority.  Baseline is an explicit legacy image-only
control; candidate-before is a native reconstruction before the merge-topology
gate; candidate is the optimized native reconstruction after that gate.
"""
from __future__ import annotations

import argparse
import copy
import hashlib
import json
import os
import shutil
import subprocess
import sys
from pathlib import Path


ROOT = Path(__file__).resolve().parent
REPO = ROOT.parents[1]
EDITABLE = REPO / "ai-ppt-editable"
COMPOSE = EDITABLE / "scripts" / "compose_pptx.py"
RENDER = EDITABLE / "scripts" / "render_pptx.py"
INSPECT = EDITABLE / "scripts" / "inspect_pptx.py"
NATIVE_VALIDATE = EDITABLE / "scripts" / "validate_native_editability.py"
SEMANTIC_AUDIT = EDITABLE / "scripts" / "semantic_object_audit.py"
CASE_AUDIT = EDITABLE / "scripts" / "case_replay_audit.py"
MANIFEST_BUILDER = EDITABLE / "scripts" / "build_object_manifest.py"
FONT_DIR = EDITABLE / "assets" / "fonts"

NAVY = "#061A35"
NAVY_2 = "#0C2B4D"
BLUE = "#1687FF"
RED = "#E60012"
RED_DARK = "#9B0B1B"
SILVER = "#F4F7FB"
MUTED = "#A9BED5"
GRID = "#355A7D"
GREEN = "#55D6A1"
AMBER = "#F3B849"


def digest(path: Path) -> str:
    value = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            value.update(chunk)
    return value.hexdigest()


def repo_relative(path: Path) -> str:
    """Return a path that remains usable from the repository root."""
    return str(path.resolve().relative_to(REPO))


def case_relative(path: Path) -> str:
    """Return a portable evidence path relative to this case package."""
    return str(path.resolve().relative_to(ROOT))


def write_json(path: Path, value):
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(value, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def run(command, cwd=REPO, allow_failure=False):
    result = subprocess.run([str(item) for item in command], cwd=cwd, capture_output=True, text=True, check=False)
    if result.returncode and not allow_failure:
        raise RuntimeError(f"command failed ({result.returncode}): {' '.join(str(item) for item in command)}\n{result.stdout}\n{result.stderr}")
    return result


def shape(object_id, shape_type, x, y, w, h, fill=None, line=None, line_width=0.8, **extra):
    value = {"object_id": object_id, "type": shape_type, "x": x, "y": y, "w": w, "h": h}
    if fill is not None:
        value["fill"] = fill
    if line is not None:
        value["line"] = line
        value["line_width"] = line_width
    value.update(extra)
    return value


def line(object_id, x, y, x2, y2, color=GRID, width=1.0):
    return {"object_id": object_id, "type": "line", "x": x, "y": y, "x2": x2, "y2": y2, "line": color, "line_width": width}


def text(object_id, value, x, y, w, h, size=14, color=SILVER, bold=False, align="left", valign="top", runs=None, **extra):
    item = {"object_id": object_id, "text": str(value), "x": x, "y": y, "w": w, "h": h, "size": size, "color": color, "bold": bold, "align": align, "valign": valign, "margin_left": 0.01, "margin_right": 0.01, "margin_top": 0.005, "margin_bottom": 0.005}
    if runs:
        item["runs"] = runs
        item.pop("text", None)
    item.update(extra)
    return item


def group(object_id, x, y, w, h, children, *, role="semantic-panel", native_required=True):
    return {"object_id": object_id, "role": role, "native_required": native_required, "x": x, "y": y, "w": w, "h": h, "children_coordinate_space": "local", "children": children, "alt_text": object_id}


def panel(object_id, x, y, w, h, accent=RED, fill=NAVY_2):
    return group(object_id, x, y, w, h, [
        shape(f"{object_id}-fill", "rounded_rect", 0, 0, 1, 1, fill=fill, line=GRID, line_width=0.8),
        shape(f"{object_id}-accent", "rect", 0, 0, 1, 0.055, fill=accent),
        shape(f"{object_id}-rule", "line", 0.03, 0.90, 0.94, 0.001, line=GRID, line_width=0.6),
    ])


def add_title(slide, title, subtitle=None):
    slide["shapes"].extend([
        shape("title-accent", "rect", 0.055, 0.065, 0.006, 0.075, RED),
        line("title-rule", 0.055, 0.17, 0.945, 0.17, GRID, 0.7),
    ])
    slide["texts"].append(text("title", title, 0.073, 0.052, 0.66, 0.09, size=25, bold=True))
    if subtitle:
        slide["texts"].append(text("subtitle", subtitle, 0.075, 0.145, 0.62, 0.035, size=8.5, color=MUTED))
    slide["texts"].append(text("case-label", "CASE REPLAY / NATIVE EDITABLE", 0.72, 0.075, 0.22, 0.04, size=7.5, color=BLUE, bold=True, align="right"))


def add_kpi(slide, object_id, x, y, w, label, value, unit, color=RED):
    slide["groups"].append(panel(object_id, x, y, w, 0.15, color, "#0B2747"))
    slide["texts"].extend([
        text(f"{object_id}-label", label, x + 0.025, y + 0.03, w * 0.42, 0.04, size=9, color=MUTED),
        text(f"{object_id}-value", value, x + w * 0.46, y + 0.018, w * 0.32, 0.065, size=20, color=SILVER, bold=True, align="right"),
        text(f"{object_id}-unit", unit, x + w * 0.79, y + 0.045, w * 0.17, 0.035, size=8, color=color, bold=True, align="right"),
    ])


def add_node(slide, object_id, x, y, w, label, index, color=BLUE):
    slide["groups"].append(panel(object_id, x, y, w, 0.105, color, "#0D2F52"))
    slide["texts"].extend([
        text(f"{object_id}-index", f"{index:02d}", x + 0.018, y + 0.028, 0.05, 0.04, size=10, color=color, bold=True),
        text(f"{object_id}-label", label, x + 0.075, y + 0.027, w - 0.09, 0.045, size=10.5, bold=True),
    ])


def add_table(slide, object_id, x, y, w, h, rows, *, merges=None, rich=False, header_fill=RED_DARK, body_fill="#F5F8FC", column_widths=None, row_heights=None):
    cell_styles = {"0": {"fill": header_fill, "color": "#FFFFFF", "bold": True, "size": 9}}
    for row_index in range(1, len(rows)):
        cell_styles[str(row_index)] = {"fill": body_fill if row_index % 2 else "#EAF1F8", "color": NAVY, "size": 8.5}
    if rich:
        for row_index, row in enumerate(rows):
            for column_index, value in enumerate(row):
                if isinstance(value, dict):
                    cell_styles[f"{row_index},{column_index}"] = {"fill": body_fill, "color": NAVY, "size": 8.5}
    spec = {"object_id": object_id, "native_required": True, "x": x, "y": y, "w": w, "h": h, "rows": rows, "columns": max(len(row) for row in rows), "header_fill": header_fill, "fill": body_fill, "border": {"all": {"color": "#9EB3C8", "width": 0.7}}, "cell_margins": {"left": 0.04, "right": 0.04, "top": 0.025, "bottom": 0.025}, "cell_styles": cell_styles, "rich_text_required": bool(rich)}
    if merges:
        spec["merges"] = merges
    if column_widths:
        spec["column_widths"] = column_widths
    if row_heights:
        spec["row_heights"] = row_heights
    slide["tables"].append(spec)


def make_base(case, run_dir):
    return {
        "project_id": "distillation-case-replay-12",
        "slide_width_in": 13.333,
        "slide_height_in": 7.5,
        "ref_width": 1920,
        "ref_height": 1080,
        "units": "fraction",
        "assets_dir": repo_relative(run_dir),
        "font_family": "Noto Sans CJK SC",
        "editable_object_policy": "native-semantic-objects",
        "theme": {"font": "Noto Sans CJK SC", "text_color": SILVER, "size": 10, "table_header_fill": RED_DARK, "table_fill": "#F5F8FC", "chart_colors": [RED, BLUE], "chart_text_color": SILVER, "chart_muted_color": MUTED, "chart_grid_color": GRID, "chart_tick_size": 8, "chart_legend_size": 8, "chart_label_size": 7},
        "slides": [{"layout_name": "Blank", "shapes": [shape("background", "rect", 0, 0, 1, 1, fill=NAVY)], "groups": [], "tables": [], "charts": [], "texts": [], "icons": []}],
    }


def framework_layout(case, run_dir, optimized):
    deck = make_base(case, run_dir)
    slide = deck["slides"][0]
    add_title(slide, case["title"], "战略牵引 × 能力筑基 × 场景突破 × 价值兑现")
    slide["groups"].extend([
        panel("layer-foundation", 0.07, 0.31, 0.86, 0.30, BLUE, "#0A2544"),
        panel("layer-value", 0.16, 0.24, 0.68, 0.17, RED, "#102F4F"),
        panel("result-cabin", 0.36, 0.71, 0.28, 0.12, RED, "#5C1120"),
    ])
    labels = ["战略牵引", "能力筑基", "场景突破", "价值兑现"]
    for index, label in enumerate(labels):
        x = 0.12 + index * 0.20
        slide["groups"].append(panel(f"card-{index + 1}", x, 0.28, 0.16, 0.14, RED if index == 0 else BLUE, "#123758"))
        slide["texts"].append(text(f"card-{index + 1}-text", label, x + 0.01, 0.325, 0.14, 0.04, size=9.5, bold=True, align="center"))
    for index, label in enumerate(["组织", "技术", "数据", "生态"]):
        x = 0.13 + index * 0.20
        slide["texts"].append(text(f"foundation-{index + 1}", label, x, 0.48, 0.16, 0.045, size=10, color=BLUE, bold=True, align="center"))
    for index in range(3):
        slide["shapes"].append(line(f"framework-arrow-{index + 1}", 0.27 + index * 0.20, 0.36, 0.31 + index * 0.20, 0.36, RED, 2.2))
    slide["texts"].append(text("result-cabin-text", "经营结果", 0.38, 0.747, 0.24, 0.04, size=14, bold=True, align="center"))
    slide["texts"].extend([
        text("framework-footnote", "组织、技术、数据、生态共同支撑四个经营跃迁层", 0.09, 0.87, 0.58, 0.035, size=8.5, color=MUTED),
        text("framework-score", "可移动卡片 / 3层结构 / 连接线", 0.70, 0.87, 0.22, 0.035, size=8, color=BLUE, align="right"),
    ])
    return deck


def dashboard_layout(case, run_dir, optimized):
    deck = make_base(case, run_dir)
    slide = deck["slides"][0]
    add_title(slide, case["title"], "以目标—路径—结果组织经营管理阅读路径")
    add_kpi(slide, "kpi-target", 0.07, 0.21, 0.25, "目标", "28.6", "亿元", RED)
    add_kpi(slide, "kpi-path", 0.375, 0.21, 0.25, "路径", "18.7", "亿元", BLUE)
    add_kpi(slide, "kpi-result", 0.68, 0.21, 0.25, "结果", "17.3", "亿元", GREEN)
    slide["texts"].append(text("flow-label", "增长路径 / FIVE-STAGE OPERATING LOOP", 0.07, 0.405, 0.4, 0.035, size=8, color=BLUE, bold=True))
    labels = ["目标洞察", "线索获取", "商机转化", "合同签约", "交付增值"]
    for index, label in enumerate(labels):
        add_node(slide, f"flow-{index + 1}", 0.07 + index * 0.17, 0.45, 0.145, label, index + 1, RED if index == 2 else BLUE)
        if index < 4:
            slide["shapes"].append(line(f"flow-link-{index + 1}", 0.216 + index * 0.17, 0.502, 0.238 + index * 0.17, 0.502, RED, 1.5))
    slide["groups"].append(panel("insight-panel", 0.84, 0.39, 0.09, 0.23, GREEN, "#103B4B"))
    slide["texts"].append(text("insight-panel-text", "TOP5\n进展", 0.85, 0.445, 0.07, 0.09, size=11, color=GREEN, bold=True, align="center"))
    table_rows = [["行业", "增长结构", "状态", "口径", "证据"]] + [[label, "TOP5", "已绑定", "亿元", "可回放"] for label in ["政府", "教育", "医疗", "交通", "制造"]]
    add_table(slide, "industry-growth-table", 0.07, 0.59, 0.42, 0.28, table_rows, column_widths=[0.08, 0.13, 0.09, 0.07, 0.05], row_heights=[0.04] + [0.048] * 5)
    project_rows = [["项目", "阶段", "风险", "进度", "证据"]] + [[f"重点项目 {index}", "经营", "可审计", "跟踪", "TOP5"] for index in range(1, 6)]
    add_table(slide, "project-progress-table", 0.52, 0.59, 0.41, 0.28, project_rows, column_widths=[0.12, 0.09, 0.09, 0.08, 0.05], row_heights=[0.04] + [0.048] * 5)
    slide["texts"].append(text("dashboard-note", "原生KPI卡 / 原生流程节点 / 2张原生表格", 0.07, 0.91, 0.5, 0.03, size=7.5, color=MUTED))
    return deck


def basic_table_layout(case, run_dir, optimized):
    deck = make_base(case, run_dir)
    slide = deck["slides"][0]
    add_title(slide, case["title"], "统一标准、清晰计提、结果可核对")
    slide["groups"].append(panel("policy-panel", 0.07, 0.22, 0.18, 0.58, RED, "#0B2747"))
    slide["texts"].extend([
        text("policy-panel-title", "POLICY\nMATRIX", 0.09, 0.31, 0.14, 0.10, size=15, color=RED, bold=True),
        text("policy-panel-copy", "场景\n收入\n标准\n周期", 0.09, 0.49, 0.12, 0.17, size=11, color=SILVER, bold=True),
    ])
    rows = [
        ["业务场景", "收入变化", "服务费标准", "计提周期"],
        ["存量升档", "可核对", "统一标准", "月度"],
        ["指定套餐", "可核对", "统一标准", "月度"],
        ["新装发展", "可核对", "统一标准", "月度"],
        ["渠道激励", "可核对", "统一标准", "月度"],
        ["增收场景", "结果留痕", "按政策", "周期锁定"],
        ["复核口径", "来源绑定", "可追溯", "可回放"],
    ]
    add_table(slide, "service-fee-table", 0.28, 0.22, 0.65, 0.58, rows, column_widths=[0.16, 0.16, 0.18, 0.15], row_heights=[0.075] + [0.084] * 6)
    slide["texts"].extend([
        text("table-callout", "TABLE / NATIVE", 0.73, 0.26, 0.15, 0.03, size=7.5, color=BLUE, bold=True, align="right"),
        text("table-footnote", "行高、列宽、网格、单元格文字均保留为可编辑语义对象", 0.28, 0.84, 0.65, 0.035, size=8.5, color=MUTED, align="right"),
    ])
    return deck


def merge_table_layout(case, run_dir, optimized):
    deck = make_base(case, run_dir)
    slide = deck["slides"][0]
    add_title(slide, case["title"], "增收有奖｜减收不罚")
    # Keep the evidence badge in the title band without colliding with the
    # routing label.  This is a case-level visual fix; it does not change the
    # native table contract below.
    case_label = next(item for item in slide["texts"] if item.get("object_id") == "case-label")
    case_label.update({"x": 0.53, "w": 0.18, "align": "left"})
    slide["groups"].append(panel("principle-badge", 0.76, 0.055, 0.17, 0.11, RED, "#5C1120"))
    slide["texts"].append(text("principle-badge-text", "增收有奖\n减收不罚", 0.775, 0.078, 0.14, 0.065, size=10, bold=True, align="center"))
    for index, label in enumerate(["核心原则", "场景", "结果"]):
        x = 0.07 + index * 0.29
        slide["groups"].append(panel(f"commission-card-{index + 1}", x, 0.22, 0.25, 0.17, RED if index == 0 else BLUE, "#0D2F52"))
        slide["texts"].extend([
            text(f"commission-card-{index + 1}-title", label, x + 0.02, 0.25, 0.21, 0.035, size=10, color=RED if index == 0 else BLUE, bold=True),
            text(f"commission-card-{index + 1}-copy", ["增收有奖", "达量补贴标准", "结果可核对"][index], x + 0.02, 0.31, 0.21, 0.04, size=9),
        ])
    slide["shapes"].append(shape("section-divider", "rect", 0.07, 0.43, 0.86, 0.025, fill=RED))
    rows = [
        ["场景", "收入变化", "服务费标准", "计提周期"],
        ["月度发展", "增收有奖", {"runs": [{"text": "达量补贴标准", "bold": True, "color": RED}, {"text": " / 已绑定", "color": NAVY}]}, "月度"],
        ["", "减收不罚", "按政策", "月度"],
        ["适用渠道", "增收有奖", {"runs": [{"text": "结果", "bold": True, "color": RED}, {"text": "可核对", "color": NAVY}]}, "周期锁定"],
        ["", "减收不罚", "统一标准", "可回放"],
        ["场景", "增收有奖", "来源绑定", "月度"],
        ["", "减收不罚", "可追溯", "月度"],
        ["结果", "可核对", "按政策", "结果"],
        ["", "可核对", "可追溯", "复核"],
    ]
    merges = case.get("data", {}).get("merge_spans", []) if optimized else []
    add_table(slide, "policy-merged-table", 0.07, 0.50, 0.51, 0.36, rows, merges=merges, rich=True, column_widths=[0.12, 0.12, 0.17, 0.10], row_heights=[0.04] + [0.04] * 8)
    incentive_rows = [["月度发展", "达量补贴标准", "适用渠道"], ["增收有奖", "结果可核对", "渠道"], ["减收不罚", "统一标准", "渠道"], ["结果", "可回放", "全部"]]
    add_table(slide, "monthly-incentive-table", 0.61, 0.50, 0.32, 0.18, incentive_rows, column_widths=[0.10, 0.14, 0.08], row_heights=[0.04] * 4)
    summary_rows = [["指标", "状态"], ["合并单元格", "已绑定"], ["富文本", "已绑定"]]
    add_table(slide, "policy-summary-table", 0.61, 0.72, 0.32, 0.14, summary_rows, column_widths=[0.15, 0.12], row_heights=[0.04] * 3)
    slide["texts"].append(text("merge-note", "3 native tables / merged first column / rich-text runs", 0.07, 0.90, 0.86, 0.03, size=7.5, color=MUTED, align="right"))
    return deck


def text_visual_layout(case, run_dir, optimized):
    deck = make_base(case, run_dir)
    slide = deck["slides"][0]
    add_title(slide, case["title"], "Service Excellence")
    slide["groups"].append(panel("service-principle", 0.07, 0.22, 0.25, 0.57, RED, "#102F4F"))
    for index, label in enumerate(["一线响应", "专业交付", "客户成功"]):
        slide["texts"].append(text(f"service-label-{index + 1}", label, 0.10, 0.31 + index * 0.12, 0.18, 0.05, size=13, bold=True))
        slide["shapes"].append(line(f"service-line-{index + 1}", 0.10, 0.38 + index * 0.12, 0.27, 0.38 + index * 0.12, RED if index == 0 else BLUE, 1.8))
    slide["groups"].append(panel("journey-panel", 0.36, 0.22, 0.34, 0.57, BLUE, "#0B2747"))
    journey = ["触达", "诊断", "交付", "复盘"]
    for index, label in enumerate(journey):
        x = 0.40 + index * 0.072
        slide["groups"].append(panel(f"journey-{index + 1}", x, 0.45, 0.06, 0.10, RED if index == 2 else BLUE, "#123758"))
        slide["texts"].append(text(f"journey-{index + 1}-text", label, x - 0.01, 0.58, 0.08, 0.04, size=8.5, bold=True, align="center"))
        if index < 3:
            slide["shapes"].append(line(f"journey-link-{index + 1}", x + 0.06, 0.50, x + 0.073, 0.50, RED, 1.5))
    slide["groups"].append(panel("closure-ring", 0.75, 0.27, 0.18, 0.28, GREEN, "#103B4B"))
    slide["shapes"].append(shape("closure-ring-circle", "oval", 0.79, 0.33, 0.20, 0.20, fill=GREEN, line=GREEN, line_width=1.5))
    slide["texts"].append(text("closure-rate", "96.8%", 0.78, 0.395, 0.22, 0.06, size=17, color=NAVY, bold=True, align="center"))
    for index, label in enumerate(["本周重点", "关键动作", "责任到人"]):
        x = 0.36 + index * 0.20
        slide["groups"].append(panel(f"responsibility-{index + 1}", x, 0.84, 0.17, 0.08, BLUE, "#0D2F52"))
        slide["texts"].append(text(f"responsibility-{index + 1}-text", label, x, 0.86, 0.17, 0.035, size=8.5, bold=True, align="center"))
    slide["texts"].append(text("bilingual-note", "CJK + English / run-level emphasis / native text", 0.72, 0.90, 0.22, 0.03, size=7.5, color=MUTED, align="right"))
    # One mixed-style text object makes the run-level color contract observable.
    slide["texts"].append(text("mixed-run-callout", "闭环率 96.8%", 0.75, 0.61, 0.18, 0.05, size=11, runs=[{"text": "闭环率 ", "color": SILVER}, {"text": "96.8%", "color": RED, "bold": True}]))
    return deck


def chart_layout(case, run_dir, optimized):
    deck = make_base(case, run_dir)
    slide = deck["slides"][0]
    # Prefer the formal copy captured by the case rather than a shortened
    # display title so the editable deck and the source contract agree.
    chart_title = (case.get("formal_text") or [case["title"]])[0]
    add_title(slide, chart_title, "政府与教育仍是规模增长主引擎")
    slide["groups"].append(panel("chart-panel", 0.07, 0.22, 0.58, 0.58, BLUE, "#0B2747"))
    data = case["data"]
    slide["charts"].append({"object_id": "industry-growth-chart", "type": "bar", "x": 0.11, "y": 0.30, "w": 0.50, "h": 0.42, "categories": data["categories"], "series": [{"name": "目标", "values": data["target"]}, {"name": "实际", "values": data["actual"]}], "colors": [RED, BLUE], "legend": True, "data_labels": True})
    slide["groups"].append(panel("chart-insight-panel", 0.69, 0.22, 0.24, 0.58, RED, "#102F4F"))
    for index, label in enumerate(["政府 8.6 / 5.2", "教育 5.2 / 3.1", "医疗 4.6 / 2.6", "交通 3.8 / 2.2", "制造 3.4 / 1.9"]):
        slide["texts"].append(text(f"chart-insight-{index + 1}", label, 0.72, 0.31 + index * 0.075, 0.22, 0.04, size=9.5, color=SILVER, bold=index < 2))
    slide["texts"].extend([
        text("chart-legend-note", "目标 / 实际（亿元）", 0.10, 0.75, 0.30, 0.03, size=8, color=MUTED),
        text("chart-source-note", "数据口径：case-suite.json / exact data snapshot", 0.58, 0.90, 0.35, 0.03, size=7.5, color=MUTED, align="right"),
    ])
    return deck


def visual_system_layout(case, run_dir, optimized):
    deck = make_base(case, run_dir)
    slide = deck["slides"][0]
    add_title(slide, case["title"], "连接 × 算力 × 感知 × 智能")
    slide["shapes"].append(shape("orbital-core", "oval", 0.38, 0.30, 0.24, 0.34, gradient={"angle": 45, "stops": [{"position": 0, "color": BLUE}, {"position": 0.55, "color": "#7B4DFF"}, {"position": 1, "color": RED}]}, line=BLUE, line_width=1.2))
    slide["shapes"].extend([shape("orbit-ring-1", "oval", 0.32, 0.25, 0.36, 0.44, line=BLUE, line_width=1.0), shape("orbit-ring-2", "oval", 0.35, 0.34, 0.30, 0.26, line=RED, line_width=0.8)])
    for index, label in enumerate(["连接", "算力", "感知", "智能"]):
        positions = [(0.12, 0.31), (0.68, 0.31), (0.12, 0.57), (0.68, 0.57)]
        x, y = positions[index]
        slide["groups"].append(panel(f"capability-{index + 1}", x, y, 0.18, 0.13, RED if index % 2 == 0 else BLUE, "#102F4F"))
        slide["texts"].append(text(f"capability-{index + 1}-label", label, x + 0.02, y + 0.045, 0.14, 0.04, size=11, bold=True, align="center"))
    for index, label in enumerate(["低时延", "大上行", "确定性网络"]):
        slide["groups"].append(panel(f"value-{index + 1}", 0.35 + index * 0.22, 0.80, 0.18, 0.08, BLUE, "#0D2F52"))
        slide["texts"].append(text(f"value-{index + 1}-text", label, 0.35 + index * 0.22, 0.822, 0.18, 0.03, size=8.2, bold=True, align="center"))
    slide["texts"].append(text("asset-note", "独立视觉资产边界：轨道、渐变、节点、标签", 0.68, 0.90, 0.25, 0.03, size=7.5, color=MUTED, align="right"))
    return deck


def irregular_layout(case, run_dir, optimized):
    deck = make_base(case, run_dir)
    slide = deck["slides"][0]
    add_title(slide, case["title"], "识别—分层—触达—转化—留存")
    for index, label in enumerate(["识别", "分层", "触达", "转化", "留存"]):
        x = 0.08 + index * 0.13
        y = 0.33 + (index % 2) * 0.08
        slide["groups"].append(panel(f"stage-{index + 1}", x, y, 0.105, 0.12, RED if index in {0, 3} else BLUE, "#0D2F52"))
        slide["texts"].append(text(f"stage-{index + 1}-text", label, x, y + 0.045, 0.105, 0.04, size=9.5, bold=True, align="center"))
        if index < 4:
            slide["shapes"].append(line(f"stage-link-{index + 1}", x + 0.105, y + 0.06, x + 0.13, 0.39 + ((index + 1) % 2) * 0.08, RED, 1.4))
    for index, label in enumerate(["高价值", "潜力用户", "风险用户"]):
        x = 0.30 + index * 0.18
        slide["groups"].append(panel(f"segment-{index + 1}", x, 0.57, 0.14, 0.16, [RED, BLUE, AMBER][index], "#123758"))
        slide["texts"].append(text(f"segment-{index + 1}-text", label, x + 0.01, 0.63, 0.12, 0.04, size=9, bold=True, align="center"))
    slide["groups"].append(panel("battle-result", 0.75, 0.33, 0.18, 0.38, RED, "#5C1120"))
    slide["texts"].append(text("battle-result-text", "作战结果", 0.77, 0.47, 0.14, 0.05, size=14, bold=True, align="center"))
    slide["texts"].extend([
        text("irregular-actions", "行动清单\n01 识别高价值\n02 分层定策略\n03 复盘沉淀", 0.08, 0.76, 0.38, 0.11, size=9, color=SILVER),
        text("irregular-note", "折角面板 / 斜切节奏 / 连接线层级 / native groups", 0.57, 0.90, 0.35, 0.03, size=7.5, color=MUTED, align="right"),
    ])
    return deck


def multi_module_layout(case, run_dir, optimized):
    deck = make_base(case, run_dir)
    slide = deck["slides"][0]
    add_title(slide, case["title"], "统一口径｜统一节奏｜统一动作")
    modules = ["目标设定", "过程经营", "问题诊断", "资源调度", "结果复盘"]
    coords = [(0.09, 0.31), (0.30, 0.20), (0.57, 0.20), (0.78, 0.31), (0.38, 0.61)]
    for index, label in enumerate(modules):
        x, y = coords[index]
        slide["groups"].append(panel(f"module-card-{index + 1}", x, y, 0.17, 0.18, RED if index == 4 else BLUE, "#0D2F52"))
        slide["texts"].extend([
            text(f"module-card-{index + 1}-index", f"0{index + 1}", x + 0.015, y + 0.025, 0.04, 0.03, size=8, color=RED if index == 4 else BLUE, bold=True),
            text(f"module-card-{index + 1}-title", label, x + 0.015, y + 0.07, 0.14, 0.04, size=9.5, bold=True, align="center"),
            text(f"module-card-{index + 1}-tag", "同源组件", x + 0.015, y + 0.13, 0.14, 0.03, size=7.5, color=MUTED, align="center"),
        ])
    slide["shapes"].append(shape("loop-center", "oval", 0.40, 0.38, 0.20, 0.18, fill="#102F4F", line=RED, line_width=1.3))
    slide["texts"].append(text("loop-center-text", "统一口径\n统一节奏\n统一动作", 0.415, 0.415, 0.17, 0.10, size=9, bold=True, align="center"))
    for index, (x, y) in enumerate(coords):
        slide["shapes"].append(line(f"loop-link-{index + 1}", 0.50, 0.47, x + 0.085, y + 0.09, RED if index == 4 else BLUE, 1.2))
    slide["texts"].append(text("consistency-note", "五个重复模块使用同一组件 token、几何和文字基线", 0.07, 0.90, 0.60, 0.03, size=8.2, color=MUTED))
    return deck


def governance_layout(case, run_dir, optimized):
    deck = make_base(case, run_dir)
    slide = deck["slides"][0]
    add_title(slide, case["title"], "未知问题不猜测，证据不足就停机")
    gates = ["输入校验", "路由确认", "原生结构", "视觉回归", "人工收口"]
    for index, label in enumerate(gates):
        add_node(slide, f"gate-{index + 1}", 0.08 + index * 0.16, 0.35, 0.13, label, index + 1, RED if index == 2 else BLUE)
        if index < 4:
            slide["shapes"].append(line(f"gate-link-{index + 1}", 0.21 + index * 0.16, 0.402, 0.235 + index * 0.16, 0.402, BLUE, 1.4))
    slide["groups"].append(panel("stop-panel", 0.70, 0.52, 0.23, 0.24, RED, "#5C1120"))
    slide["texts"].extend([
        text("stop-panel-title", "立即停机", 0.73, 0.57, 0.17, 0.05, size=16, color=RED, bold=True, align="center"),
        text("stop-panel-copy", "未知问题\n不得静默回退", 0.73, 0.65, 0.17, 0.06, size=10, bold=True, align="center"),
    ])
    slide["shapes"].append(line("failure-branch", 0.72, 0.43, 0.80, 0.52, RED, 2.0))
    slide["texts"].append(text("evidence-chain", "baseline → candidate → object audit → visual compare → mutation smoke", 0.08, 0.72, 0.55, 0.04, size=8.5, color=MUTED))
    slide["texts"].append(text("governance-note", "5 gates / 1 failure branch / hard stop", 0.66, 0.90, 0.27, 0.03, size=7.5, color=MUTED, align="right"))
    return deck


def package_layout(case, run_dir, optimized):
    deck = make_base(case, run_dir)
    slide = deck["slides"][0]
    add_title(slide, case["title"], "同一源文件 → 同一证据 → 可复现交付")
    layers = ["技能包", "运行时", "字体", "渲染器", "PPTX包"]
    for index, label in enumerate(layers):
        x = 0.08 + index * 0.045
        y = 0.28 + index * 0.055
        w = 0.42 - index * 0.015
        slide["groups"].append(panel(f"layer-{index + 1}", x, y, w, 0.10, RED if index == 0 else BLUE, "#0D2F52"))
        slide["texts"].append(text(f"layer-{index + 1}-text", label, x + 0.03, y + 0.032, w - 0.06, 0.035, size=10, bold=True))
    rows = [["环境", "字体", "渲染器", "PPTX包"], ["桌面WPS", "可见", "一致", "通过"], ["移动端WPS", "可见", "一致", "通过"], ["CI Runner", "可见", "一致", "哈希锁定"]]
    add_table(slide, "environment-matrix", 0.57, 0.28, 0.36, 0.30, rows, column_widths=[0.13, 0.08, 0.08, 0.10], row_heights=[0.06] + [0.08] * 3)
    for index, label in enumerate(["包完整", "字体可见", "渲染一致", "哈希锁定"]):
        slide["groups"].append(panel(f"lock-{index + 1}", 0.08 + index * 0.22, 0.72, 0.18, 0.09, RED if index == 3 else BLUE, "#102F4F"))
        slide["texts"].append(text(f"lock-{index + 1}-text", label, 0.08 + index * 0.22, 0.748, 0.18, 0.03, size=8.5, bold=True, align="center"))
    slide["texts"].append(text("package-note", "5-layer architecture / 3 target environments / native matrix", 0.52, 0.90, 0.41, 0.03, size=7.5, color=MUTED, align="right"))
    return deck


def pipeline_layout(case, run_dir, optimized):
    deck = make_base(case, run_dir)
    slide = deck["slides"][0]
    add_title(slide, case["title"], "页面级缓存、并行QA、重复运行结果稳定")
    for index, (label, color) in enumerate(zip(["源文件哈希", "页面级缓存", "受影响区域", "并行QA", "重复运行"], [RED, BLUE, BLUE, BLUE, GREEN])):
        x = 0.07 + index * 0.17
        add_node(slide, f"pipeline-{index + 1}", x, 0.39, 0.14, label, index + 1, color)
        if index < 4:
            slide["shapes"].append(line(f"pipeline-link-{index + 1}", x + 0.14, 0.442, x + 0.17, 0.442, color, 1.6))
    branches = ["视觉比较", "对象审计", "字体检查"]
    for index, label in enumerate(branches):
        y = 0.58 + index * 0.08
        slide["groups"].append(panel(f"parallel-{index + 1}", 0.34, y, 0.16, 0.055, BLUE, "#102F4F"))
        slide["texts"].append(text(f"parallel-{index + 1}-text", label, 0.34, y + 0.01, 0.16, 0.03, size=8, bold=True, align="center"))
        slide["shapes"].append(line(f"parallel-link-{index + 1}", 0.41, 0.495, 0.42, y, BLUE, 1.0))
    for index, (label, value, color) in enumerate([("缓存命中率", "86%", BLUE), ("返工次数", "1", RED), ("结果稳定", "100%", GREEN)]):
        x = 0.64 + index * 0.10
        slide["groups"].append(panel(f"metric-{index + 1}", x, 0.60, 0.08, 0.14, color, "#103B4B"))
        slide["texts"].append(text(f"metric-{index + 1}-value", value, x, 0.625, 0.08, 0.04, size=12, color=color, bold=True, align="center"))
        slide["texts"].append(text(f"metric-{index + 1}-label", label, x - 0.015, 0.69, 0.11, 0.03, size=7, color=MUTED, align="center"))
    slide["texts"].append(text("pipeline-timeline", "首轮 → 修复轮 → 回归轮", 0.07, 0.86, 0.36, 0.03, size=8.5, color=MUTED))
    slide["texts"].append(text("pipeline-note", "5-stage pipeline / 3 parallel QA branches / idempotent replay", 0.49, 0.90, 0.44, 0.03, size=7.5, color=MUTED, align="right"))
    return deck


def build_layout(case, run_dir, optimized):
    page_type = case["page_type"]
    builders = {
        "dashboard": dashboard_layout,
        "framework": framework_layout,
        "table": basic_table_layout,
        "policy-table": merge_table_layout,
        "text-visual": text_visual_layout,
        "chart": chart_layout,
        "visual-system": visual_system_layout,
        "irregular-framework": irregular_layout,
        "multi-module-overview": multi_module_layout,
        "governance-flow": governance_layout,
        "architecture-matrix": package_layout,
        "pipeline": pipeline_layout,
    }
    return builders[page_type](case, run_dir, optimized)


def add_expected_merges(manifest, case):
    spans = case.get("data", {}).get("merge_spans", [])
    if not spans:
        return manifest
    for slide in manifest.get("slides", []):
        for obj in slide.get("objects", []):
            if obj.get("object_id") == "policy-merged-table":
                obj["merges"] = copy.deepcopy(spans)
    return manifest


def make_text_manifest(layout, out):
    slides = []
    for slide_no, slide in enumerate(layout.get("slides", []), 1):
        specs = []
        for item in slide.get("texts", []):
            content = item.get("text")
            if content is None and item.get("runs"):
                content = "".join(str(run.get("text", "")) for run in item["runs"])
            specs.append({"object_id": item["object_id"], "text_id": item["object_id"], "content": str(content or ""), "source_ref": "case-suite.json"})
        slides.append({"slide_no": slide_no, "text_specs": specs})
    write_json(out, {"schema": "ai-ppt-plus/text-layout-manifest/v1", "slides": slides})


def make_baseline_layout(case, run_dir, reference):
    return {"project_id": "distillation-case-replay-12", "slide_width_in": 13.333, "slide_height_in": 7.5, "ref_width": 1920, "ref_height": 1080, "units": "fraction", "assets_dir": repo_relative(run_dir), "slides": [{"layout_name": "Blank", "frame": os.path.relpath(reference, run_dir), "frame_object_id": "legacy-full-slide-frame", "frame_role": "legacy-flattened-control"}]}


def mutate_deck(source, target, case):
    from pptx import Presentation

    prs = Presentation(str(source))
    slide = prs.slides[0]
    table_shapes = [shape for shape in slide.shapes if getattr(shape, "has_table", False)]
    table_changed = False
    if table_shapes:
        table_shapes[0].table.cell(0, 0).text = "MUTATION_OK"
        table_changed = True
    panel_changed = False
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in slide.shapes:
        if str(shape.name) == "background" or getattr(shape, "has_table", False) or getattr(shape, "has_chart", False):
            continue
        if getattr(shape, "shape_type", None) in {MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.TEXT_BOX}:
            shape.left = int(shape.left + prs.slide_width * 0.012)
            panel_changed = True
            break
    prs.save(str(target))
    return {"table_cell_changed": table_changed, "panel_or_group_moved": panel_changed}


def image_change(before, after):
    from PIL import Image, ImageChops

    with Image.open(before).convert("RGB") as first, Image.open(after).convert("RGB") as second:
        if first.size != second.size:
            second = second.resize(first.size)
        diff = ImageChops.difference(first, second)
        bbox = diff.getbbox()
        extrema = diff.getextrema()
        changed = sum(1 for channel in extrema if channel[1] > 0)
        return {"changed": bool(bbox), "bbox": list(bbox) if bbox else None, "changed_channels": changed}


def evaluate_visual(case, run_dir, deck, reference, label):
    render_dir = run_dir / "render"
    render_report = run_dir / "render-report.json"
    run([sys.executable, RENDER, str(deck), "--output-dir", str(render_dir), "--dpi", "126", "--font-dir", str(FONT_DIR), "--report", str(render_report)], allow_failure=True)
    rendered = render_dir / "slide-1.png"
    if not rendered.exists():
        rendered = next(render_dir.glob("*.png"), None)
    inspect_report = run_dir / "inspect.json"
    run([sys.executable, INSPECT, str(deck), "--report", str(inspect_report)], allow_failure=True)
    visual_report = run_dir / "visual-compare.json"
    visual_result = run([sys.executable, EDITABLE / "scripts" / "compare_visual.py", str(rendered), str(reference), "--expected-ratio", "1.7777777778", "--raw-slide", "--report", str(visual_report)], allow_failure=True)
    object_manifest = run_dir / "object-manifest.json"
    native_report = run_dir / "native-editability.json"
    semantic_report = run_dir / "semantic-audit.json"
    case_report = run_dir / "case-replay-audit.json"
    native_result = run([sys.executable, NATIVE_VALIDATE, str(deck), "--object-manifest", str(object_manifest), "--require-native-structure", "--require-native-panels", "--require-native-tables", "--forbid-whole-slide-pictures", "--require-complete-manifest", "--report", str(native_report)], allow_failure=True)
    semantic_result = run([sys.executable, SEMANTIC_AUDIT, str(deck), "--object-manifest", str(object_manifest), "--text-manifest", str(run_dir / "text-manifest.json"), "--require-independent-text-manifest", "--report", str(semantic_report)], allow_failure=True)
    case_result = run([sys.executable, CASE_AUDIT, str(deck), "--object-manifest", str(object_manifest), "--report", str(case_report)], allow_failure=True)
    mutation_path = run_dir / "mutation.pptx"
    mutation_render_dir = run_dir / "mutation-render"
    mutation_render = run_dir / "mutation-render-report.json"
    mutation_action = mutate_deck(deck, mutation_path, case)
    run([sys.executable, RENDER, str(mutation_path), "--output-dir", str(mutation_render_dir), "--dpi", "126", "--font-dir", str(FONT_DIR), "--report", str(mutation_render)], allow_failure=True)
    mutated_rendered = mutation_render_dir / "slide-1.png"
    if not mutated_rendered.exists():
        mutated_rendered = next(mutation_render_dir.glob("*.png"), None)
    mutation_evidence = {**mutation_action, "rendered": bool(mutated_rendered and mutated_rendered.exists())}
    if mutated_rendered and mutated_rendered.exists() and rendered and rendered.exists():
        mutation_evidence["pixel_change"] = image_change(rendered, mutated_rendered)
    else:
        mutation_evidence["pixel_change"] = {"changed": False, "bbox": None}
    inspect = json.loads(inspect_report.read_text(encoding="utf-8")) if inspect_report.exists() else {}
    visual = json.loads(visual_report.read_text(encoding="utf-8")) if visual_report.exists() else {}
    native = json.loads(native_report.read_text(encoding="utf-8")) if native_report.exists() else {}
    semantic = json.loads(semantic_report.read_text(encoding="utf-8")) if semantic_report.exists() else {}
    case_audit = json.loads(case_report.read_text(encoding="utf-8")) if case_report.exists() else {}
    return {
        "label": label,
        "deck": case_relative(deck),
        "deck_sha256": digest(deck),
        "rendered": case_relative(rendered) if rendered else None,
        "rendered_sha256": digest(rendered) if rendered and rendered.exists() else None,
        "reference": case_relative(reference),
        "reference_sha256": digest(reference),
        "inspect": {"slide": (inspect.get("slides") or [{}])[0], "issues": inspect.get("issues", [])},
        "objects": {"native_table_count": native.get("native_table_count", 0), "native_panel_count": native.get("native_panel_count", 0), "formal_text_count": case_audit.get("formal_text_count", 0), "formal_text_native_count": case_audit.get("formal_text_native_count", 0), "a_tbl_count": case_audit.get("a_tbl_count", 0), "whole_slide_pictures": case_audit.get("whole_slide_pictures", []), "native_editability_valid": native.get("valid", False), "semantic_audit_valid": semantic.get("valid", False), "case_replay_audit_valid": case_audit.get("valid", False), "semantic_errors": semantic.get("errors", []), "case_replay_errors": case_audit.get("errors", [])},
        "visual": {"valid": visual.get("valid", False), "metrics": visual.get("metrics", {}), "issues": visual.get("issues", [])},
        "mutation_smoke": mutation_evidence,
        "technical_status": "passed" if native.get("valid") and semantic.get("valid") and case_audit.get("valid") and mutation_evidence.get("pixel_change", {}).get("changed") else "blocked",
    }


def legacy_evaluate(case, run_dir, reference):
    layout = make_baseline_layout(case, run_dir, reference)
    layout_path = run_dir / "layout.json"
    write_json(layout_path, layout)
    deck = run_dir / "legacy-image-only.pptx"
    run([sys.executable, COMPOSE, str(layout_path), str(deck)], allow_failure=False)
    render_dir = run_dir / "render"
    render_report = run_dir / "render-report.json"
    run([sys.executable, RENDER, str(deck), "--output-dir", str(render_dir), "--dpi", "126", "--font-dir", str(FONT_DIR), "--report", str(render_report)], allow_failure=True)
    rendered = render_dir / "slide-1.png"
    if not rendered.exists():
        rendered = next(render_dir.glob("*.png"), None)
    visual_report = run_dir / "visual-compare.json"
    run([sys.executable, EDITABLE / "scripts" / "compare_visual.py", str(rendered), str(reference), "--expected-ratio", "1.7777777778", "--raw-slide", "--report", str(visual_report)], allow_failure=True)
    inspect_report = run_dir / "inspect.json"
    run([sys.executable, INSPECT, str(deck), "--report", str(inspect_report)], allow_failure=True)
    inspect = json.loads(inspect_report.read_text(encoding="utf-8")) if inspect_report.exists() else {}
    visual = json.loads(visual_report.read_text(encoding="utf-8")) if visual_report.exists() else {}
    return {"label": "legacy-image-only-control", "deck": case_relative(deck), "deck_sha256": digest(deck), "rendered": case_relative(rendered) if rendered else None, "rendered_sha256": digest(rendered) if rendered and rendered.exists() else None, "reference": case_relative(reference), "reference_sha256": digest(reference), "inspect": {"slide": (inspect.get("slides") or [{}])[0], "issues": inspect.get("issues", [])}, "objects": {"native_table_count": 0, "native_panel_count": 0, "formal_text_count": 0, "formal_text_native_count": 0, "a_tbl_count": 0, "whole_slide_pictures": [{"name": "legacy-full-slide-frame"}], "native_editability_valid": False, "semantic_audit_valid": False, "case_replay_audit_valid": False, "semantic_errors": [{"code": "legacy_full_slide_picture"}], "case_replay_errors": [{"code": "native_structure_missing"}, {"code": "formal_text_rasterized"}]}, "visual": {"valid": visual.get("valid", False), "metrics": visual.get("metrics", {}), "issues": visual.get("issues", [])}, "mutation_smoke": {"status": "blocked", "reason": "legacy image-only control has no native table/panel to mutate"}, "technical_status": "blocked"}


def native_evaluate(case, run_dir, reference, optimized):
    layout = build_layout(case, run_dir, optimized)
    layout_path = run_dir / "layout.json"
    write_json(layout_path, layout)
    deck = run_dir / "editable.pptx"
    run([sys.executable, COMPOSE, str(layout_path), str(deck), "--strict-input", "--require-native-structure"], allow_failure=False)
    manifest_result = run([sys.executable, MANIFEST_BUILDER, str(layout_path), "--output", str(run_dir / "object-manifest.json")], allow_failure=False)
    manifest_path = run_dir / "object-manifest.json"
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    if not optimized:
        manifest = add_expected_merges(manifest, case)
        write_json(manifest_path, manifest)
    make_text_manifest(layout, run_dir / "text-manifest.json")
    result = evaluate_visual(case, run_dir, deck, reference, "optimized-native-reconstruction" if optimized else "pre-fix-native-reconstruction")
    result["layout_sha256"] = digest(layout_path)
    result["object_manifest_sha256"] = digest(manifest_path)
    result["text_manifest_sha256"] = digest(run_dir / "text-manifest.json")
    return result


def aggregate(cases, runs):
    return {
        "schema": "ai-ppt-plus/case-replay-evaluation/v1",
        "suite_id": "distillation-case-replay-12",
        "case_count": len(cases),
        "cases": runs,
        "rollup": {
            "legacy_visual_pass": sum(bool(item["baseline"]["visual"].get("valid")) for item in runs),
            "pre_fix_native_technical_pass": sum(item["pre_fix"]["technical_status"] == "passed" for item in runs),
            "optimized_native_technical_pass": sum(item["candidate"]["technical_status"] == "passed" for item in runs),
            "legacy_native_table_total": sum(item["baseline"]["objects"].get("native_table_count", 0) for item in runs),
            "optimized_native_table_total": sum(item["candidate"]["objects"].get("native_table_count", 0) for item in runs),
            "optimized_a_tbl_total": sum(item["candidate"]["objects"].get("a_tbl_count", 0) for item in runs),
            "optimized_formal_text_native_ratio": round(sum(item["candidate"]["objects"].get("formal_text_native_count", 0) for item in runs) / max(1, sum(item["candidate"]["objects"].get("formal_text_count", 0) for item in runs)), 6),
            "mutation_smoke_pass": sum(bool(item["candidate"]["mutation_smoke"].get("pixel_change", {}).get("changed")) for item in runs),
        },
        "human_visual_review_required": True,
        "release_eligible": False,
    }


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--suite", default=str(ROOT / "case-suite.json"))
    parser.add_argument("--output-dir", default=str(ROOT / "runs"))
    parser.add_argument("--strict", action="store_true", help="fail unless every optimized case passes the technical replay gates")
    args = parser.parse_args()
    suite = json.loads(Path(args.suite).read_text(encoding="utf-8"))
    output_dir = Path(args.output_dir).resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    all_runs = []
    for case in suite["cases"]:
        case_id = case["case_id"]
        reference = ROOT / "visual" / f"{case_id}-reference.png"
        baseline_dir = output_dir / "baseline" / case_id
        pre_dir = output_dir / "candidate-before" / case_id
        candidate_dir = output_dir / "candidate" / case_id
        baseline_dir.mkdir(parents=True, exist_ok=True)
        pre_dir.mkdir(parents=True, exist_ok=True)
        candidate_dir.mkdir(parents=True, exist_ok=True)
        baseline = legacy_evaluate(case, baseline_dir, reference)
        pre_fix = native_evaluate(case, pre_dir, reference, optimized=False)
        candidate = native_evaluate(case, candidate_dir, reference, optimized=True)
        all_runs.append({"case_id": case_id, "title": case["title"], "priority": case["priority"], "responsibility": case["responsibility"], "baseline": baseline, "pre_fix": pre_fix, "candidate": candidate, "improved": {"native_table_count_delta": candidate["objects"].get("native_table_count", 0) - baseline["objects"].get("native_table_count", 0), "a_tbl_count_delta": candidate["objects"].get("a_tbl_count", 0) - baseline["objects"].get("a_tbl_count", 0), "formal_text_native_delta": candidate["objects"].get("formal_text_native_count", 0) - baseline["objects"].get("formal_text_native_count", 0), "technical_status_changed": pre_fix["technical_status"] != candidate["technical_status"], "merge_gate_fixed": any(error.get("code") in {"case_replay_merge_topology_mismatch", "case_replay_table_not_native"} for error in pre_fix["objects"].get("case_replay_errors", [])) and candidate["technical_status"] == "passed"}})
    baseline_eval = {"schema": "ai-ppt-plus/baseline-evaluation/v1", "evaluation_kind": "pre-distillation legacy image-only control replay", "suite_id": suite["suite_id"], "source_case_suite_sha256": digest(Path(args.suite)), "cases": [{"case_id": item["case_id"], "reference_sha256": item["baseline"]["reference_sha256"], "deck": item["baseline"]["deck"], "deck_sha256": item["baseline"]["deck_sha256"], "rendered": item["baseline"]["rendered"], "rendered_sha256": item["baseline"]["rendered_sha256"], "failure_codes": item["baseline"]["objects"].get("case_replay_errors", []), "original_pptx_object_counts": item["baseline"]["objects"], "rendered_visual_metrics": item["baseline"]["visual"], "table_panel_text_audit": {"tables": 0, "panels": 0, "formal_text_native": False}, "mutation_smoke": item["baseline"]["mutation_smoke"]} for item in all_runs], "rollup": {"case_count": len(all_runs), "native_table_total": 0, "full_slide_picture_cases": len(all_runs), "formal_text_native_cases": 0, "visual_pass_cases": sum(bool(item["baseline"]["visual"].get("valid")) for item in all_runs)}, "human_visual_review_required": True, "release_eligible": False}
    rollup = aggregate(suite["cases"], all_runs)["rollup"]
    strict_failures = []
    if args.strict:
        if rollup["legacy_visual_pass"] != len(all_runs):
            strict_failures.append("legacy visual replay")
        if rollup["optimized_native_technical_pass"] != len(all_runs):
            strict_failures.append("optimized native technical replay")
        if rollup["mutation_smoke_pass"] != len(all_runs):
            strict_failures.append("mutation smoke")
    candidate_eval = {"schema": "ai-ppt-plus/candidate-evaluation/v1", "evaluation_kind": "post-distillation actual 12-case replay", "suite_id": suite["suite_id"], "source_case_suite_sha256": digest(Path(args.suite)), "skill_revision": "ai-ppt-editable-native-structure-plus-merge-topology-gate", "cases": [{"case_id": item["case_id"], "title": item["title"], "priority": item["priority"], "responsibility": item["responsibility"], "pre_fix": item["pre_fix"], "candidate": item["candidate"], "improvement": item["improved"]} for item in all_runs], "rollup": rollup, "strict_gate": {"requested": bool(args.strict), "passed": not strict_failures, "failures": strict_failures}, "human_visual_review_required": True, "release_eligible": False}
    write_json(ROOT / "baseline-evaluation.json", baseline_eval)
    write_json(ROOT / "candidate-evaluation.json", candidate_eval)
    write_json(ROOT / "case-improvement.json", {"schema": "ai-ppt-plus/case-improvement/v1", "suite_id": suite["suite_id"], "cases": all_runs, "rollup": aggregate(suite["cases"], all_runs)["rollup"], "human_visual_review_required": True, "release_eligible": False})
    result = {"cases": len(all_runs), "baseline": str(ROOT / "baseline-evaluation.json"), "candidate": str(ROOT / "candidate-evaluation.json"), "rollup": rollup, "strict_gate": candidate_eval["strict_gate"]}
    print(json.dumps(result, ensure_ascii=False))
    if strict_failures:
        raise SystemExit("case replay strict gate failed: " + ", ".join(strict_failures))


if __name__ == "__main__":
    main()
